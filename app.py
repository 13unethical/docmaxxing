"""
Flask entrypoint. All document logic lives under `formatter/` for clarity.
"""

from __future__ import annotations

from pathlib import Path

try:
    from dotenv import load_dotenv

    load_dotenv(Path(__file__).resolve().parent / ".env")
except ImportError:
    pass

import io
import json
import os
import re
import threading
import time
import traceback
import uuid
from typing import Any

import requests
from docx import Document
from flask import Flask, jsonify, redirect, render_template, request, send_file, session, url_for

from formatter import FormatJob, format_document_full
from formatter.document_reconstruction import reconstruct_document_before_format
from formatter.cover_page import CoverPageData, prepend_cover_document, prepend_cover_page
from formatter.heading_plan import ParagraphHeadingAssignment
from formatter.preview_html import build_formatted_preview_html
from formatter.references_section import append_references_section
from formatter.document_io import (
    build_document_from_inputs,
    build_document_from_upload,
    document_has_visible_content,
    extract_text_from_document_bytes,
    is_supported_document_upload,
    upload_extension,
)
from services.reference_list_formatter import prepare_reference_section
from services.citation_engine import CITATION_STYLES as ENGINE_CITATION_STYLES
from services.citation_engine import generate_citation
from services.document_checker import MAX_TEXT_CHARS, check_document
from services.document_structure_engine import infer_assignment_title, paragraphs_from_text, recover_structure
from services.intext_citations import generate_intext
from services.requirements_ocr import extract_text_from_image_stream
from services.requirements_parser import form_autofill_from_parsed, parse_requirements
from services.assignment_pipeline import (
    PIPELINE_STAGE_SPECS,
    AssignmentPipelineService,
    PipelineStage,
)
from services.assignment_pipeline.handlers import StageResult
from services.assignment_pipeline.models import utc_now
from services.assignment_project import ProjectService
from services.assignment_project.paths import assignment_storage_root, project_files_dir
from services.assignment_project.store import ProjectStore
from services.assignment_project.trace_log import trace, trace_startup
from services.research_engine import ResearchEngineService
from services.research_engine.models import ParsedDocument
from services.blueprint_engine import BlueprintEngineService
from services.writer_engine import WriterEngineService
from services.reviewer_engine import GeminiAcademicReviewer, ReviewerEngineService
from services.revision_engine import GeminiSectionReviser, RevisionEngineService
from services.humanizer_engine import HumanizerEngineService
from services.humanizer_engine.mock_validator import ZeroGPTParagraphValidator
from services.humanizer_engine.stealthwriter_humanizer import StealthWriterTextHumanizer
from services.humanizer_engine.zerogpt_humanizer import count_words
from services.ai_detection_engine import AIDetectionEngineService
from services.ai_detection_engine.zerogpt_detector import ZeroGPTAIDetector
from services.delivery_engine import DeliveryEngineService
from services.llm_errors import llm_error_http_status, user_friendly_llm_error
from services.zerogpt_business import (
    ZeroGPTClient,
    ZeroGPTDetectionProvider,
    ZeroGPTError,
    ZeroGPTHumanizerProvider,
    ZeroGPTProviderError,
    orchestrator_review,
)
from services.citation_service import CitationService, CrossrefProvider
from services.economy import (
    FEATURE_LABELS,
    TOPUP_PACKAGES,
    WELCOME_BONUS,
    InsufficientCoins,
    WalletService,
    feature_cost,
    init_db as economy_init_db,
    package as economy_package,
)
from services.economy import auth as economy_auth
from services.economy.admin import AdminError, AdminService, bootstrap_admin_from_env
from services.economy.avatar_upload import AvatarUploadError, validate_and_store_avatar
from services.economy.disposable_email import DisposableEmailError
from services.economy.email_verify import OTP_TTL_MINUTES
from services.economy.paddle_purchases import (
    PaddlePurchaseService,
)
from services.economy.paddle_gateway import (
    PaddleGatewayError,
    PaddleSignatureError,
    apply_paid_purchase_atomic,
    create_checkout,
    handle_webhook_event,
    mock_topup_allowed,
    paddle_client_token,
    paddle_configured,
    paddle_environment,
    verify_paddle_signature,
)
from services.economy.cryptomus_gateway import (
    CryptomusGatewayError,
    CryptomusSignatureError,
    create_invoice as cryptomus_create_invoice,
    cryptomus_configured,
    handle_webhook as cryptomus_handle_webhook,
    parse_webhook_payload as cryptomus_parse_webhook_payload,
)
from services.economy.gumroad_gateway import (
    GumroadGatewayError,
    handle_ping as gumroad_handle_ping,
)
from services.economy.lemon_squeezy_gateway import (
    LemonSqueezyGatewayError,
    LemonSqueezySignatureError,
    handle_webhook_event as lemon_squeezy_handle_webhook,
    verify_lemon_squeezy_signature,
)
from services.economy.pricing import USD_TO_COINS
from services.economy.usage import (
    FEATURE_ASSIGNMENT,
    FEATURE_DETECTION,
    FEATURE_HUMANIZER,
    FEATURE_TURNITIN,
    UsageService,
)
from services.turnitin_service import TurnitinService, init_db as turnitin_init_db

app = Flask(__name__)
# Document / assignment uploads may be large; avatar route enforces its own 2 MB cap.
app.config["MAX_CONTENT_LENGTH"] = 16 * 1024 * 1024
_REPO_ROOT = Path(__file__).resolve().parent
_AVATAR_MAX_BYTES = 2 * 1024 * 1024


def _env_bool(name: str, default: bool = False) -> bool:
    raw = (os.environ.get(name) or "").strip().lower()
    if not raw:
        return default
    if raw in ("1", "true", "yes", "y", "on"):
        return True
    if raw in ("0", "false", "no", "n", "off"):
        return False
    return default


# SMTP / Resend — loaded from .env via load_dotenv above.
app.config["MAIL_SERVER"] = (os.environ.get("MAIL_SERVER") or "").strip()
app.config["MAIL_PORT"] = int((os.environ.get("MAIL_PORT") or "587").strip() or "587")
app.config["MAIL_USERNAME"] = (os.environ.get("MAIL_USERNAME") or "").strip()
app.config["MAIL_PASSWORD"] = (os.environ.get("MAIL_PASSWORD") or "").strip()
app.config["MAIL_USE_TLS"] = _env_bool("MAIL_USE_TLS", True)
app.config["MAIL_FROM"] = (os.environ.get("MAIL_FROM") or "").strip()


def _require_strong_secret_key() -> str:
    """Refuse startup if SECRET_KEY is missing or obviously weak."""
    raw = (os.environ.get("SECRET_KEY") or "").strip()
    weak_markers = (
        "dev-insecure-change-me",
        "change-me",
        "changeme",
        "secret",
        "development key",
        "your-secret",
        "replace-me",
    )
    lowered = raw.lower()
    if not raw:
        raise SystemExit(
            "FATAL: SECRET_KEY is required. Set a strong random value in the environment."
        )
    if len(raw) < 32:
        raise SystemExit(
            "FATAL: SECRET_KEY must be at least 32 characters."
        )
    if any(marker in lowered for marker in weak_markers):
        raise SystemExit(
            "FATAL: SECRET_KEY looks weak or is a placeholder. "
            "Generate a strong random secret (e.g. python -c \"import secrets; print(secrets.token_urlsafe(48))\")."
        )
    return raw


app.secret_key = _require_strong_secret_key()
app.config["PERMANENT_SESSION_LIFETIME"] = 60 * 60 * 24 * 30  # 30 days

economy_init_db()
turnitin_init_db()
try:
    from services.economy.telegram_poller import try_start_telegram_poller

    try_start_telegram_poller()
except Exception:  # noqa: BLE001
    import logging as _logging

    _logging.getLogger(__name__).exception("telegram poller failed to start")
bootstrap_admin_from_env()
wallet = WalletService()
paddle_purchases = PaddlePurchaseService()
usage_service = UsageService()
admin_service = AdminService(wallet, paddle_purchases, usage_service)
turnitin_service = TurnitinService()


def _charge_current_user(feature: str, cost: int, *, ref_id: str | None = None, meta: dict | None = None):
    """Debit the logged-in user. Returns (user_id, tx) or an error response.

    Callers use it as::

        charged = _charge_current_user("humanize", cost)
        if isinstance(charged, tuple) is False:
            return charged  # error Response
        user_id, _tx = charged
    """
    from services.economy.referral import apply_pro_discount, consume_free_turnitin_report

    user = economy_auth.current_user()
    if user is None:
        resp = jsonify(
            {"success": False, "error": "AUTH_REQUIRED", "message": "Please sign in to continue."}
        )
        resp.status_code = 401
        return resp
    if not economy_auth.user_email_verified(user):
        resp = jsonify(
            {
                "success": False,
                "error": "EMAIL_NOT_VERIFIED",
                "message": "Please verify your email before continuing.",
            }
        )
        resp.status_code = 403
        return resp

    effective = int(cost)
    used_free_turnitin = False
    if feature == "turnitin" and int(user.get("free_turnitin_reports") or 0) > 0:
        if consume_free_turnitin_report(int(user["id"])):
            used_free_turnitin = True
            effective = 0
            # Bust cached user so free report count refreshes
            try:
                from flask import g

                g._economy_user = None
            except RuntimeError:
                pass

    if effective > 0:
        effective = apply_pro_discount(user, effective)
        try:
            tx = wallet.debit(user["id"], effective, feature, ref_id=ref_id, meta=meta)
        except InsufficientCoins as exc:
            resp = jsonify(
                {
                    "success": False,
                    "error": "INSUFFICIENT_COINS",
                    "message": (
                        f"Not enough credits. This requires {exc.required} credits; "
                        f"you have {exc.balance}."
                    ),
                    "required": exc.required,
                    "balance": exc.balance,
                }
            )
            resp.status_code = 402
            return resp
        return (user["id"], tx)

    return (
        user["id"],
        {
            "amount": 0,
            "free_turnitin": used_free_turnitin,
            "feature": feature,
            "ref_id": ref_id,
        },
    )


def _send_user_verification_email(user: dict, *, background: bool = False) -> str | None:
    """Issue a fresh OTP, persist it, and email the 6-digit code (never a magic link)."""
    user_id = int(user["id"])
    try:
        code = economy_auth.issue_verification_otp(user_id)
    except economy_auth.AuthError:
        app.logger.exception("OTP issue failed user_id=%s", user_id)
        return None

    # Re-read from DB so the emailed value matches what verify will check.
    from services.economy.db import connect as economy_connect

    with economy_connect() as conn:
        row = conn.execute(
            "SELECT email, verification_code FROM users WHERE id = ?",
            (user_id,),
        ).fetchone()
    if row is None or not row["verification_code"]:
        app.logger.error("OTP missing after issue user_id=%s", user_id)
        return None
    code = str(row["verification_code"]).strip()
    to_email = str(row["email"] or user.get("email") or "").strip()
    name = user.get("name")

    def _do_send() -> None:
        # Import inside the worker so a stale reloader never keeps the old
        # magic-link implementation bound in a closure.
        from services.economy.email_verify import (
            EmailVerifyError as _EmailVerifyError,
            send_verification_otp_email as _send_otp,
        )

        try:
            _send_otp(to_email=to_email, code=code, name=name)
        except _EmailVerifyError:
            app.logger.exception("OTP email failed user_id=%s to=%s", user_id, to_email)

    if background:
        def _run() -> None:
            with app.app_context():
                _do_send()

        threading.Thread(target=_run, daemon=True).start()
    else:
        _do_send()
    return code


def _load_owned_project(project_id: str):
    """IDOR guard: return (bundle, None) or (None, flask_response)."""
    user = economy_auth.current_user()
    if user is None:
        return None, (
            jsonify(
                {
                    "success": False,
                    "error": "AUTH_REQUIRED",
                    "message": "Please sign in to continue.",
                }
            ),
            401,
        )
    try:
        bundle = project_service.get_project(project_id)
    except KeyError:
        return None, (jsonify({"error": "Project not found"}), 404)
    owner = str(getattr(bundle.project, "user_id", None) or "").strip()
    if economy_auth.is_admin(user):
        return bundle, None
    if not owner or owner != str(user["id"]):
        return None, (jsonify({"error": "Forbidden"}), 403)
    return bundle, None


def _require_project_owner_response(project_id: str):
    """Return error response if caller cannot access project, else None."""
    _bundle, err = _load_owned_project(project_id)
    return err


# Pages + APIs that must never be usable before email verification.
_EMAIL_VERIFIED_PREFIXES = (
    "/workspace",
    "/editor",
    "/humanizer",
    "/turnitin",
    "/assignment",
    "/assignments",
    "/api/workspace",
    "/api/humanizer",
    "/api/turnitin",
    "/api/assignment",
    "/api/browser/providers/stealthwriter/humanize",
)


def _path_requires_email_verification(path: str) -> bool:
    for prefix in _EMAIL_VERIFIED_PREFIXES:
        if path == prefix or path.startswith(prefix + "/"):
            return True
    return False


@app.before_request
def _email_verification_wall():
    """Strict wall: Workspace, Humanizer, Turnitin, Assignment + their APIs."""
    path = request.path or ""
    if not _path_requires_email_verification(path):
        return None
    return economy_auth.email_verification_gate()


@app.before_request
def _assignment_project_ownership_guard():
    """IDOR prevention: every /api/assignment/projects/<id>/… call must own the project."""
    path = request.path or ""
    marker = "/api/assignment/projects/"
    if not path.startswith(marker):
        return None
    rest = path[len(marker) :].lstrip("/")
    if not rest:
        return None
    project_id = rest.split("/", 1)[0]
    if project_id in {"upload"}:
        return None
    if len(project_id) < 10:
        return None
    return _require_project_owner_response(project_id)


def _refund_safe(user_id: int, cost: int, feature: str, *, ref_id: str | None = None) -> None:
    """Best-effort refund; never raises into the request path."""
    try:
        wallet.refund(user_id, int(cost), feature, ref_id=ref_id, meta={"reason": "auto-refund"})
    except Exception:  # noqa: BLE001
        app.logger.exception("refund failed for user=%s feature=%s", user_id, feature)


def _record_usage_safe(
    user_id: int,
    *,
    feature: str,
    credits_used: int,
    provider: str | None = None,
    provider_cost: float | None = None,
    latency: int | None = None,
    request_id: str | None = None,
) -> None:
    """Best-effort Usage row; never breaks the request path."""
    try:
        usage_service.record(
            user_id=int(user_id),
            feature=feature,
            credits_used=int(credits_used),
            provider=provider,
            provider_cost=provider_cost,
            latency=latency,
            request_id=request_id,
        )
    except Exception:  # noqa: BLE001
        app.logger.exception("usage record failed for user=%s feature=%s", user_id, feature)

assignment_pipeline = AssignmentPipelineService()
research_engine = ResearchEngineService()
blueprint_engine = BlueprintEngineService()
writer_engine = WriterEngineService()
reviewer_engine = ReviewerEngineService(reviewer=GeminiAcademicReviewer())
revision_engine = RevisionEngineService(
    draft_store=writer_engine.drafts,
    reviser=GeminiSectionReviser(),
)
zerogpt_client = ZeroGPTClient()
ai_detection_engine = AIDetectionEngineService(detector=ZeroGPTAIDetector(client=zerogpt_client))
delivery_engine = DeliveryEngineService()
citation_service = CitationService(CrossrefProvider())


def _zerogpt_configured() -> bool:
    api_key = (os.environ.get("ZEROGPT_API_KEY") or "").strip()
    email = (os.environ.get("ZEROGPT_EMAIL") or "").strip()
    password = (os.environ.get("ZEROGPT_PASSWORD") or "").strip()
    return bool(api_key or (email and password))


def _build_humanizer_engine() -> HumanizerEngineService:
    # Assignment humanization always goes through StealthWriter Legacy 5.1.
    # ZeroGPT remains available for detection / validation when configured.
    humanizer = StealthWriterTextHumanizer()
    if _zerogpt_configured():
        return HumanizerEngineService(
            humanizer=humanizer,
            validator=ZeroGPTParagraphValidator(),
        )
    return HumanizerEngineService(humanizer=humanizer)


humanizer_engine = _build_humanizer_engine()
zerogpt_detection_provider = ZeroGPTDetectionProvider(client=zerogpt_client)
zerogpt_humanizer_provider = ZeroGPTHumanizerProvider(client=zerogpt_client)
PROJECT_STORAGE_ROOT = assignment_storage_root()
project_service = ProjectService(
    store=ProjectStore(root=PROJECT_STORAGE_ROOT),
    pipeline=assignment_pipeline,
    research=research_engine,
    blueprint=blueprint_engine,
    writer=writer_engine,
    reviewer=reviewer_engine,
    revision=revision_engine,
    humanizer=humanizer_engine,
    ai_detection=ai_detection_engine,
    delivery=delivery_engine,
)
trace(
    "app.project_service.init",
    storage_root=str(PROJECT_STORAGE_ROOT),
    store_root=str(project_service.store.storage_root),
)
trace_startup()


def _assignment_not_found(endpoint: str, project_id: str, exc: KeyError) -> tuple[Any, int]:
    reason = str(exc)
    diagnostics = project_service.store.lookup_diagnostics(project_id)
    trace(
        "api.project_not_found",
        endpoint=endpoint,
        reason=reason,
        **diagnostics,
    )
    return jsonify({"error": "Project not found"}), 404

ALLOWED_FONTS = frozenset(
    {
        "Times New Roman",
        "Arial",
        "Calibri",
        "Cambria",
        "Georgia",
        "Verdana",
        "Tahoma",
    }
)
ALLOWED_FONT_SIZES = frozenset({10, 11, 12, 13, 14, 16, 18, 20})
LINE_SPACING_MAP = {"1.0": 1.0, "1.15": 1.15, "1.5": 1.5, "2.0": 2.0}
PAGE_POSITIONS = frozenset(
    {"none", "top_left", "top_right", "bottom_left", "bottom_right"}
)
MARGIN_PRESETS = frozenset({"normal", "narrow", "wide"})
ALIGNMENTS = frozenset({"left", "justify"})
CITATION_STYLES = frozenset({"APA", "MLA", "Harvard", "Chicago", "IEEE", "Vancouver"})
REQUIREMENTS_IMAGE_EXT = frozenset({".jpg", ".jpeg", ".png"})
REQUIREMENTS_TEXT_EXT = frozenset({".txt", ".md"})
REQUIREMENTS_DOC_EXT = frozenset({".docx", ".pdf"}) | REQUIREMENTS_TEXT_EXT | REQUIREMENTS_IMAGE_EXT

# Telegram Bot API: TELEGRAM_TOKEN + CHAT_ID (also accepts TELEGRAM_BOT_TOKEN + TELEGRAM_CHAT_ID).
TELEGRAM_SEND_MESSAGE_TIMEOUT_S = 12
TELEGRAM_TEXT_MAX_LEN = 4096


def _project_api_payload(bundle, *, include_pipeline: bool = True) -> dict[str, Any]:
    payload = bundle.to_dict()
    if include_pipeline:
        try:
            payload["pipeline"] = assignment_pipeline.get_project(bundle.project.id).to_dict()
        except Exception:  # noqa: BLE001
            payload["pipeline"] = None

    def _safe(key: str, loader):
        try:
            payload[key] = loader().to_dict()
        except Exception:  # noqa: BLE001
            payload[key] = None

    _safe("research_plan", lambda: project_service.get_research_plan(bundle.project.id))
    _safe("blueprint", lambda: project_service.get_blueprint(bundle.project.id))
    _safe("writer_session", lambda: project_service.get_writer_session(bundle.project.id))
    _safe("draft", lambda: project_service.get_draft(bundle.project.id))
    _safe("review_report", lambda: project_service.get_review_report(bundle.project.id))
    try:
        payload["revision_history"] = project_service.get_revision_history(bundle.project.id).to_dict()
    except Exception:  # noqa: BLE001
        payload["revision_history"] = None
    _safe("humanizer_session", lambda: project_service.get_humanizer_session(bundle.project.id))
    _safe("humanized_draft", lambda: project_service.get_humanized_draft(bundle.project.id))
    _safe("detection_session", lambda: project_service.get_detection_session(bundle.project.id))
    _safe("detection_report", lambda: project_service.get_detection_report(bundle.project.id))
    _safe("delivery_package", lambda: project_service.get_delivery_package(bundle.project.id))
    try:
        payload["chat_transcript"] = project_service.get_chat_transcript(bundle.project.id)
    except Exception:  # noqa: BLE001
        payload["chat_transcript"] = []
    return payload


def _truthy(form: Any, key: str) -> bool:
    val = form.get(key, "off")
    if val is True:
        return True
    if val is False:
        return False
    return str(val).lower() in {"on", "true", "1", "yes"}


def _int_clamped(raw: str, default: int = 0, lo: int = 0, hi: int = 72) -> int:
    try:
        v = int(raw)
    except (TypeError, ValueError):
        return default
    return max(lo, min(hi, v))


def document_has_visible_text(doc: Document) -> bool:
    """Reject completely empty inputs with a friendly message."""
    return document_has_visible_content(doc)


def parse_job(form) -> FormatJob:
    """Validate form fields and return a FormatJob."""
    font = form.get("font_family", "Times New Roman")
    if font not in ALLOWED_FONTS:
        font = "Times New Roman"

    size = _int_clamped(str(form.get("font_size", "12")), 12, 8, 24)
    if size not in ALLOWED_FONT_SIZES:
        size = 12

    ls_key = form.get("line_spacing", "1.5")
    line_spacing = LINE_SPACING_MAP.get(ls_key, 1.5)

    alignment = form.get("alignment", "left")
    if alignment not in ALIGNMENTS:
        alignment = "left"

    page_pos = (form.get("page_number_position") or "none").lower()
    if page_pos not in PAGE_POSITIONS:
        page_pos = "none"

    margin = form.get("margin_preset", "normal")
    if margin not in MARGIN_PRESETS:
        margin = "normal"

    before_pt = _int_clamped(str(form.get("space_before_pt", "0")), 0, 0, 72)
    after_pt = _int_clamped(str(form.get("space_after_pt", "0")), 0, 0, 72)

    format_style = (
        form.get("format_style")
        or form.get("style_preset")
        or form.get("citation_style")
        or "harvard"
    )

    return FormatJob(
        font_family=font,
        font_size_pt=size,
        line_spacing=line_spacing,
        alignment=alignment,
        first_line_indent=_truthy(form, "first_line_indent"),
        space_before_pt=before_pt,
        space_after_pt=after_pt,
        margin_preset=margin,
        page_number_position=page_pos,
        auto_headings=_truthy(form, "auto_headings"),
        heading_all_caps=_truthy(form, "heading_all_caps"),
        auto_justify_refs=_truthy(form, "auto_justify_refs"),
        requirement_headings=_truthy(form, "requirement_headings"),
        heading_size_pt=_int_clamped(str(form.get("heading_size_pt", "16")), 16, 12, 24),
        format_style=str(format_style),
    )


def parse_cover_page(form, *, fallback_paragraphs: list[str] | None = None) -> CoverPageData | None:
    """Build cover page data when the user enabled the title page toggle."""
    if not _truthy(form, "include_cover_page"):
        return None

    title = (form.get("cover_assignment_title") or "").strip()
    if not title and fallback_paragraphs:
        title = infer_assignment_title(fallback_paragraphs)

    cover = CoverPageData(
        assignment_title=title or "Assignment",
        student_name=(form.get("cover_student_name") or "").strip(),
        student_id=(form.get("cover_student_id") or "").strip(),
        university=(form.get("cover_university") or "").strip(),
        module=(form.get("cover_module") or "").strip(),
        lecturer=(form.get("cover_lecturer") or "").strip(),
        submission_date=(form.get("cover_submission_date") or "").strip(),
    )
    return cover


@app.context_processor
def inject_account():
    """Expose the current user and coin balance to every template."""
    user = economy_auth.current_user()
    balance = wallet.get_balance(user["id"]) if user else 0
    try:
        from services.economy.site_settings import get_current_humanizer_discount_status

        discount = get_current_humanizer_discount_status()
    except Exception:
        discount = {
            "active": False,
            "percent": 50,
            "source": "none",
        }
    return {
        "current_user": user,
        "coin_balance": balance,
        "welcome_bonus": WELCOME_BONUS,
        "is_admin": bool(user and user.get("is_admin")),
        "is_humanizer_discount_active": bool(discount.get("active"))
        and int(discount.get("percent") or 0) > 0,
        "humanizer_discount_percent": int(discount.get("percent") or 0),
    }


# ---------------------------------------------------------------- auth routes


@app.route("/register", methods=["GET", "POST"])
def register():
    existing = economy_auth.current_user()
    if existing:
        if not economy_auth.user_email_verified(existing):
            return redirect(url_for("verify_email_code"))
        return redirect(url_for("index"))
    ref_code = (request.values.get("ref") or request.values.get("referral_code") or "").strip()
    if request.method == "GET":
        if ref_code:
            session["pending_referral_code"] = ref_code
        return render_template(
            "register.html",
            nav_active=None,
            referral_code=ref_code or session.get("pending_referral_code") or "",
        )

    email = (request.form.get("email") or "").strip()
    name = (request.form.get("name") or "").strip()
    password = request.form.get("password") or ""
    password_confirm = (
        request.form.get("password_confirm") or request.form.get("confirm_password") or ""
    )
    fingerprint = (request.form.get("device_fingerprint") or "").strip()
    referral_code = (
        (request.form.get("referral_code") or request.form.get("ref") or "").strip()
        or session.get("pending_referral_code")
        or ""
    )
    if password != password_confirm:
        return render_template(
            "register.html",
            nav_active=None,
            error="Passwords do not match.",
            form={"email": email, "name": name},
            referral_code=referral_code,
        ), 400
    try:
        user = economy_auth.create_user(
            email,
            password,
            name=name,
            referral_code=referral_code or None,
            ip_address=economy_auth.client_ip_from_request(),
            device_fingerprint=fingerprint or None,
        )
    except DisposableEmailError as exc:
        return render_template(
            "register.html",
            nav_active=None,
            error=str(exc),
            form={"email": email, "name": name},
            referral_code=referral_code,
        ), 400
    except economy_auth.DuplicateEmail:
        return render_template(
            "register.html",
            nav_active=None,
            error="An account with this email already exists. Try signing in.",
            form={"email": email, "name": name},
            referral_code=referral_code,
        ), 409
    except economy_auth.AuthError as exc:
        return render_template(
            "register.html",
            nav_active=None,
            error=str(exc),
            form={"email": email, "name": name},
            referral_code=referral_code,
        ), 400

    session.pop("pending_referral_code", None)
    economy_auth.login_user(user["id"])
    # Always land on the notice page until the address is verified
    # (bootstrap admin is created already verified and may skip).
    if not economy_auth.user_email_verified(user):
        _send_user_verification_email(user, background=True)
        return redirect(url_for("verify_email_code"))
    return redirect(url_for("index"))


@app.route("/verify-email/code", methods=["GET", "POST"])
@app.route("/verify-email/notice", methods=["GET", "POST"])
@economy_auth.login_required
def verify_email_code():
    """Enter the 6-digit OTP sent by email (legacy /notice URL kept as alias)."""
    user = economy_auth.current_user()
    if user and economy_auth.user_email_verified(user):
        return redirect(url_for("workspace"))

    error = None
    resent = False

    if request.method == "POST":
        action = (request.form.get("action") or "verify").strip().lower()
        if action == "resend":
            if user:
                code = _send_user_verification_email(user)
                resent = True
                if code and (os.environ.get("EXPOSE_VERIFY_CODE") or "").strip().lower() in (
                    "1",
                    "true",
                    "yes",
                ):
                    # Local/debug only — never enable in production.
                    app.logger.info("EXPOSE_VERIFY_CODE user=%s code=%s", user.get("id"), code)
            return render_template(
                "verify_email.html",
                nav_active=None,
                email=(user or {}).get("email"),
                resent=True,
                otp_minutes=OTP_TTL_MINUTES,
            )

        submitted = (request.form.get("code") or request.form.get("verification_code") or "").strip()
        try:
            assert user is not None
            economy_auth.verify_email_otp(int(user["id"]), submitted)
            return redirect(url_for("workspace"))
        except economy_auth.AuthError as exc:
            error = str(exc)

    return render_template(
        "verify_email.html",
        nav_active=None,
        email=(user or {}).get("email"),
        error=error,
        resent=resent,
        otp_minutes=OTP_TTL_MINUTES,
    )


@app.post("/verify-email/resend")
@economy_auth.login_required
def verify_email_resend():
    """Compatibility endpoint — issues a new OTP and shows the code form."""
    user = economy_auth.current_user()
    if user and economy_auth.user_email_verified(user):
        return redirect(url_for("workspace"))
    if user:
        _send_user_verification_email(user)
    return redirect(url_for("verify_email_code", resent=1))


@app.get("/verify-email/<token>")
def verify_email_legacy_token(token: str):
    """Old magic-link URLs → OTP page (cross-device links no longer verify)."""
    _ = token
    user = economy_auth.current_user()
    if user and economy_auth.user_email_verified(user):
        return redirect(url_for("workspace"))
    if user:
        return redirect(url_for("verify_email_code"))
    return redirect(url_for("login", next=url_for("verify_email_code")))


@app.route("/login", methods=["GET", "POST"])
def login():
    next_url = request.values.get("next") or ""
    existing = economy_auth.current_user()
    if existing:
        if not economy_auth.user_email_verified(existing):
            return redirect(url_for("verify_email_code"))
        return redirect(next_url or url_for("index"))
    if request.method == "GET":
        return render_template("login.html", nav_active=None, next_url=next_url)

    email = (request.form.get("email") or "").strip()
    password = request.form.get("password") or ""
    user = economy_auth.verify_credentials(email, password)
    if user is None:
        return render_template(
            "login.html",
            nav_active=None,
            error="Incorrect email or password.",
            form={"email": email},
            next_url=next_url,
        ), 401

    economy_auth.login_user(user["id"])
    if not economy_auth.user_email_verified(user):
        return redirect(url_for("verify_email_code"))
    if next_url.startswith("/") and not next_url.startswith("//"):
        return redirect(next_url)
    return redirect(url_for("index"))


@app.post("/logout")
def logout():
    economy_auth.logout_user()
    return redirect(url_for("index"))


@app.route("/account", methods=["GET", "POST"])
@economy_auth.login_required
def account():
    profile_message = None
    profile_error = None
    password_message = None
    password_error = None
    user = economy_auth.current_user()

    if request.method == "POST":
        action = (request.form.get("action") or "").strip()
        if action == "profile":
            try:
                user = economy_auth.update_profile(
                    int(user["id"]),
                    name=request.form.get("name") or "",
                )
                profile_message = "Display name updated."
            except economy_auth.AuthError as exc:
                profile_error = str(exc)
        elif action == "password":
            current_password = request.form.get("current_password") or ""
            new_password = request.form.get("new_password") or ""
            confirm_password = request.form.get("confirm_password") or ""
            if new_password != confirm_password:
                password_error = "New password and confirmation do not match."
            else:
                try:
                    economy_auth.change_password(
                        int(user["id"]),
                        current_password=current_password,
                        new_password=new_password,
                    )
                    password_message = "Password updated."
                except economy_auth.AuthError as exc:
                    password_error = str(exc)

    return render_template(
        "account.html",
        nav_active=None,
        profile_message=profile_message,
        profile_error=profile_error,
        password_message=password_message,
        password_error=password_error,
    )


def _render_info(template: str, *, info_active: str | None = None):
    return render_template(template, nav_active=None, info_active=info_active)


@app.get("/account-info")
def account_info():
    return _render_info("info/account.html", info_active="account")


@app.get("/credits")
def credits():
    return _render_info("info/credits.html", info_active="credits")


@app.get("/about")
def about():
    return _render_info("info/about.html", info_active="about")


@app.get("/privacy")
def privacy():
    return _render_info("info/privacy.html", info_active="privacy")


@app.get("/terms")
def terms():
    return _render_info("info/terms.html", info_active="terms")


@app.get("/disclaimer")
def disclaimer():
    return _render_info("info/disclaimer.html", info_active="disclaimer")


@app.get("/payment-policy")
def payment_policy():
    return _render_info("info/payment.html", info_active="payment")


@app.get("/delivery-policy")
def delivery_policy():
    return _render_info("info/delivery.html", info_active="delivery")


@app.get("/refund-policy")
def refund_policy():
    return _render_info("info/refund.html", info_active="refund")


@app.get("/contact")
def contact():
    return _render_info("info/contact.html", info_active="contact")


@app.get("/faq")
def faq():
    return _render_info("info/faq.html", info_active="faq")


@app.get("/changelog")
def changelog():
    return _render_info("info/changelog.html", info_active="changelog")


# Legacy /legal/* URLs → dedicated pages
@app.get("/legal")
def legal():
    return redirect(url_for("about"), code=301)


@app.get("/legal/terms")
def legal_terms():
    return redirect(url_for("terms"), code=301)


@app.get("/legal/privacy")
def legal_privacy():
    return redirect(url_for("privacy"), code=301)


@app.get("/legal/refund")
def legal_refund():
    return redirect(url_for("refund_policy"), code=301)


# ---------------------------------------------------------------- admin


@app.route("/admin")
@economy_auth.admin_required
def admin_panel():
    return render_template("admin.html", nav_active="admin")


@app.get("/api/admin/daily-stats")
@economy_auth.admin_required
def api_admin_daily_stats():
    from services.economy.site_settings import get_admin_dashboard_stats

    try:
        payload = get_admin_dashboard_stats()
        return jsonify({"success": True, **payload})
    except Exception as exc:  # noqa: BLE001
        app.logger.exception("admin daily-stats failed")
        return jsonify(
            {
                "success": True,
                "today": {
                    "date": None,
                    "humanizer_requests_count": 0,
                    "humanizer_daily_limit": 50,
                    "humanizer_remaining": 50,
                    "turnitin_global_balance": 0,
                },
                "settings": {},
                "discount": {"active": False, "percent": 0, "source": "none"},
                "warning": str(exc),
            }
        )


@app.patch("/api/admin/site-settings")
@economy_auth.admin_required
def api_admin_site_settings():
    from services.economy.site_settings import (
        get_admin_dashboard_stats,
        update_site_settings,
    )

    body = request.get_json(silent=True) or {}
    kwargs: dict = {}
    if "is_humanizer_discount_active" in body:
        kwargs["is_humanizer_discount_active"] = bool(body["is_humanizer_discount_active"])
    if "humanizer_discount_percent" in body:
        try:
            kwargs["humanizer_discount_percent"] = int(body["humanizer_discount_percent"])
        except (TypeError, ValueError):
            return jsonify(
                {"success": False, "error": "humanizer_discount_percent must be an integer."}
            ), 400
    if "humanizer_daily_limit" in body:
        try:
            kwargs["humanizer_daily_limit"] = int(body["humanizer_daily_limit"])
        except (TypeError, ValueError):
            return jsonify(
                {"success": False, "error": "humanizer_daily_limit must be an integer."}
            ), 400
    if "turnitin_global_balance" in body:
        try:
            kwargs["turnitin_global_balance"] = int(body["turnitin_global_balance"])
        except (TypeError, ValueError):
            return jsonify(
                {"success": False, "error": "turnitin_global_balance must be an integer."}
            ), 400
    if "auto_discount_enabled" in body:
        kwargs["auto_discount_enabled"] = bool(body["auto_discount_enabled"])
    if "auto_discount_time" in body:
        kwargs["auto_discount_time"] = str(body["auto_discount_time"] or "20:00")
    if "auto_discount_min_remaining" in body:
        try:
            kwargs["auto_discount_min_remaining"] = int(body["auto_discount_min_remaining"])
        except (TypeError, ValueError):
            return jsonify(
                {
                    "success": False,
                    "error": "auto_discount_min_remaining must be an integer.",
                }
            ), 400
    if not kwargs:
        return jsonify({"success": False, "error": "No settings provided."}), 400
    try:
        settings = update_site_settings(**kwargs)
        dash = get_admin_dashboard_stats()
        return jsonify(
            {
                "success": True,
                "settings": settings,
                "today": dash.get("today"),
                "discount": dash.get("discount"),
            }
        )
    except Exception as exc:  # noqa: BLE001
        app.logger.exception("admin site-settings update failed")
        return jsonify({"success": False, "error": str(exc)}), 500


@app.get("/api/admin/users")
@economy_auth.admin_required
def api_admin_users():
    search = (request.args.get("q") or request.args.get("search") or "").strip()
    try:
        limit = int(request.args.get("limit", "100"))
    except ValueError:
        limit = 100
    try:
        offset = int(request.args.get("offset", "0"))
    except ValueError:
        offset = 0
    payload = admin_service.list_users(search=search, limit=limit, offset=offset)
    return jsonify({"success": True, **payload})


@app.patch("/api/admin/users/<int:user_id>/balance")
@economy_auth.admin_required
def api_admin_set_balance(user_id: int):
    """Set the user's wallet to an exact balance (absolute assign, not add)."""
    body = request.get_json(silent=True) or {}
    if "balance" not in body:
        return jsonify({"success": False, "error": "balance is required."}), 400

    raw = body["balance"]
    # Reject bool (int subclass), non-numeric types, and non-integer floats.
    if isinstance(raw, bool):
        return jsonify({"success": False, "error": "balance must be an integer."}), 400
    if isinstance(raw, int):
        new_balance = raw
    elif isinstance(raw, float):
        if not raw.is_integer():
            return jsonify({"success": False, "error": "balance must be an integer."}), 400
        new_balance = int(raw)
    elif isinstance(raw, str):
        s = raw.strip().replace(",", "").replace(" ", "")
        if not s or s.startswith("+"):
            # Disallow "+1500" to avoid ambiguous / concatenated payloads.
            return jsonify({"success": False, "error": "balance must be an integer."}), 400
        try:
            new_balance = int(s, 10)
        except ValueError:
            return jsonify({"success": False, "error": "balance must be an integer."}), 400
    else:
        return jsonify({"success": False, "error": "balance must be an integer."}), 400

    if new_balance < 0:
        return jsonify({"success": False, "error": "balance cannot be negative."}), 400

    reason = (body.get("reason") or "").strip() or None
    admin_id = economy_auth.current_user_id()
    try:
        result = admin_service.set_balance(
            user_id,
            new_balance,
            admin_id=int(admin_id),
            reason=reason,
        )
    except AdminError as exc:
        return jsonify({"success": False, "error": str(exc)}), 400
    return jsonify({"success": True, **result})


@app.patch("/api/admin/users/<int:user_id>/admin")
@economy_auth.admin_required
def api_admin_set_role(user_id: int):
    body = request.get_json(silent=True) or {}
    if "is_admin" not in body:
        return jsonify({"success": False, "error": "is_admin is required."}), 400
    is_admin_flag = bool(body["is_admin"])
    actor_id = economy_auth.current_user_id()
    try:
        result = admin_service.set_admin(
            user_id,
            is_admin=is_admin_flag,
            actor_id=int(actor_id),
        )
    except AdminError as exc:
        return jsonify({"success": False, "error": str(exc)}), 400
    return jsonify({"success": True, **result})


@app.delete("/api/admin/users/<int:user_id>")
@app.post("/api/admin/users/<int:user_id>/delete")
@economy_auth.admin_required
def api_admin_delete_user(user_id: int):
    actor_id = economy_auth.current_user_id()
    try:
        result = admin_service.delete_user(user_id, actor_id=int(actor_id))
    except AdminError as exc:
        return jsonify({"success": False, "error": str(exc)}), 400
    return jsonify({"success": True, **result})


@app.get("/api/admin/users/<int:user_id>/ledger")
@economy_auth.admin_required
def api_admin_user_ledger(user_id: int):
    """CreditTransaction journal for one user (admin audit view)."""
    try:
        limit = int(request.args.get("limit", "100"))
    except ValueError:
        limit = 100
    try:
        offset = int(request.args.get("offset", "0"))
    except ValueError:
        offset = 0
    try:
        payload = admin_service.get_ledger(user_id, limit=limit, offset=offset)
    except AdminError as exc:
        return jsonify({"success": False, "error": str(exc)}), 404
    return jsonify({"success": True, **payload})


@app.get("/api/admin/users/<int:user_id>/purchases")
@economy_auth.admin_required
def api_admin_user_purchases(user_id: int):
    """PaddlePurchase rows for one user."""
    try:
        limit = int(request.args.get("limit", "100"))
    except ValueError:
        limit = 100
    try:
        offset = int(request.args.get("offset", "0"))
    except ValueError:
        offset = 0
    try:
        payload = admin_service.get_purchases(user_id, limit=limit, offset=offset)
    except AdminError as exc:
        return jsonify({"success": False, "error": str(exc)}), 404
    return jsonify({"success": True, **payload})


@app.get("/api/admin/purchases")
@economy_auth.admin_required
def api_admin_purchases():
    """All PaddlePurchase rows (optional search / status filter)."""
    search = (request.args.get("q") or request.args.get("search") or "").strip()
    status = (request.args.get("status") or "").strip()
    try:
        limit = int(request.args.get("limit", "100"))
    except ValueError:
        limit = 100
    try:
        offset = int(request.args.get("offset", "0"))
    except ValueError:
        offset = 0
    try:
        payload = admin_service.list_purchases(
            search=search, status=status, limit=limit, offset=offset
        )
    except AdminError as exc:
        return jsonify({"success": False, "error": str(exc)}), 400
    return jsonify({"success": True, **payload})


@app.get("/api/admin/users/<int:user_id>/usage")
@economy_auth.admin_required
def api_admin_user_usage(user_id: int):
    """Usage events (AI launches) for one user."""
    try:
        limit = int(request.args.get("limit", "100"))
    except ValueError:
        limit = 100
    try:
        offset = int(request.args.get("offset", "0"))
    except ValueError:
        offset = 0
    try:
        payload = admin_service.get_usage(user_id, limit=limit, offset=offset)
    except AdminError as exc:
        return jsonify({"success": False, "error": str(exc)}), 404
    return jsonify({"success": True, **payload})


@app.get("/api/admin/usage")
@economy_auth.admin_required
def api_admin_usage():
    """All usage events (optional search / feature filter)."""
    search = (request.args.get("q") or request.args.get("search") or "").strip()
    feature = (request.args.get("feature") or "").strip()
    try:
        limit = int(request.args.get("limit", "100"))
    except ValueError:
        limit = 100
    try:
        offset = int(request.args.get("offset", "0"))
    except ValueError:
        offset = 0
    payload = admin_service.list_usage(
        search=search, feature=feature, limit=limit, offset=offset
    )
    return jsonify({"success": True, **payload})


@app.get("/api/admin/analytics")
@economy_auth.admin_required
def api_admin_analytics():
    """Aggregate KPIs: sold/used credits, revenue, top customers/countries."""
    try:
        top_limit = int(request.args.get("top", "10"))
    except ValueError:
        top_limit = 10
    payload = admin_service.get_analytics(top_limit=top_limit)
    return jsonify({"success": True, **payload})


@app.get("/api/admin/dataset-stats")
@economy_auth.admin_required
def api_admin_dataset_stats():
    """Humanizer + AI-detection ML dataset collection counters."""
    from services.dataset_logger import get_dataset_recent_samples, get_dataset_stats

    try:
        stats = get_dataset_stats()
        samples = get_dataset_recent_samples(limit=10)
    except Exception as exc:  # noqa: BLE001
        app.logger.exception("dataset-stats failed")
        return jsonify({"success": False, "error": str(exc)}), 500
    return jsonify({"success": True, **stats, **samples})


@app.route("/admin/dataset-stats")
@economy_auth.admin_required
def admin_dataset_stats_page():
    """Dedicated admin view for humanizer + detector dataset progress."""
    return render_template("admin_dataset_stats.html", nav_active="admin")


# ------------------------------------------------- JSON auth (register modal)


def _auth_success_payload(user: dict) -> dict:
    return {
        "success": True,
        "user": {
            "id": user["id"],
            "email": user["email"],
            "name": user.get("name"),
            "is_verified": bool(user.get("is_verified")),
            "avatar_file": user.get("avatar_file"),
        },
        "balance": wallet.get_balance(user["id"]),
        "welcome_bonus_granted": bool(user.get("welcome_bonus_granted")),
    }


@app.post("/api/auth/register")
def api_auth_register():
    payload = request.get_json(silent=True) or {}
    email = str(payload.get("email") or "").strip()
    name = str(payload.get("name") or "").strip()
    password = str(payload.get("password") or "")
    password_confirm = str(
        payload.get("password_confirm") or payload.get("confirm_password") or ""
    )
    fingerprint = str(payload.get("device_fingerprint") or "").strip()
    referral_code = str(
        payload.get("referral_code") or payload.get("ref") or session.get("pending_referral_code") or ""
    ).strip()
    if password != password_confirm:
        return jsonify({"success": False, "error": "Passwords do not match."}), 400
    try:
        user = economy_auth.create_user(
            email,
            password,
            name=name,
            referral_code=referral_code or None,
            ip_address=economy_auth.client_ip_from_request(),
            device_fingerprint=fingerprint or None,
        )
    except DisposableEmailError as exc:
        return jsonify({"success": False, "error": str(exc)}), 400
    except economy_auth.DuplicateEmail as exc:
        return jsonify({"success": False, "error": str(exc)}), 409
    except economy_auth.AuthError as exc:
        return jsonify({"success": False, "error": str(exc)}), 400
    session.pop("pending_referral_code", None)
    economy_auth.login_user(user["id"])
    verify_code = None
    needs_verify = not economy_auth.user_email_verified(user)
    if needs_verify:
        verify_code = _send_user_verification_email(user, background=True)
    payload_out = _auth_success_payload(user)
    payload_out["email_verification_required"] = needs_verify
    if verify_code and (os.environ.get("EXPOSE_VERIFY_CODE") or "").strip().lower() in (
        "1",
        "true",
        "yes",
    ):
        payload_out["verify_code"] = verify_code
    return jsonify(payload_out)


@app.post("/api/auth/login")
def api_auth_login():
    payload = request.get_json(silent=True) or {}
    email = str(payload.get("email") or "").strip()
    password = str(payload.get("password") or "")
    user = economy_auth.verify_credentials(email, password)
    if user is None:
        return jsonify({"success": False, "error": "Incorrect email or password."}), 401
    economy_auth.login_user(user["id"])
    payload_out = _auth_success_payload(user)
    payload_out["email_verification_required"] = not economy_auth.user_email_verified(user)
    return jsonify(payload_out)


@app.post("/api/auth/logout")
def api_auth_logout():
    economy_auth.logout_user()
    return jsonify({"success": True})


@app.post("/api/upload_avatar")
@economy_auth.login_required
def api_upload_avatar():
    """Secure avatar upload — 2 MB, png/jpg/jpeg/webp only."""
    user = economy_auth.current_user()
    assert user is not None
    upload = request.files.get("avatar") or request.files.get("file")
    content_length = request.content_length or 0
    if content_length and content_length > _AVATAR_MAX_BYTES + 4096:
        return jsonify({"success": False, "error": "Avatar must be 2 MB or smaller."}), 413
    try:
        relative = validate_and_store_avatar(
            upload,
            user_id=int(user["id"]),
            repo_root=_REPO_ROOT,
        )
        updated = economy_auth.set_avatar_file(int(user["id"]), relative)
    except AvatarUploadError as exc:
        return jsonify({"success": False, "error": str(exc)}), 400
    except economy_auth.AuthError as exc:
        return jsonify({"success": False, "error": str(exc)}), 400
    return jsonify(
        {
            "success": True,
            "avatar_file": updated.get("avatar_file"),
            "avatar_url": url_for("static", filename=relative),
        }
    )


# ------------------------------------------------------------- economy API


@app.get("/api/economy/balance")
def api_economy_balance():
    user = economy_auth.current_user()
    if user is None:
        return jsonify({"authenticated": False, "balance": 0})
    return jsonify(
        {"authenticated": True, "balance": wallet.get_balance(user["id"])}
    )


@app.post("/api/economy/quote")
def api_economy_quote():
    payload = request.get_json(silent=True) or {}
    feature = str(payload.get("feature") or "").strip().lower()
    params = payload.get("params") if isinstance(payload.get("params"), dict) else {}
    try:
        cost = feature_cost(feature, **params)
    except KeyError:
        return jsonify({"error": f"Unknown feature: {feature}"}), 400
    user = economy_auth.current_user()
    balance = wallet.get_balance(user["id"]) if user else 0
    return jsonify(
        {
            "feature": feature,
            "label": FEATURE_LABELS.get(feature, feature),
            "cost": cost,
            "balance": balance,
            "authenticated": user is not None,
            "affordable": (user is not None) and balance >= cost,
        }
    )


@app.get("/api/economy/transactions")
@economy_auth.login_required
def api_economy_transactions():
    user = economy_auth.current_user()
    try:
        limit = int(request.args.get("limit") or 50)
    except (TypeError, ValueError):
        limit = 50
    return jsonify(
        {
            "balance": wallet.get_balance(user["id"]),
            "transactions": wallet.history(user["id"], limit=limit),
        }
    )


@app.get("/api/economy/packages")
def api_economy_packages():
    return jsonify(
        {
            "packages": list(TOPUP_PACKAGES.values()),
            "paddle_configured": paddle_configured(),
            "paddle_environment": paddle_environment(),
            "client_token": paddle_client_token() or None,
            "cryptomus_configured": cryptomus_configured(),
            "payment_provider": "cryptomus" if cryptomus_configured() else (
                "paddle" if paddle_configured() else None
            ),
        }
    )


@app.post("/api/economy/checkout")
@economy_auth.login_required
def api_economy_checkout():
    """Create a Paddle Checkout for a credit package. Returns checkout_url / txn id."""
    user = economy_auth.current_user()
    payload = request.get_json(silent=True) or {}
    package_id = str(payload.get("package") or "").strip()
    if not package_id:
        return jsonify({"error": "package is required"}), 400
    if economy_package(package_id) is None:
        return jsonify({"error": "Unknown package"}), 400
    if not paddle_configured():
        return (
            jsonify(
                {
                    "error": "PADDLE_NOT_CONFIGURED",
                    "message": "Paddle API key is not set. Add PADDLE_API_KEY and price ids.",
                }
            ),
            503,
        )
    try:
        checkout = create_checkout(
            user_id=int(user["id"]),
            package_id=package_id,
            customer_email=str(user.get("email") or "") or None,
        )
    except PaddleGatewayError as exc:
        return jsonify({"error": str(exc)}), 400
    return jsonify({"success": True, **checkout})


@app.post("/api/webhooks/paddle")
def api_paddle_webhook():
    """Paddle Billing webhook — signature verified, fulfillment idempotent."""
    raw = request.get_data(cache=True, as_text=False)
    signature = request.headers.get("Paddle-Signature")
    try:
        verify_paddle_signature(raw, signature)
    except PaddleSignatureError as exc:
        app.logger.warning("paddle webhook signature failed: %s", exc)
        return jsonify({"error": "invalid_signature"}), 400

    try:
        event = json.loads(raw.decode("utf-8"))
    except (UnicodeDecodeError, ValueError):
        return jsonify({"error": "invalid_json"}), 400

    try:
        result = handle_webhook_event(event if isinstance(event, dict) else {})
    except PaddleGatewayError as exc:
        app.logger.exception("paddle webhook fulfill failed: %s", exc)
        # 500 so Paddle retries; idempotency protects double-credit on success path.
        return jsonify({"error": str(exc)}), 500

    return jsonify({"ok": True, **result}), 200


@app.post("/api/webhooks/lemon-squeezy")
def api_lemon_squeezy_webhook():
    """Lemon Squeezy webhook — HMAC X-Signature verified, order_created → coins.

    Configure the store webhook URL to this endpoint and set
    ``LEMON_SQUEEZY_WEBHOOK_SECRET`` to the signing secret from Lemon.
    Pass ``user_id`` in checkout ``custom_data`` so credits land on the right account.
    """
    raw = request.get_data(cache=True, as_text=False)
    signature = request.headers.get("X-Signature")
    try:
        verify_lemon_squeezy_signature(raw, signature)
    except LemonSqueezySignatureError as exc:
        app.logger.warning("lemon-squeezy webhook signature failed: %s", exc)
        return jsonify({"error": "invalid_signature", "message": str(exc)}), 403

    try:
        payload = json.loads(raw.decode("utf-8"))
    except (UnicodeDecodeError, ValueError):
        app.logger.warning("lemon-squeezy webhook: invalid JSON body")
        return jsonify({"error": "invalid_json"}), 400

    if not isinstance(payload, dict):
        return jsonify({"error": "invalid_payload"}), 400

    try:
        result = lemon_squeezy_handle_webhook(payload)
    except LemonSqueezyGatewayError as exc:
        app.logger.warning("lemon-squeezy webhook rejected: %s", exc)
        # 400 for bad custom_data / unknown variant — Lemon should not infinite-retry.
        return jsonify({"error": str(exc)}), 400
    except Exception as exc:  # noqa: BLE001
        app.logger.exception("lemon-squeezy webhook fulfill failed: %s", exc)
        # 500 so Lemon retries transient DB/outage failures (idempotent on success).
        return jsonify({"error": "fulfillment_failed"}), 500

    # Spec: successful credit path returns {"status": "success"}.
    if result.get("status") == "success":
        return jsonify({"status": "success", **{k: v for k, v in result.items() if k != "status"}}), 200
    return jsonify(result), 200


@app.post("/api/payments/create")
@economy_auth.login_required
def api_payments_create():
    """Create a Cryptomus invoice for a credit package. Returns payment_url.

    Only ``package`` (id) is accepted from the client. Amount and credits are
    resolved exclusively from the server-side TOPUP catalog.
    """
    user = economy_auth.current_user()
    payload = request.get_json(silent=True) or {}
    # Ignore any client-supplied amount/credits/usd — never trusted.
    package_id = str(payload.get("package") or payload.get("package_id") or "").strip()
    if not package_id:
        return jsonify({"error": "package is required"}), 400
    if economy_package(package_id) is None:
        return jsonify({"error": "Unknown package"}), 400
    if not cryptomus_configured():
        return (
            jsonify(
                {
                    "error": "CRYPTOMUS_NOT_CONFIGURED",
                    "message": "Cryptomus payments are not available right now.",
                }
            ),
            503,
        )
    try:
        invoice = cryptomus_create_invoice(
            user_id=int(user["id"]),
            package_id=package_id,
        )
    except CryptomusGatewayError as exc:
        app.logger.warning("cryptomus create invoice failed: %s", exc)
        return jsonify({"error": str(exc)}), 400
    return jsonify({"success": True, **invoice})


@app.post("/api/payments/cryptomus/webhook")
def api_cryptomus_webhook():
    """Cryptomus payment webhook — signature verified, fulfillment idempotent.

    Reads raw body first (official PHP: file_get_contents('php://input')).
    """
    raw = request.get_data(cache=True, as_text=False)
    try:
        payload = cryptomus_parse_webhook_payload(raw)
    except CryptomusGatewayError:
        return jsonify({"error": "invalid_json"}), 400

    # Prefer original client IP behind nginx (docs IP allowlist).
    remote = (
        (request.headers.get("X-Real-IP") or "").strip()
        or (request.headers.get("X-Forwarded-For") or "").split(",")[0].strip()
        or request.remote_addr
    )

    try:
        result = cryptomus_handle_webhook(payload, remote_addr=remote)
    except CryptomusSignatureError as exc:
        app.logger.warning("cryptomus webhook rejected: %s", exc)
        return jsonify({"error": "invalid_signature"}), 400
    except CryptomusGatewayError as exc:
        app.logger.exception("cryptomus webhook fulfill failed: %s", exc)
        msg = str(exc)
        if "Unknown order_id" in msg:
            # Signed but unknown — acknowledge so Cryptomus stops retrying forever.
            return jsonify({"error": msg}), 200
        return jsonify({"error": msg}), 500

    # Cryptomus expects HTTP 200 on successful receipt.
    return jsonify({"ok": True, **result}), 200


@app.post("/api/webhooks/gumroad")
def api_gumroad_webhook():
    """Gumroad Ping — form-urlencoded sale notification → credit top-up.

    Always returns plain ``200 OK`` so Gumroad stops retrying. Fulfillment is
    idempotent on ``sale_id``. Pass ``user_id`` as a checkout URL param
    (``?user_id=123``) so credits land on the right account.
    """
    # Gumroad emergency fallback — keep disabled (do not delete body below).
    return jsonify({"error": "Payment method temporarily disabled"}), 403

    form = request.form
    try:
        result = gumroad_handle_ping(form)
        app.logger.info("gumroad ping ok: %s", result)
    except GumroadGatewayError as exc:
        app.logger.warning("gumroad ping ignored/failed: %s", exc)
    except Exception as exc:  # noqa: BLE001
        app.logger.exception("gumroad ping unexpected error: %s", exc)
    return "OK", 200

@app.post("/api/economy/topup")
@economy_auth.login_required
def api_economy_topup():
    """Dev-only mock top-up. Disabled in live mode; sandbox needs PADDLE_ALLOW_MOCK_TOPUP=1."""
    import uuid

    if not mock_topup_allowed():
        return (
            jsonify(
                {
                    "error": "MOCK_TOPUP_DISABLED",
                    "message": "Use /api/economy/checkout for real Paddle payments.",
                }
            ),
            403,
        )

    user = economy_auth.current_user()
    payload = request.get_json(silent=True) or {}
    pkg = economy_package(str(payload.get("package") or ""))
    if pkg is None:
        return jsonify({"error": "Unknown package"}), 400

    mock_txn = f"mock_{uuid.uuid4().hex[:16]}"
    try:
        result = apply_paid_purchase_atomic(
            user_id=int(user["id"]),
            paddle_transaction_id=mock_txn,
            product_id=str(pkg["id"]),
            price_id=f"price_{pkg['id']}",
            credits=int(pkg["coins"]),
            amount=float(pkg["usd"]),
            currency="USD",
            meta={"mock": True, "package": pkg["id"], "usd": pkg["usd"]},
        )
    except PaddleGatewayError as exc:
        return jsonify({"error": str(exc)}), 400

    return jsonify(
        {
            "success": True,
            "coins_added": int(result.get("credits_added") or pkg["coins"]),
            "balance": result.get("balance"),
            "package": pkg,
            "purchase": result.get("purchase"),
        }
    )


@app.route("/")
def index():
    """Format tab. V2 UI when ``FORMATTER_V2_ENABLED``; otherwise the legacy page."""
    if formatter_v2_enabled():
        return render_template("format_v2.html", nav_active="home")
    return render_template("index.html", nav_active="home")


@app.route("/check")
def check():
    """Academic Check UI is temporarily gated — keep templates/static intact for later."""
    return render_template(
        "soon.html",
        nav_active="check",
        feature="Academic Check",
    )


@app.route("/templates")
def templates():
    return redirect(url_for("index"), code=302)


@app.route("/references")
def references():
    return redirect(url_for("index"), code=302)


@app.route("/workspace")
@economy_auth.email_verified_required
def workspace():
    """Full document workspace — editor, humanize, AI, cite, comments."""
    return render_template("workspace.html", nav_active="workspace")


@app.route("/presentation")
def presentation():
    """Presentation UI is temporarily gated."""
    return render_template(
        "soon.html",
        nav_active="presentation",
        feature="Presentation",
    )


@app.route("/earn")
@economy_auth.login_required
def earn_share():
    return render_template("earn.html", nav_active="earn")


@app.get("/api/referral/me")
@economy_auth.login_required
def api_referral_me():
    from services.economy.referral import ReferralError, get_referral_profile

    user_id = economy_auth.current_user_id()
    try:
        profile = get_referral_profile(int(user_id))
    except ReferralError as exc:
        return jsonify({"success": False, "error": str(exc)}), 400
    link = url_for("register", ref=profile["referral_code"], _external=True)
    profile["referral_link"] = link
    return jsonify({"success": True, **profile})


@app.post("/api/referral/convert")
@economy_auth.login_required
def api_referral_convert():
    from services.economy.referral import ReferralError, convert_balance_to_credits

    payload = request.get_json(silent=True) or {}
    try:
        amount = float(payload.get("amount_usd") or payload.get("amount") or 0)
    except (TypeError, ValueError):
        return jsonify({"success": False, "error": "amount_usd must be a number."}), 400
    try:
        result = convert_balance_to_credits(int(economy_auth.current_user_id()), amount)
    except ReferralError as exc:
        return jsonify({"success": False, "error": str(exc)}), 400
    return jsonify({"success": True, **result})


@app.post("/api/referral/withdraw")
@economy_auth.login_required
def api_referral_withdraw():
    from services.economy.referral import ReferralError, create_withdrawal

    payload = request.get_json(silent=True) or {}
    try:
        amount = float(payload.get("amount_usd") or payload.get("amount") or 0)
    except (TypeError, ValueError):
        return jsonify({"success": False, "error": "amount_usd must be a number."}), 400
    wallet_details = str(payload.get("wallet_details") or payload.get("wallet") or "")
    try:
        result = create_withdrawal(
            int(economy_auth.current_user_id()),
            amount_usd=amount,
            wallet_details=wallet_details,
        )
    except ReferralError as exc:
        return jsonify({"success": False, "error": str(exc)}), 400
    return jsonify({"success": True, **result})


@app.get("/api/admin/withdrawals")
@economy_auth.admin_required
def api_admin_withdrawals():
    from services.economy.referral import list_withdrawals

    status = (request.args.get("status") or "pending").strip() or None
    if status == "all":
        status = None
    try:
        limit = int(request.args.get("limit", "100"))
    except ValueError:
        limit = 100
    try:
        offset = int(request.args.get("offset", "0"))
    except ValueError:
        offset = 0
    payload = list_withdrawals(status=status, limit=limit, offset=offset)
    return jsonify({"success": True, **payload})


@app.post("/api/admin/withdrawals/<int:request_id>/approve")
@economy_auth.admin_required
def api_admin_withdrawal_approve(request_id: int):
    from services.economy.referral import ReferralError, resolve_withdrawal

    payload = request.get_json(silent=True) or {}
    try:
        result = resolve_withdrawal(
            request_id,
            approve=True,
            admin_id=economy_auth.current_user_id(),
            note=payload.get("note"),
        )
    except ReferralError as exc:
        return jsonify({"success": False, "error": str(exc)}), 400
    return jsonify({"success": True, **result})


@app.post("/api/admin/withdrawals/<int:request_id>/reject")
@economy_auth.admin_required
def api_admin_withdrawal_reject(request_id: int):
    from services.economy.referral import ReferralError, resolve_withdrawal

    payload = request.get_json(silent=True) or {}
    try:
        result = resolve_withdrawal(
            request_id,
            approve=False,
            admin_id=economy_auth.current_user_id(),
            note=payload.get("note"),
        )
    except ReferralError as exc:
        return jsonify({"success": False, "error": str(exc)}), 400
    return jsonify({"success": True, **result})


@app.route("/editor")
def editor():
    """Legacy URL → workspace."""
    return redirect(url_for("workspace"), code=302)


@app.route("/humanizer")
@economy_auth.email_verified_required
def humanizer():
    from services.economy.pricing import FEATURE_COSTS
    from services.economy.site_settings import humanize_credit_cost

    pricing = humanize_credit_cost(FEATURE_COSTS["humanize"])
    base_cost = int(pricing["original_price"])
    discount_percent = int(pricing["discount_percent"] or 0)
    is_discount_active = bool(pricing["discount_active"]) and discount_percent > 0
    discounted_cost = int(pricing["charged"])
    return render_template(
        "humanizer.html",
        nav_active="humanizer",
        humanize_cost=discounted_cost if is_discount_active else base_cost,
        base_cost=base_cost,
        is_discount_active=is_discount_active,
        discount_percent=discount_percent,
        discounted_cost=discounted_cost,
    )


@app.route("/assignment")
@economy_auth.email_verified_required
def assignment():
    return render_template("assignment.html", nav_active="assignment")


@app.route("/assignments")
def assignments_history():
    """Legacy URL — history lives on the Assignment page panel."""
    return redirect(url_for("assignment"))


def _parse_pipeline_stage(value: str) -> PipelineStage | None:
    try:
        return PipelineStage(value.strip().lower())
    except ValueError:
        return None


@app.get("/api/assignment/pipeline/stages")
def api_assignment_pipeline_stages():
    """Pipeline stage definitions and future provider integration slots."""
    return jsonify({"stages": [spec.to_dict() for spec in PIPELINE_STAGE_SPECS]})


@app.get("/api/assignment/projects")
@economy_auth.login_required
def api_assignment_projects_list():
    user_id = economy_auth.current_user_id()
    projects = project_service.list_projects_for_user(str(user_id))
    return jsonify({"projects": projects})


@app.post("/api/assignment/projects")
@economy_auth.login_required
def api_assignment_project_create():
    """Create an assignment project with files, requirement shell, and pipeline."""
    payload = request.get_json(silent=True) or {}
    upload_manifest = payload.get("upload_manifest")
    if upload_manifest is not None and not isinstance(upload_manifest, dict):
        return jsonify({"error": "upload_manifest must be an object"}), 400
    files = payload.get("files")
    if files is not None and not isinstance(files, list):
        return jsonify({"error": "files must be an array"}), 400

    session_uid = economy_auth.current_user_id()
    # Never trust client-supplied user_id (IDOR).
    user_id = str(session_uid)

    try:
        bundle = project_service.create_project(
            user_id=user_id,
            title=payload.get("title"),
            university=payload.get("university"),
            deadline=payload.get("deadline"),
            note=payload.get("note"),
            files=files,
            upload_manifest=upload_manifest,
        )
    except ValueError as exc:
        return jsonify({"error": str(exc)}), 400

    return jsonify(_project_api_payload(bundle)), 201


@app.get("/api/assignment/projects/<project_id>")
@economy_auth.login_required
def api_assignment_project_get(project_id: str):
    bundle, err = _load_owned_project(project_id)
    if err is not None:
        return err
    try:
        return jsonify(_project_api_payload(bundle))
    except Exception as exc:  # noqa: BLE001
        trace(
            "api.project_get.error",
            project_id=project_id,
            error=str(exc),
            error_type=type(exc).__name__,
        )
        return jsonify({"error": "Failed to load project state"}), 500


@app.post("/api/assignment/projects/<project_id>/files")
@economy_auth.login_required
def api_assignment_project_add_file(project_id: str):
    denied = _require_project_owner_response(project_id)
    if denied is not None:
        return denied
    payload = request.get_json(silent=True) or {}
    original_filename = (payload.get("original_filename") or payload.get("name") or "").strip()
    file_type = (payload.get("file_type") or payload.get("source") or "").strip()
    if not original_filename or not file_type:
        return jsonify({"error": "file_type and original_filename are required"}), 400
    try:
        file_record = project_service.add_file(
            project_id,
            file_type=file_type,
            original_filename=original_filename,
            storage_path=payload.get("storage_path"),
            parsed=bool(payload.get("parsed", False)),
        )
    except KeyError:
        return jsonify({"error": "Project not found"}), 404
    except ValueError as exc:
        return jsonify({"error": str(exc)}), 400
    return jsonify(file_record.to_dict()), 201


@app.post("/api/assignment/projects/upload")
@economy_auth.login_required
def api_assignment_project_upload():
    """Create a project and persist uploaded assignment files."""
    note = (request.form.get("note") or request.form.get("lecture_notes") or "").strip()
    deadline = (request.form.get("deadline") or "").strip() or None
    title = (request.form.get("title") or "").strip() or "Assignment Project"

    has_upload = any(
        upload and upload.filename
        for field in ("assignment_brief", "rubric", "additional_files", "lecture_notes", "files")
        for upload in request.files.getlist(field)
    )
    if not has_upload and not note:
        return jsonify({"error": "Upload at least one file or add a note."}), 400

    trace(
        "api.upload.received",
        storage_root=str(PROJECT_STORAGE_ROOT),
        store_root=str(project_service.store.storage_root),
    )
    try:
        session_uid = economy_auth.current_user_id()
        bundle = project_service.create_project(
            user_id=str(session_uid),
            title=title,
            deadline=deadline,
            note=note or None,
        )
        project_id = bundle.project.id
        bundle_path = project_service.store._bundle_path(project_id)
        trace(
            "api.upload.created",
            project_id=project_id,
            bundle_path=str(bundle_path.resolve()),
            bundle_exists=bundle_path.is_file(),
            save_ok=bundle_path.is_file(),
        )
        saved_files = _attach_multipart_uploads(project_id)
        if note:
            note_path = _write_debug_text_file(bundle.project.id, "project_note.txt", note)
            project_service.add_file(
                bundle.project.id,
                file_type="professor_notes",
                original_filename="project_note.txt",
                storage_path=note_path,
                parsed=True,
            )
    except ValueError as exc:
        trace("api.upload.failed", error=str(exc))
        return jsonify({"error": str(exc)}), 400

    bundle_path = project_service.store._bundle_path(bundle.project.id)
    trace(
        "api.upload.completed",
        project_id=bundle.project.id,
        bundle_path=str(bundle_path.resolve()),
        bundle_exists=bundle_path.is_file(),
        uploaded_files=len(saved_files),
    )
    payload = _project_api_payload(bundle)
    payload["uploaded_files"] = saved_files
    return jsonify(payload), 201


@app.post("/api/assignment/projects/<project_id>/analyze-requirements")
def api_assignment_project_analyze_requirements(project_id: str):
    """Run requirement analysis via Gemini.

    Uploading a brief and filling the form is free/anonymous; running the
    analysis (which leads to pricing) requires an account.
    """
    if economy_auth.current_user() is None:
        return (
            jsonify(
                {
                    "success": False,
                    "error": "REGISTER_REQUIRED",
                    "message": "Create a free account to analyze and price your assignment.",
                }
            ),
            403,
        )
    trace(
        "api.analyze.received",
        **project_service.store.lookup_diagnostics(project_id),
    )
    try:
        bundle = project_service.analyze_requirements(project_id)
    except KeyError as exc:
        return _assignment_not_found("analyze-requirements", project_id, exc)
    except ValueError as exc:
        trace("api.analyze.failed", project_id=project_id, error=str(exc))
        return jsonify({"error": str(exc)}), 502
    except Exception as exc:  # noqa: BLE001
        trace("api.analyze.error", project_id=project_id, error=str(exc), error_type=type(exc).__name__)
        raise
    trace(
        "api.analyze.completed",
        project_id=project_id,
        price=bundle.project.price,
        assignment_type=bundle.requirement.assignment_type,
    )
    return jsonify(_project_api_payload(bundle))


@app.post("/api/assignment/projects/<project_id>/pricing")
def api_assignment_project_pricing(project_id: str):
    payload = request.get_json(silent=True) or {}
    priority = str(payload.get("priority") or "standard").strip().lower()
    trace(
        "api.pricing.received",
        priority=priority,
        **project_service.store.lookup_diagnostics(project_id),
    )
    try:
        bundle = project_service.calculate_pricing(project_id, priority=priority)
    except KeyError as exc:
        return _assignment_not_found("pricing", project_id, exc)
    except ValueError as exc:
        trace("api.pricing.failed", project_id=project_id, error=str(exc))
        return jsonify({"error": str(exc)}), 400
    except Exception as exc:  # noqa: BLE001
        trace("api.pricing.error", project_id=project_id, error=str(exc), error_type=type(exc).__name__)
        raise
    trace(
        "api.pricing.completed",
        project_id=project_id,
        price=bundle.project.price,
        bundle_path=str(project_service.store._bundle_path(project_id).resolve()),
    )
    return jsonify(_project_api_payload(bundle))


@app.post("/api/assignment/projects/<project_id>/confirm-payment")
@economy_auth.login_required
def api_assignment_project_confirm_payment(project_id: str):
    trace(
        "api.confirm_payment.received",
        **project_service.store.lookup_diagnostics(project_id),
    )
    user_id = economy_auth.current_user_id()

    # Determine coin cost before confirming (skip charge if already paid).
    already_paid = False
    coins = 0
    try:
        existing = project_service.store.require_bundle(project_id)
        already_paid = bool(existing.project.artifacts.get("payment_confirmed"))
        price_usd = existing.project.price
        if price_usd is None:
            pricing_artifact = existing.project.artifacts.get("pricing") or {}
            price_usd = pricing_artifact.get("amount_usd")
        if price_usd is not None:
            coins = max(1, int(round(float(price_usd) * USD_TO_COINS)))
    except KeyError as exc:
        return _assignment_not_found("confirm-payment", project_id, exc)

    charged_here = False
    if not already_paid and coins > 0:
        from services.economy.referral import apply_pro_discount

        user = economy_auth.current_user()
        coins = apply_pro_discount(user, coins)
        try:
            wallet.debit(user_id, coins, "assignment", ref_id=project_id)
            charged_here = True
        except InsufficientCoins as exc:
            return (
                jsonify(
                    {
                        "error": "INSUFFICIENT_COINS",
                        "message": f"Not enough coins. This project costs {exc.required}; you have {exc.balance}.",
                        "required": exc.required,
                        "balance": exc.balance,
                    }
                ),
                402,
            )

    try:
        bundle = project_service.confirm_payment(project_id)
    except KeyError as exc:
        if charged_here:
            _refund_safe(user_id, coins, "assignment", ref_id=project_id)
        return _assignment_not_found("confirm-payment", project_id, exc)
    except ValueError as exc:
        if charged_here:
            _refund_safe(user_id, coins, "assignment", ref_id=project_id)
        trace(
            "api.confirm_payment.failed",
            error=str(exc),
            **project_service.store.lookup_diagnostics(project_id),
        )
        return jsonify({"error": str(exc)}), 400
    except Exception:
        if charged_here:
            _refund_safe(user_id, coins, "assignment", ref_id=project_id)
        raise
    trace(
        "api.confirm_payment.completed",
        project_id=project_id,
        price=bundle.project.price,
        payment_confirmed=bool(bundle.project.artifacts.get("payment_confirmed")),
    )
    if charged_here and coins > 0:
        _record_usage_safe(
            user_id,
            feature=FEATURE_ASSIGNMENT,
            credits_used=coins,
            provider="AssignmentPipeline",
            request_id=project_id,
        )
    payload = _project_api_payload(bundle)
    if isinstance(payload, dict):
        payload["coins_charged"] = coins if charged_here else 0
        payload["balance"] = wallet.get_balance(user_id)
    return jsonify(payload)


@app.post("/api/assignment/projects/<project_id>/research")
def api_assignment_project_research(project_id: str):
    """Build Research Plan from Requirement JSON and parsed documents only."""
    trace(
        "api.research.received",
        **project_service.store.lookup_diagnostics(project_id),
    )
    payload = request.get_json(silent=True) or {}
    parsed_documents = payload.get("parsed_documents")
    if parsed_documents is not None and not isinstance(parsed_documents, list):
        return jsonify({"error": "parsed_documents must be an array"}), 400
    try:
        plan = project_service.run_research(project_id, parsed_documents=parsed_documents)
    except KeyError as exc:
        return _assignment_not_found("research", project_id, exc)
    except ValueError as exc:
        trace("api.research.failed", project_id=project_id, error=str(exc))
        return jsonify({"error": str(exc)}), 502
    except Exception as exc:  # noqa: BLE001
        trace(
            "api.research.error",
            project_id=project_id,
            error=str(exc),
            error_type=type(exc).__name__,
        )
        return jsonify({"error": "Research planning failed. Please try again."}), 502
    trace(
        "api.research.completed",
        project_id=project_id,
        plan_id=plan.id,
        engine=plan.engine_version,
    )
    return jsonify({"research_plan": plan.to_dict()})


@app.get("/api/assignment/projects/<project_id>/research-plan")
def api_assignment_project_research_plan(project_id: str):
    try:
        plan = project_service.get_research_plan(project_id)
    except KeyError:
        return jsonify({"error": "Research plan not found"}), 404
    return jsonify(plan.to_dict())


@app.patch("/api/assignment/projects/<project_id>/research-plan")
def api_assignment_project_research_plan_update(project_id: str):
    payload = request.get_json(silent=True) or {}
    if not isinstance(payload, dict):
        return jsonify({"error": "Invalid payload"}), 400
    try:
        plan = project_service.update_research_plan(project_id, payload)
    except KeyError:
        return jsonify({"error": "Research plan not found"}), 404
    return jsonify(plan.to_dict())


@app.post("/api/assignment/projects/<project_id>/blueprint")
def api_assignment_project_blueprint(project_id: str):
    """Build writing Blueprint from Requirement JSON + Research Plan only."""
    payload = request.get_json(silent=True) or {}
    research_plan = payload.get("research_plan")
    if research_plan is not None and not isinstance(research_plan, dict):
        return jsonify({"error": "research_plan must be an object"}), 400
    trace(
        "api.blueprint.received",
        has_client_research_plan=isinstance(research_plan, dict),
        **project_service.store.lookup_diagnostics(project_id),
    )
    try:
        blueprint = project_service.run_blueprint(project_id, research_plan=research_plan)
    except KeyError as exc:
        return _assignment_not_found("blueprint", project_id, exc)
    except ValueError as exc:
        trace("api.blueprint.failed", project_id=project_id, error=str(exc))
        return jsonify({"error": str(exc)}), 502
    except Exception as exc:  # noqa: BLE001
        trace(
            "api.blueprint.error",
            project_id=project_id,
            error=str(exc),
            error_type=type(exc).__name__,
        )
        return jsonify({"error": "Blueprint generation failed. Please try again."}), 502
    trace(
        "api.blueprint.completed",
        project_id=project_id,
        blueprint_id=blueprint.id,
        engine=blueprint.engine_version,
    )
    return jsonify({"blueprint": blueprint.to_dict()})


@app.get("/api/assignment/projects/<project_id>/blueprint")
def api_assignment_project_blueprint_get(project_id: str):
    try:
        blueprint = project_service.get_blueprint(project_id)
    except KeyError:
        return jsonify({"error": "Blueprint not found"}), 404
    return jsonify(blueprint.to_dict())


@app.patch("/api/assignment/projects/<project_id>/blueprint")
def api_assignment_project_blueprint_update(project_id: str):
    payload = request.get_json(silent=True) or {}
    if not isinstance(payload, dict):
        return jsonify({"error": "Invalid payload"}), 400
    try:
        blueprint = project_service.update_blueprint(project_id, payload)
    except KeyError:
        return jsonify({"error": "Blueprint not found"}), 404
    return jsonify(blueprint.to_dict())


@app.get("/api/debug/blueprint/<project_id>")
def api_debug_blueprint_get(project_id: str):
    """Debug endpoint: return stored Blueprint JSON without transformations."""
    try:
        blueprint = project_service.get_blueprint(project_id)
    except KeyError:
        return jsonify({"error": "Blueprint not found"}), 404
    return jsonify(blueprint.to_dict())


@app.get("/api/browser/runtime-test")
def api_browser_runtime_test():
    """Temporary diagnostics: verify BrowserService can serve a page via CDP."""
    from services.browser.browser_service import BrowserService

    def _work():
        service = BrowserService.instance()
        service.ensure_running()
        page = service.new_page()
        started = time.monotonic()
        page.goto("https://example.com", wait_until="networkidle")
        load_time_ms = int((time.monotonic() - started) * 1000)
        browser = page.context.browser
        version = browser.version if browser is not None else "unknown"
        title = page.title()
        url = page.url
        health = service.health()
        try:
            page.close()
        except Exception:  # noqa: BLE001
            pass
        return {
            "success": True,
            "browser": {"engine": "chromium", "version": version, "mode": "cdp", "headless": False},
            "page": {"title": title, "url": url, "load_time_ms": load_time_ms},
            "runtime": {
                "pages_open": health.get("pages", 0),
                "cdp_url": service.cdp_url,
                "connected": bool(health.get("connected")),
            },
        }

    try:
        return jsonify(_browser_submit(_work, timeout=90))
    except Exception as exc:  # noqa: BLE001
        return (
            jsonify(
                {
                    "success": False,
                    "error": str(exc),
                    "exception_type": type(exc).__name__,
                    "traceback": traceback.format_exc(),
                }
            ),
            500,
        )


@app.get("/api/browser/health")
def api_browser_health():
    """Aggregate Browser Service + provider diagnostics (Task 10 health)."""
    from services.browser.browser_service import BrowserService

    ensure_engine_started()
    try:
        data = _browser_submit(lambda: BrowserService.instance().health(), timeout=15)
    except Exception:  # noqa: BLE001  (worker busy or error → degraded, no Playwright)
        from services.browser.chrome_launcher import ChromeLauncher

        launcher = ChromeLauncher()
        data = {
            "browser_running": launcher.is_cdp_available(),
            "connected": None,
            "chrome_pid": launcher.pid,
            "memory_usage": launcher.memory_usage(),
            "cdp_url": launcher.cdp_url,
            "note": "degraded snapshot (worker busy)",
        }
    data["active_jobs"] = job_manager.active_count() if job_manager is not None else 0
    if health_monitor is not None:
        data["last_monitor_check"] = health_monitor.last
    return jsonify(data)


@app.get("/api/browser/cdp-diagnostics")
def api_browser_cdp_diagnostics():
    """Diagnose the CDP connection layer (Chrome/Playwright versions, download-behavior support).

    Read-only: uses HTTP + a throwaway Playwright instance, never touches the
    long-lived BrowserService connection.
    """
    from services.browser.cdp_compat import cdp_diagnostics
    from services.browser.chrome_launcher import ChromeLauncher

    try:
        cdp_url = ChromeLauncher().cdp_url
        report = cdp_diagnostics(cdp_url)
        return jsonify({"success": bool(report.get("connect_ok")), **report})
    except Exception as exc:  # noqa: BLE001
        return (
            jsonify(
                {
                    "success": False,
                    "error": str(exc),
                    "exception_type": type(exc).__name__,
                    "traceback": traceback.format_exc(),
                }
            ),
            500,
        )


@app.get("/api/browser/connect")
def api_browser_connect():
    """Start/reuse the long-lived Browser Service (auto-launches Chrome + CDP)."""
    from services.browser.browser_service import BrowserService

    try:
        ensure_engine_started()
        _browser_submit(lambda: BrowserService.instance().start(), timeout=90)
        health = _browser_submit(lambda: BrowserService.instance().health(), timeout=15)
        return jsonify(
            {
                "success": True,
                "connected": bool(health.get("connected")),
                "contexts": int(health.get("contexts") or 0),
                "pages": int(health.get("pages") or 0),
                "browser": "Google Chrome",
            }
        )
    except Exception as exc:  # noqa: BLE001
        return (
            jsonify(
                {
                    "success": False,
                    "connected": False,
                    "error": str(exc),
                    "exception_type": type(exc).__name__,
                    "traceback": traceback.format_exc(),
                }
            ),
            500,
        )


@app.get("/api/browser/providers/stealthwriter/health")
def api_browser_stealthwriter_health():
    """Temporary diagnostics: load StealthWriter public homepage (no login)."""
    from services.browser.providers.stealthwriter import StealthWriterProvider

    def _work():
        provider = StealthWriterProvider()
        provider.initialize()
        return provider.health_check()

    try:
        health = _browser_submit(_work, timeout=90)
        details = health.details or {}
        return jsonify(
            {
                "success": bool(health.healthy),
                "provider": "stealthwriter",
                "title": details.get("title"),
                "url": details.get("url"),
                "logged_in": bool(details.get("logged_in")),
                "turnstile": bool(details.get("turnstile")),
                "cookies": int(details.get("cookies") or 0),
                "page_load_ms": details.get("page_load_ms"),
                "screenshot": details.get("screenshot"),
                "http_status": details.get("http_status"),
                "login_button": details.get("login_button"),
                "dashboard": details.get("dashboard"),
                "humanizer": details.get("humanizer"),
                "localStorage_keys": details.get("localStorage_keys"),
                "sessionStorage_keys": details.get("sessionStorage_keys"),
                "message": health.message,
            }
        )
    except Exception as exc:  # noqa: BLE001
        return (
            jsonify(
                {
                    "success": False,
                    "provider": "stealthwriter",
                    "error": str(exc),
                    "exception_type": type(exc).__name__,
                    "traceback": traceback.format_exc(),
                }
            ),
            500,
        )


@app.get("/api/browser/providers/stealthwriter/session")
def api_browser_stealthwriter_session():
    """Temporary diagnostics: check CDP StealthWriter session (verbose failure reporting)."""
    from services.browser.browser_service import BrowserService

    diag: dict[str, Any] = {
        "cdp_url": None,
        "runtime_initialized": False,
        "cdp_connected": False,
        "context_created": False,
        "page_created": False,
        "current_url": None,
    }

    def _work():
        service = BrowserService.instance()
        diag["cdp_url"] = service.cdp_url
        service.start()
        health = service.health()
        diag["runtime_initialized"] = True
        diag["cdp_connected"] = bool(health.get("connected"))
        diag["context_created"] = int(health.get("contexts") or 0) > 0

        page = service.get_or_create_page("stealthwriter")
        diag["page_created"] = True
        page.goto("https://stealthwriter.ai/dashboard", wait_until="networkidle")
        diag["current_url"] = page.url

        current_url = page.url
        title = page.title()
        redirected = "/sign-in" in current_url.lower() or "/signin" in current_url.lower()
        logged_in = (not redirected) and ("/dashboard" in current_url.lower())
        return {
            "success": True,
            "logged_in": logged_in,
            "cdp_url": diag["cdp_url"],
            "dashboard_loaded": logged_in,
            "current_url": current_url,
            "title": title,
            "redirected": redirected,
            "diagnostics": diag,
        }

    try:
        return jsonify(_browser_submit(_work, timeout=90))
    except Exception as exc:  # noqa: BLE001
        tb = traceback.format_exc()
        return (
            jsonify(
                {
                    "success": False,
                    "exception": str(exc),
                    "exception_type": type(exc).__name__,
                    "traceback": tb,
                    "diagnostics": diag,
                }
            ),
            500,
        )


@app.get("/api/browser/providers/stealthwriter/login")
def api_browser_stealthwriter_login():
    """Navigate attached Chrome to sign-in for interactive manual login."""
    from services.browser.providers.stealthwriter import start_interactive_login

    try:
        result = _browser_submit(start_interactive_login, timeout=60)
        return jsonify(
            {
                "success": True,
                "message": result.get("message") or "Attached to Chrome. Please login manually.",
                "profile": result.get("profile"),
                "cdp_url": result.get("cdp_url"),
                "already_open": result.get("already_open", False),
            }
        )
    except Exception as exc:  # noqa: BLE001
        return (
            jsonify(
                {
                    "success": False,
                    "error": str(exc),
                    "exception_type": type(exc).__name__,
                    "traceback": traceback.format_exc(),
                }
            ),
            500,
        )


@app.get("/api/browser/providers/stealthwriter/check-login")
def api_browser_stealthwriter_check_login():
    """Verify interactive login; shut down runtime only after success."""
    from services.browser.providers.stealthwriter import check_interactive_login

    try:
        result = _browser_submit(check_interactive_login, timeout=60)
        return jsonify(
            {
                "success": True,
                "logged_in": bool(result.get("logged_in")),
                "dashboard_loaded": bool(result.get("dashboard_loaded")),
                "current_url": result.get("current_url"),
                "title": result.get("title"),
                "redirected": result.get("redirected"),
                "profile": result.get("profile"),
                "runtime_shutdown": result.get("runtime_shutdown"),
                "message": result.get("message"),
            }
        )
    except Exception as exc:  # noqa: BLE001
        return (
            jsonify(
                {
                    "success": False,
                    "logged_in": False,
                    "exception": str(exc),
                    "exception_type": type(exc).__name__,
                    "traceback": traceback.format_exc(),
                }
            ),
            500,
        )


@app.post("/api/browser/providers/stealthwriter/humanize")
@economy_auth.email_verified_required
def api_browser_stealthwriter_humanize():
    """Humanize via the production job engine (retries, timeout, recovery, logs).

    Backward compatible: still returns humanized_text synchronously, plus job_id.
    """
    from services.browser.jobs.retry import MAX_RETRIES

    payload = request.get_json(silent=True) or {}
    text = payload.get("text")
    if not isinstance(text, str) or not text.strip():
        return jsonify({"success": False, "error": "text is required"}), 400

    raw_source = str(payload.get("source") or "standalone").strip().lower()
    dataset_source = (
        raw_source
        if raw_source in ("standalone", "workspace_partial", "assignment")
        else "standalone"
    )

    from services.economy.pricing import FEATURE_COSTS
    from services.economy.site_settings import humanize_credit_cost

    pricing = humanize_credit_cost(FEATURE_COSTS["humanize"])
    cost = int(pricing["charged"])
    charged = _charge_current_user(
        "humanize",
        cost,
        meta={
            "original_price": pricing["original_price"],
            "discount_active": pricing["discount_active"],
            "discount_percent": pricing["discount_percent"],
            "discount_source": pricing["discount_source"],
        },
    )
    if not isinstance(charged, tuple):
        return charged
    user_id, _tx = charged

    try:
        ensure_engine_started()
        job = job_manager.create(
            "stealthwriter", "humanize", {"text": text}, max_retries=MAX_RETRIES
        )
        max_wait = (browser_worker.job_timeout + 30) * (job.max_retries + 1) + 30
        job_manager.wait(job.id, timeout=max_wait)
        job = job_manager.get(job.id)
        if job is None:
            _refund_safe(user_id, cost, "humanize")
            return jsonify({"success": False, "error": "job not found"}), 500

        status = job.status.value
        if status == "COMPLETED":
            res = job.result or {}
            elapsed = res.get("elapsed_seconds")
            latency_ms = int(float(elapsed) * 1000) if elapsed is not None else None
            humanized_out = res.get("humanized_text")
            _record_usage_safe(
                user_id,
                feature=FEATURE_HUMANIZER,
                credits_used=cost,
                provider="StealthWriter",
                latency=latency_ms,
                request_id=job.id,
            )
            try:
                from services.economy.site_settings import record_humanizer_success

                record_humanizer_success()
            except Exception:
                app.logger.exception("daily_stats humanizer increment failed")
            try:
                from services.dataset_logger import log_humanization_event

                log_humanization_event(
                    user_id,
                    dataset_source,
                    text,
                    humanized_out if isinstance(humanized_out, str) else "",
                )
            except Exception:
                app.logger.exception("dataset_logger stealthwriter hook failed")
            return jsonify(
                {
                    "success": True,
                    "humanized_text": humanized_out,
                    "elapsed_seconds": res.get("elapsed_seconds"),
                    "job_id": job.id,
                    "coins_charged": cost,
                    "original_price": pricing["original_price"],
                    "discount_active": pricing["discount_active"],
                    "discount_percent": pricing["discount_percent"],
                    "balance": wallet.get_balance(user_id),
                }
            )

        # Any non-completed terminal/failure state refunds the charge.
        _refund_safe(user_id, cost, "humanize", ref_id=job.id)

        if status == "CANCELLED":
            return jsonify({"success": False, "error": "cancelled", "job_id": job.id}), 409

        code = job.error_code
        if code == "LOGIN_REQUIRED":
            return jsonify(
                {
                    "success": False,
                    "error": "LOGIN_REQUIRED",
                    "message": job.error
                    or "StealthWriter session is not logged in. Upload browser_profiles/sessions/stealthwriter.json.",
                    "job_id": job.id,
                }
            ), 401
        if code == "NO_CHANGE":
            return (
                jsonify(
                    {
                        "success": False,
                        "error": "NO_CHANGE",
                        "message": job.error,
                        "job_id": job.id,
                    }
                ),
                409,
            )
        if code == "TIMEOUT":
            return jsonify({"success": False, "error": "timeout", "job_id": job.id}), 504
        if code == "AUTOMATION_ERROR" and job.error_details:
            body = {
                "success": False,
                "error": job.error,
                "job_id": job.id,
                "diagnostics": job.error_details,
            }
            for key in (
                "current_url",
                "page_title",
                "visible_buttons",
                "textarea_count",
                "dom_snippet",
                "step",
            ):
                if key in job.error_details:
                    body[key] = job.error_details[key]
            return jsonify(body), 422
        if not job.is_terminal:
            return (
                jsonify(
                    {"success": False, "error": "still running", "status": status, "job_id": job.id}
                ),
                202,
            )
        return jsonify({"success": False, "error": job.error or "failed", "job_id": job.id}), 500
    except Exception as exc:  # noqa: BLE001
        _refund_safe(user_id, cost, "humanize")
        return (
            jsonify(
                {
                    "success": False,
                    "error": str(exc),
                    "exception_type": type(exc).__name__,
                    "traceback": traceback.format_exc(),
                }
            ),
            500,
        )


@app.get("/api/browser/providers/stealthwriter/status")
def api_browser_stealthwriter_status():
    """Return StealthWriter session status from the dashboard sidebar."""
    from services.browser.providers.stealthwriter import get_session_status

    try:
        status = _browser_submit(get_session_status, timeout=120)
        return jsonify(
            {
                "logged_in": bool(status.get("logged_in")),
                "current_url": status.get("current_url"),
                "plan": status.get("plan"),
                "username": status.get("username"),
            }
        )
    except Exception as exc:  # noqa: BLE001
        return (
            jsonify(
                {
                    "logged_in": False,
                    "current_url": None,
                    "plan": None,
                    "username": None,
                    "error": str(exc),
                    "exception_type": type(exc).__name__,
                    "traceback": traceback.format_exc(),
                }
            ),
            500,
        )


@app.get("/api/browser/providers/stealthwriter/verify-production")
def api_browser_stealthwriter_verify_production():
    """End-to-end check: session restored + real Humanize request completes."""
    from services.browser.browser_service import BrowserService
    from services.browser.providers import stealthwriter as sw

    sample_text = "This text was generated by artificial intelligence and should be rewritten."

    def _work() -> dict[str, Any]:
        sessions = BrowserService.instance().sessions
        status = sw.get_session_status()
        result: dict[str, Any] = {
            "session_file_present": sessions.has(sw.PROVIDER_NAME),
            "current_url": status.get("current_url"),
            "logged_in": bool(status.get("logged_in")),
            "username": status.get("username"),
            "plan": status.get("plan"),
            "input_length": len(sample_text),
            "output_length": 0,
            "processing_time_ms": 0,
            "success": False,
        }

        if not status.get("logged_in") or sw._is_sign_in_url(status.get("current_url") or ""):
            result["error"] = "LOGIN_REQUIRED"
            return result

        started = time.monotonic()
        try:
            humanize_result = sw.humanize_text(sample_text)
        except sw.StealthWriterAutomationError as exc:
            result["processing_time_ms"] = int((time.monotonic() - started) * 1000)
            result["error"] = "HUMANIZE_FAILED"
            result["message"] = str(exc)
            return result

        result["processing_time_ms"] = int((time.monotonic() - started) * 1000)
        result["current_url"] = humanize_result.get("current_url") or result.get("current_url")
        output_text = (humanize_result.get("humanized_text") or "").strip()
        result["output_length"] = len(output_text)

        if sw._is_sign_in_url(result.get("current_url") or ""):
            result["logged_in"] = False
            result["error"] = "LOGIN_REQUIRED"
            return result

        if humanize_result.get("error") == "LOGIN_REQUIRED":
            result["logged_in"] = False
            result["error"] = "LOGIN_REQUIRED"
            return result

        if not humanize_result.get("success"):
            result["error"] = humanize_result.get("error") or "HUMANIZE_FAILED"
            result["message"] = humanize_result.get("message")
            return result

        if not output_text:
            result["error"] = "OUTPUT_MISSING"
            result["message"] = "Humanize completed but returned empty output."
            return result

        if output_text == sample_text:
            result["error"] = "OUTPUT_UNCHANGED"
            result["message"] = "Humanized output is identical to the input."
            return result

        result["success"] = True
        return result

    try:
        return jsonify(_browser_submit(_work, timeout=180))
    except Exception as exc:  # noqa: BLE001
        return (
            jsonify(
                {
                    "success": False,
                    "logged_in": False,
                    "error": str(exc),
                    "exception_type": type(exc).__name__,
                    "traceback": traceback.format_exc(),
                }
            ),
            500,
        )


@app.get("/api/browser/providers/stealthwriter/open")
def api_browser_stealthwriter_open():
    """Navigate attached Chrome to sign-in for one-time manual login (no automation)."""
    from services.browser.providers.stealthwriter import open_manual_login_browser

    try:
        result = _browser_submit(open_manual_login_browser, timeout=60)
        return jsonify(
            {
                "success": True,
                "message": result["message"],
                "profile": result.get("profile"),
                "already_open": result.get("already_open", False),
            }
        )
    except Exception as exc:  # noqa: BLE001
        return (
            jsonify(
                {
                    "success": False,
                    "error": str(exc),
                    "exception_type": type(exc).__name__,
                    "traceback": traceback.format_exc(),
                }
            ),
            500,
        )


@app.get("/api/jobs")
def api_jobs_list():
    """List all jobs (most recent last)."""
    ensure_engine_started()
    return jsonify({"jobs": [job.to_dict() for job in job_manager.list()]})


@app.get("/api/jobs/<job_id>")
def api_jobs_get(job_id: str):
    """Return a single job by id."""
    ensure_engine_started()
    job = job_manager.get(job_id)
    if job is None:
        return jsonify({"error": "job not found"}), 404
    return jsonify(job.to_dict())


@app.delete("/api/jobs/<job_id>")
def api_jobs_cancel(job_id: str):
    """Cancel a job (immediate if still queued; cooperative otherwise)."""
    ensure_engine_started()
    if job_manager.get(job_id) is None:
        return jsonify({"error": "job not found"}), 404
    cancelled = job_manager.cancel(job_id)
    return jsonify({"success": bool(cancelled), "job_id": job_id})


@app.post("/api/browser/restart")
def api_browser_restart():
    """Force a full Browser Service restart (Chrome relaunch + reconnect + restore)."""
    from services.browser.browser_service import BrowserService

    try:
        ensure_engine_started()
        result = _browser_submit(lambda: BrowserService.instance().restart(), timeout=120)
        if browser_metrics is not None:
            browser_metrics.record_browser_restart()
        return jsonify({"success": True, "restart": result})
    except Exception as exc:  # noqa: BLE001
        return (
            jsonify(
                {
                    "success": False,
                    "error": str(exc),
                    "exception_type": type(exc).__name__,
                    "traceback": traceback.format_exc(),
                }
            ),
            500,
        )


@app.get("/api/browser/metrics")
def api_browser_metrics():
    """Engine metrics: uptime, job counts, avg exec time, restarts, retries."""
    ensure_engine_started()
    data = browser_metrics.snapshot() if browser_metrics is not None else {}
    data["active_jobs"] = job_manager.active_count() if job_manager is not None else 0
    data["total_jobs"] = len(job_manager.list()) if job_manager is not None else 0
    data["worker_ready"] = bool(browser_worker.is_ready()) if browser_worker is not None else False
    return jsonify(data)


@app.post("/api/blueprint/build")
def api_blueprint_build():
    """Standalone Blueprint Engine entrypoint."""
    payload = request.get_json(silent=True) or {}
    requirement_json = payload.get("requirement_json")
    research_plan = payload.get("research_plan")
    if not isinstance(requirement_json, dict):
        return jsonify({"error": "requirement_json object is required"}), 400
    if not isinstance(research_plan, dict):
        return jsonify({"error": "research_plan object is required"}), 400
    blueprint = blueprint_engine.build_blueprint(
        requirement_json=requirement_json,
        research_plan=research_plan,
        project_id=payload.get("project_id"),
    )
    return jsonify(blueprint.to_dict()), 201


@app.post("/api/writer/session")
def api_writer_session_create():
    payload = request.get_json(silent=True) or {}
    for key in ("requirement_json", "research_plan", "blueprint"):
        if not isinstance(payload.get(key), dict):
            return jsonify({"error": f"{key} object is required"}), 400
    session = writer_engine.create_session(
        requirement_json=payload["requirement_json"],
        research_plan=payload["research_plan"],
        blueprint=payload["blueprint"],
        project_id=payload.get("project_id"),
    )
    return jsonify(session.to_dict()), 201


@app.get("/api/writer/session/<session_id>")
def api_writer_session_get(session_id: str):
    try:
        session = writer_engine.get_session(session_id)
    except KeyError:
        return jsonify({"error": "Writer session not found"}), 404
    return jsonify(session.to_dict())


@app.post("/api/writer/session/<session_id>/advance")
def api_writer_session_advance(session_id: str):
    try:
        session = writer_engine.advance_section(session_id)
    except KeyError:
        return jsonify({"error": "Writer session not found"}), 404
    return jsonify(session.to_dict())


@app.post("/api/writer/session/<session_id>/revise")
def api_writer_session_revise(session_id: str):
    payload = request.get_json(silent=True) or {}
    try:
        session = writer_engine.revise_section(session_id, payload.get("section_id"))
    except KeyError:
        return jsonify({"error": "Writer session not found"}), 404
    except ValueError as exc:
        return jsonify({"error": str(exc)}), 400
    return jsonify(session.to_dict())


@app.post("/api/writer/session/<session_id>/merge")
def api_writer_session_merge(session_id: str):
    payload = request.get_json(silent=True) or {}
    try:
        draft = writer_engine.merge_draft(session_id, title=payload.get("title"))
    except KeyError:
        return jsonify({"error": "Writer session not found"}), 404
    except ValueError as exc:
        return jsonify({"error": str(exc)}), 400
    return jsonify(draft.to_dict()), 201


@app.post("/api/assignment/projects/<project_id>/writer/start")
def api_assignment_writer_start(project_id: str):
    trace(
        "api.writer.start.received",
        **project_service.store.lookup_diagnostics(project_id),
    )
    try:
        session = project_service.start_writer(project_id)
    except KeyError:
        return jsonify({"error": "Project not found"}), 404
    except ValueError as exc:
        trace("api.writer.start.failed", project_id=project_id, error=str(exc))
        return jsonify({"error": str(exc)}), 400
    except Exception as exc:  # noqa: BLE001
        trace(
            "api.writer.start.error",
            project_id=project_id,
            error=str(exc),
            error_type=type(exc).__name__,
        )
        return jsonify({"error": "Writer setup failed. Please try again."}), 502
    trace(
        "api.writer.start.completed",
        project_id=project_id,
        session_id=session.id,
        sections=len(session.sections),
        engine=session.engine_version,
    )
    return jsonify(session.to_dict()), 201


@app.post("/api/assignment/projects/<project_id>/writer/advance")
def api_assignment_writer_advance(project_id: str):
    payload = request.get_json(silent=True) or {}
    writer_session = payload.get("writer_session")
    if writer_session is not None and not isinstance(writer_session, dict):
        return jsonify({"error": "writer_session must be an object"}), 400
    trace(
        "api.writer.advance.received",
        has_client_writer_session=isinstance(writer_session, dict),
        **project_service.store.lookup_diagnostics(project_id),
    )
    try:
        session = project_service.advance_writer(project_id, writer_session=writer_session)
    except KeyError:
        return jsonify({"error": "Writer session not found"}), 404
    except ValueError as exc:
        message = str(exc)
        trace("api.writer.advance.failed", project_id=project_id, error=message)
        return jsonify({"error": user_friendly_llm_error(message)}), llm_error_http_status(message)
    except Exception as exc:  # noqa: BLE001
        trace(
            "api.writer.advance.error",
            project_id=project_id,
            error=str(exc),
            error_type=type(exc).__name__,
        )
        return jsonify({"error": "Writing step failed. Please try again."}), 502
    current = session.current_section()
    trace(
        "api.writer.advance.completed",
        project_id=project_id,
        session_id=session.id,
        progress=session.progress,
        status=session.status.value,
        current_section=current.title if current else None,
        completed_sections=len(session.completed_section_ids),
    )
    return jsonify(session.to_dict())


@app.post("/api/assignment/projects/<project_id>/writer/revise")
def api_assignment_writer_revise(project_id: str):
    payload = request.get_json(silent=True) or {}
    writer_session = payload.get("writer_session")
    if writer_session is not None and not isinstance(writer_session, dict):
        return jsonify({"error": "writer_session must be an object"}), 400
    try:
        session = project_service.revise_writer_section(
            project_id,
            payload.get("section_id"),
            writer_session=writer_session,
        )
    except KeyError:
        return jsonify({"error": "Writer session not found"}), 404
    except ValueError as exc:
        return jsonify({"error": str(exc)}), 400
    except Exception as exc:  # noqa: BLE001
        trace(
            "api.writer.revise.error",
            project_id=project_id,
            error=str(exc),
            error_type=type(exc).__name__,
        )
        return jsonify({"error": "Writing revision failed. Please try again."}), 502
    return jsonify(session.to_dict())


@app.post("/api/assignment/projects/<project_id>/writer/merge")
def api_assignment_writer_merge(project_id: str):
    payload = request.get_json(silent=True) or {}
    writer_session = payload.get("writer_session")
    if writer_session is not None and not isinstance(writer_session, dict):
        return jsonify({"error": "writer_session must be an object"}), 400
    trace(
        "api.writer.merge.received",
        has_client_writer_session=isinstance(writer_session, dict),
        **project_service.store.lookup_diagnostics(project_id),
    )
    try:
        draft = project_service.merge_writer_draft(project_id, writer_session=writer_session)
    except KeyError:
        return jsonify({"error": "Writer session not found"}), 404
    except ValueError as exc:
        trace("api.writer.merge.failed", project_id=project_id, error=str(exc))
        return jsonify({"error": str(exc)}), 400
    except Exception as exc:  # noqa: BLE001
        trace(
            "api.writer.merge.error",
            project_id=project_id,
            error=str(exc),
            error_type=type(exc).__name__,
        )
        return jsonify({"error": "Draft merge failed. Please try again."}), 502
    trace(
        "api.writer.merge.completed",
        project_id=project_id,
        draft_id=draft.id,
        total_words=draft.total_words,
        model=draft.model,
    )
    return jsonify(draft.to_dict()), 201


@app.get("/api/assignment/projects/<project_id>/writer")
def api_assignment_writer_get(project_id: str):
    try:
        session = project_service.get_writer_session(project_id)
    except KeyError:
        return jsonify({"error": "Writer session not found"}), 404
    except Exception as exc:  # noqa: BLE001
        trace(
            "api.writer.get.error",
            project_id=project_id,
            error=str(exc),
            error_type=type(exc).__name__,
        )
        return jsonify({"error": "Failed to load writer session. Please retry."}), 500
    return jsonify(session.to_dict())


@app.get("/api/assignment/projects/<project_id>/draft")
def api_assignment_draft_get(project_id: str):
    try:
        draft = project_service.get_draft(project_id)
    except KeyError:
        return jsonify({"error": "Draft not found"}), 404
    return jsonify(draft.to_dict())


@app.get("/api/debug/draft/<project_id>")
def api_debug_draft_get(project_id: str):
    """Debug endpoint: return stored Draft JSON without transformations."""
    try:
        draft = project_service.get_draft(project_id)
    except KeyError:
        return jsonify({"error": "Draft not found"}), 404
    return jsonify(draft.to_dict())


@app.get("/api/debug/section-review/<project_id>")
def api_debug_section_review_get(project_id: str):
    """Return all section review results for project draft without transformations."""
    try:
        draft = project_service.get_draft(project_id).to_dict()
    except KeyError:
        return jsonify({"error": "Draft not found"}), 404
    sections = list(draft.get("sections") or [])
    review_results = [
        {
            "title": section.get("title"),
            "review_result": section.get("review_result"),
        }
        for section in sections
    ]
    return jsonify(
        {
            "project_id": project_id,
            "sections_review": review_results,
        }
    )


def _debug_now_iso() -> str:
    return utc_now().isoformat()


def _extract_text_from_pptx_bytes(raw: bytes) -> tuple[str, int]:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(raw))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if isinstance(text, str) and text.strip():
                parts.append(text.strip())
    return ("\n\n".join(parts).strip(), len(prs.slides))


def _extract_uploaded_text(raw: bytes, filename: str, mimetype: str | None = None) -> tuple[str, int | None]:
    ext = upload_extension(filename, mimetype).lower()
    if ext in {".docx", ".pdf"}:
        text = extract_text_from_document_bytes(raw, filename, mimetype)
        pages: int | None = None
        if ext == ".pdf":
            try:
                from pypdf import PdfReader

                pages = len(PdfReader(io.BytesIO(raw)).pages)
            except Exception:  # noqa: BLE001
                pages = None
        return text, pages
    if ext == ".txt":
        return raw.decode("utf-8", errors="replace").strip(), None
    if ext == ".pptx":
        text, slides = _extract_text_from_pptx_bytes(raw)
        return text, slides
    raise ValueError("Unsupported file type. Allowed: PDF, DOCX, TXT, PPTX")


def _collect_debug_input_payload() -> tuple[dict[str, Any], list[dict[str, Any]], list[dict[str, Any]]]:
    if request.mimetype == "application/json" or request.is_json:
        payload = request.get_json(silent=True) or {}
        if not isinstance(payload, dict):
            raise ValueError("Invalid JSON payload")
        return payload, [], []

    payload: dict[str, Any] = {
        "assignment_brief": request.form.get("assignment_brief", ""),
        "rubric": request.form.get("rubric", ""),
        "lecture_notes": request.form.get("lecture_notes", ""),
        "uploaded_files": [],
    }
    parsed_files: list[dict[str, Any]] = []
    file_errors: list[dict[str, Any]] = []

    field_to_type = {
        "assignment_brief": "assignment_brief",
        "rubric": "rubric",
        "lecture_notes": "lecture_slides",
        "professor_notes": "professor_notes",
        "reading_materials": "reading_material",
        "additional_files": "additional_file",
    }
    aggregate_fields = {"lecture_notes", "reading_materials"}

    aggregated_by_field: dict[str, list[str]] = {field: [] for field in aggregate_fields}
    aggregated_meta: dict[str, list[dict[str, Any]]] = {field: [] for field in aggregate_fields}

    for field, file_type in field_to_type.items():
        files = request.files.getlist(field)
        for upload in files:
            if not upload or not upload.filename:
                continue
            filename = upload.filename
            try:
                raw = upload.read()
                if not raw:
                    raise ValueError("Uploaded file is empty")
                text, pages = _extract_uploaded_text(raw, filename, upload.mimetype)
                characters = len(text)
                words = len([w for w in text.split() if w.strip()])
                preview = text[:500]
                file_info = {
                    "filename": filename,
                    "file_type": file_type,
                    "pages": pages,
                    "characters": characters,
                    "words": words,
                    "extracted_text_preview": preview,
                }
                parsed_files.append(file_info)
                if field in aggregate_fields:
                    aggregated_by_field[field].append(text)
                    aggregated_meta[field].append({"filename": filename, "text": text, "file_type": file_type})
                elif field in {"assignment_brief", "rubric"}:
                    existing = str(payload.get(field) or "").strip()
                    payload[field] = f"{existing}\n\n{text}".strip() if existing else text
                else:
                    payload["uploaded_files"].append(
                        {
                            "filename": filename,
                            "file_type": file_type,
                            "text": text,
                        }
                    )
            except Exception as exc:  # noqa: BLE001
                file_errors.append(
                    {
                        "filename": filename,
                        "file_type": file_type,
                        "message": str(exc),
                        "traceback": traceback.format_exc(),
                    }
                )

    if aggregated_by_field["lecture_notes"]:
        payload["lecture_notes"] = "\n\n".join(aggregated_by_field["lecture_notes"]).strip()
    for entry in aggregated_meta["reading_materials"]:
        payload["uploaded_files"].append(entry)

    return payload, parsed_files, file_errors


def _write_debug_text_file(project_id: str, filename: str, content: str) -> str:
    base = project_files_dir(project_id)
    base.mkdir(parents=True, exist_ok=True)
    path = base / filename
    path.write_text(content or "", encoding="utf-8")
    return str(path)


def _write_project_binary_file(project_id: str, filename: str, raw: bytes) -> str:
    base = project_files_dir(project_id)
    base.mkdir(parents=True, exist_ok=True)
    safe_name = Path(filename).name
    path = base / f"{uuid.uuid4()}_{safe_name}"
    path.write_bytes(raw)
    return str(path)


def _attach_multipart_uploads(project_id: str) -> list[dict[str, Any]]:
    """Persist uploaded brief/rubric/materials for a project.

    Chat UX may send undifferentiated files via the ``files`` field; the first
    document becomes ``assignment_brief`` and the rest ``additional_file``.
    Legacy field names remain supported.
    """
    field_to_type = {
        "assignment_brief": "assignment_brief",
        "rubric": "rubric",
        "lecture_notes": "lecture_slides",
        "professor_notes": "professor_notes",
        "reading_materials": "reading_material",
        "additional_files": "additional_file",
    }
    saved: list[dict[str, Any]] = []
    for field, file_type in field_to_type.items():
        for upload in request.files.getlist(field):
            if not upload or not upload.filename:
                continue
            raw = upload.read()
            if not raw:
                raise ValueError(f"{upload.filename} is empty")
            storage_path = _write_project_binary_file(project_id, upload.filename, raw)
            record = project_service.add_file(
                project_id,
                file_type=file_type,
                original_filename=upload.filename,
                storage_path=storage_path,
                parsed=False,
            )
            saved.append(record.to_dict())

    # Chat composer: single ``files`` pool (Gemini classifies content later).
    generic = [u for u in request.files.getlist("files") if u and u.filename]
    for idx, upload in enumerate(generic):
        raw = upload.read()
        if not raw:
            raise ValueError(f"{upload.filename} is empty")
        file_type = "assignment_brief" if idx == 0 and not saved else "additional_file"
        # If a legacy brief already exists, treat every generic file as material.
        if any(r.get("file_type") == "assignment_brief" for r in saved):
            file_type = "additional_file"
        elif idx == 0:
            file_type = "assignment_brief"
        storage_path = _write_project_binary_file(project_id, upload.filename, raw)
        record = project_service.add_file(
            project_id,
            file_type=file_type,
            original_filename=upload.filename,
            storage_path=storage_path,
            parsed=False,
        )
        saved.append(record.to_dict())
    return saved


def _prepare_debug_project(payload: dict[str, Any]) -> str:
    bundle = project_service.create_project(
        title="Debug Full Pipeline Project",
        files=[],
    )
    project_id = bundle.project.id

    assignment_brief = str(payload.get("assignment_brief") or "").strip()
    rubric = str(payload.get("rubric") or "").strip()
    lecture_notes = str(payload.get("lecture_notes") or "").strip()
    uploaded_files = payload.get("uploaded_files") or []

    if assignment_brief:
        project_service.add_file(
            project_id,
            file_type="assignment_brief",
            original_filename="assignment_brief.txt",
            storage_path=_write_debug_text_file(project_id, "assignment_brief.txt", assignment_brief),
            parsed=True,
        )
    if rubric:
        project_service.add_file(
            project_id,
            file_type="rubric",
            original_filename="rubric.txt",
            storage_path=_write_debug_text_file(project_id, "rubric.txt", rubric),
            parsed=True,
        )
    if lecture_notes:
        project_service.add_file(
            project_id,
            file_type="lecture_slides",
            original_filename="lecture_notes.txt",
            storage_path=_write_debug_text_file(project_id, "lecture_notes.txt", lecture_notes),
            parsed=True,
        )

    if isinstance(uploaded_files, list):
        for idx, entry in enumerate(uploaded_files, start=1):
            if isinstance(entry, str):
                text = entry.strip()
                if not text:
                    continue
                file_type = "additional_file"
                original_filename = f"uploaded_{idx}.txt"
            elif isinstance(entry, dict):
                text = str(entry.get("text") or "").strip()
                if not text:
                    continue
                file_type = str(entry.get("file_type") or "additional_file")
                original_filename = str(entry.get("filename") or f"uploaded_{idx}.txt")
            else:
                continue
            project_service.add_file(
                project_id,
                file_type=file_type,
                original_filename=original_filename,
                storage_path=_write_debug_text_file(project_id, original_filename, text),
                parsed=True,
            )
    return project_id


def _debug_run_stage(stage: str, project_id: str) -> dict[str, Any]:
    stage_start = _debug_now_iso()
    started = time.perf_counter()
    try:
        if stage == "requirement":
            output = project_service.analyze_requirements(project_id).requirement.to_dict()
        elif stage == "research":
            bundle = project_service.get_project(project_id)
            if not bundle.project.artifacts.get("payment_confirmed"):
                project_service.calculate_pricing(project_id, priority="standard")
                project_service.confirm_payment(project_id)
            output = project_service.run_research(project_id).to_dict()
        elif stage == "blueprint":
            output = project_service.run_blueprint(project_id).to_dict()
        elif stage == "writer":
            session = project_service.start_writer(project_id)
            while session.status.value == "active":
                session = project_service.advance_writer(project_id)
            output = project_service.merge_writer_draft(project_id).to_dict()
        else:
            raise ValueError(f"Unsupported stage: {stage}")

        ended = time.perf_counter()
        return {
            "success": True,
            "start_time": stage_start,
            "end_time": _debug_now_iso(),
            "duration_ms": int((ended - started) * 1000),
            "output": output,
        }
    except Exception as exc:  # noqa: BLE001
        ended = time.perf_counter()
        return {
            "success": False,
            "start_time": stage_start,
            "end_time": _debug_now_iso(),
            "duration_ms": int((ended - started) * 1000),
            "output": None,
            "error": {"message": str(exc), "traceback": traceback.format_exc()},
        }


@app.post("/api/debug/full-pipeline")
def api_debug_full_pipeline():
    try:
        payload, parsed_files, file_errors = _collect_debug_input_payload()
    except ValueError as exc:
        return jsonify({"error": str(exc)}), 400

    total_start = _debug_now_iso()
    started = time.perf_counter()
    errors: list[dict[str, Any]] = list(file_errors)
    stages: dict[str, Any] = {}

    try:
        project_id = _prepare_debug_project(payload)
    except Exception as exc:  # noqa: BLE001
        return (
            jsonify(
                {
                    "success": False,
                    "project_id": None,
                    "total_time_ms": int((time.perf_counter() - started) * 1000),
                    "models_used": {},
                    "stages": {},
                    "parsed_files": parsed_files,
                    "errors": [
                        {
                            "stage": "setup",
                            "message": str(exc),
                            "traceback": traceback.format_exc(),
                        }
                    ],
                }
            ),
            500,
        )

    stage_order = ["requirement", "research", "blueprint", "writer"]
    for stage in stage_order:
        result = _debug_run_stage(stage, project_id)
        stages[stage] = result
        if not result["success"]:
            errors.append(
                {
                    "stage": stage,
                    "message": result["error"]["message"],
                    "traceback": result["error"]["traceback"],
                }
            )
            break

    ended = time.perf_counter()

    models_used = {
        "requirement": (stages.get("requirement", {}).get("output") or {}).get("analyzer_version"),
        "research": (stages.get("research", {}).get("output") or {}).get("engine_version"),
        "blueprint": (stages.get("blueprint", {}).get("output") or {}).get("engine_version"),
        "writer": (stages.get("writer", {}).get("output") or {}).get("model"),
    }

    return jsonify(
        {
            "success": not errors,
            "project_id": project_id,
            "start_time": total_start,
            "end_time": _debug_now_iso(),
            "total_time_ms": int((ended - started) * 1000),
            "models_used": models_used,
            "stages": stages,
            "parsed_files": parsed_files,
            "errors": errors,
        }
    )


@app.post("/api/debug/run-stage")
def api_debug_run_stage():
    payload = request.get_json(silent=True) or {}
    if not isinstance(payload, dict):
        return jsonify({"error": "Invalid JSON payload"}), 400
    stage = str(payload.get("stage") or "").strip().lower()
    if stage not in {"requirement", "research", "blueprint", "writer"}:
        return jsonify({"error": "stage must be one of: requirement, research, blueprint, writer"}), 400

    project_id = str(payload.get("project_id") or "").strip()
    if not project_id:
        if stage != "requirement":
            return jsonify({"error": "project_id is required for stages other than requirement"}), 400
        project_id = _prepare_debug_project(payload)

    result = _debug_run_stage(stage, project_id)
    errors = []
    if not result["success"]:
        errors.append(
            {
                "stage": stage,
                "message": result["error"]["message"],
                "traceback": result["error"]["traceback"],
            }
        )

    return jsonify(
        {
            "success": result["success"],
            "project_id": project_id,
            "stage": stage,
            "result": result,
            "errors": errors,
        }
    )


@app.get("/api/debug/project/<project_id>")
def api_debug_project_get(project_id: str):
    data: dict[str, Any] = {"project_id": project_id}
    errors: list[dict[str, str]] = []

    try:
        data["requirement"] = project_service.get_project(project_id).requirement.to_dict()
    except Exception as exc:  # noqa: BLE001
        data["requirement"] = None
        errors.append({"stage": "requirement", "message": str(exc)})
    try:
        data["research"] = project_service.get_research_plan(project_id).to_dict()
    except Exception as exc:  # noqa: BLE001
        data["research"] = None
        errors.append({"stage": "research", "message": str(exc)})
    try:
        data["blueprint"] = project_service.get_blueprint(project_id).to_dict()
    except Exception as exc:  # noqa: BLE001
        data["blueprint"] = None
        errors.append({"stage": "blueprint", "message": str(exc)})
    try:
        data["draft"] = project_service.get_draft(project_id).to_dict()
    except Exception as exc:  # noqa: BLE001
        data["draft"] = None
        errors.append({"stage": "writer", "message": str(exc)})

    data["errors"] = errors
    return jsonify(data)


@app.post("/api/reviewer/review")
def api_reviewer_review():
    payload = request.get_json(silent=True) or {}
    for key in ("requirement_json", "research_plan", "blueprint", "draft"):
        if not isinstance(payload.get(key), dict):
            return jsonify({"error": f"{key} object is required"}), 400
    report = reviewer_engine.review_draft(
        requirement_json=payload["requirement_json"],
        research_plan=payload["research_plan"],
        blueprint=payload["blueprint"],
        draft=payload["draft"],
        project_id=payload.get("project_id"),
    )
    return jsonify(report.to_dict()), 201


@app.post("/api/assignment/projects/<project_id>/review")
def api_assignment_project_review(project_id: str):
    try:
        report = project_service.run_academic_review(project_id)
        bundle = project_service.get_project(project_id)
        project = bundle.project
    except KeyError:
        return jsonify({"error": "Project or draft not found"}), 404
    except ValueError as exc:
        return jsonify({"error": str(exc)}), 400
    return jsonify({
        "review_report": report.to_dict(),
        "pass_number": int(project.artifacts.get("review_pass_number", 1)),
        "issues_found": int(project.artifacts.get("last_review_issues_found", len(report.issues))),
        "issues_fixed": int(project.artifacts.get("last_issues_fixed", 0)),
    })


@app.get("/api/assignment/projects/<project_id>/review-report")
def api_assignment_project_review_report(project_id: str):
    try:
        report = project_service.get_review_report(project_id)
    except KeyError:
        return jsonify({"error": "Review report not found"}), 404
    return jsonify(report.to_dict())


@app.post("/api/revision/revise")
def api_revision_revise():
    payload = request.get_json(silent=True) or {}
    for key in ("requirement_json", "research_plan", "blueprint", "draft", "review_report"):
        if not isinstance(payload.get(key), dict):
            return jsonify({"error": f"{key} object is required"}), 400
    try:
        result = revision_engine.revise_draft(
            requirement_json=payload["requirement_json"],
            research_plan=payload["research_plan"],
            blueprint=payload["blueprint"],
            draft=payload["draft"],
            review_report=payload["review_report"],
            project_id=payload.get("project_id"),
        )
    except ValueError as exc:
        return jsonify({"error": str(exc)}), 400
    return jsonify(result.to_dict()), 201


@app.post("/api/assignment/projects/<project_id>/revision")
def api_assignment_project_revision(project_id: str):
    payload = request.get_json(silent=True) or {}
    review_report = payload.get("review_report")
    if review_report is not None and not isinstance(review_report, dict):
        return jsonify({"error": "review_report must be an object"}), 400
    try:
        result = project_service.run_revision(project_id, review_report=review_report)
    except KeyError as exc:
        return jsonify({"error": str(exc)}), 404
    except ValueError as exc:
        return jsonify({"error": str(exc)}), 400
    payload_out = result.to_dict() if hasattr(result, "to_dict") else result
    return jsonify({"revision_result": payload_out})


@app.get("/api/assignment/projects/<project_id>/revision-history")
def api_assignment_project_revision_history(project_id: str):
    history = project_service.get_revision_history(project_id)
    return jsonify(history.to_dict())


@app.post("/api/assignment/projects/<project_id>/draft/restore")
def api_assignment_project_restore_draft(project_id: str):
    payload = request.get_json(silent=True) or {}
    version = payload.get("version")
    if not isinstance(version, int):
        return jsonify({"error": "version integer is required"}), 400
    try:
        draft = project_service.restore_draft_version(project_id, version)
    except KeyError:
        return jsonify({"error": "Draft version not found"}), 404
    return jsonify({"draft": draft.to_dict()})


@app.post("/api/humanizer/run")
@economy_auth.email_verified_required
def api_humanizer_run():
    """Humanize via StealthWriter browser automation (legacy alias for the humanizer page)."""
    return api_browser_stealthwriter_humanize()


def _extract_flagged_sentences(raw: dict) -> list[str]:
    """Pull the AI-flagged sentence strings from a ZeroGPT detect payload."""
    data = raw.get("data") if isinstance(raw, dict) else None
    container = data if isinstance(data, dict) else raw
    for key in ("h", "highlighted_sentences", "highlightedSentences", "ai_sentences"):
        value = container.get(key) if isinstance(container, dict) else None
        if isinstance(value, list):
            return [str(s).strip() for s in value if str(s).strip()]
    return []


def _detect_number(container: dict, *keys) -> float | int | None:
    for key in keys:
        value = container.get(key)
        if isinstance(value, (int, float)):
            return value
    return None


@app.post("/api/workspace/detect")
@economy_auth.email_verified_required
def api_workspace_detect():
    """Run the existing ZeroGPT AI detector on editor text and return spans.

    This is a thin frontend adapter over ``zerogpt_client.detect`` — no browser
    automation is involved. Cost: 1 credit per 100 words (minimum 1).
    """
    payload = request.get_json(silent=True) or {}
    text = str(payload.get("text") or "").strip()
    if not text:
        return jsonify({"error": "text is required"}), 400
    if not _zerogpt_configured():
        return jsonify({"error": "ZeroGPT is not configured. Set ZEROGPT_API_KEY in .env"}), 503

    word_count = count_words(text)
    cost = feature_cost("detect", word_count=word_count)
    charged = _charge_current_user(
        "detect",
        cost,
        meta={"word_count": word_count, "pricing": "1_per_100_words"},
    )
    if not isinstance(charged, tuple):
        return charged
    user_id, _tx = charged

    started = time.monotonic()
    try:
        raw = zerogpt_client.detect(text)
    except (ZeroGPTError, ZeroGPTProviderError) as exc:
        _refund_safe(user_id, cost, "detect")
        return jsonify({"error": str(exc)}), 502
    except Exception as exc:  # noqa: BLE001
        _refund_safe(user_id, cost, "detect")
        app.logger.exception("workspace/detect failed")
        return jsonify({"error": f"Detection failed: {exc}"}), 502

    latency_ms = int((time.monotonic() - started) * 1000)
    data = raw.get("data") if isinstance(raw, dict) else {}
    container = data if isinstance(data, dict) else raw
    flagged = _extract_flagged_sentences(raw)
    ai_percentage = _detect_number(container, "fakePercentage", "aiPercentage", "ai_percentage", "score") or 0
    text_words = _detect_number(container, "textWords", "text_words", "words") or word_count
    ai_words = _detect_number(container, "aiWords", "ai_words")
    if ai_words is None:
        ai_words = sum(count_words(s) for s in flagged)

    _record_usage_safe(
        user_id,
        feature=FEATURE_DETECTION,
        credits_used=cost,
        provider="ZeroGPT",
        latency=latency_ms,
        request_id=None,
    )

    # Passive detector corpus: AI > 20% with sentence-level highlights.
    try:
        if float(ai_percentage) > 20.0 and flagged:
            from services.dataset_logger import infer_human_segments, log_detection_event

            log_detection_event(
                user_id,
                text,
                float(ai_percentage),
                flagged,
                infer_human_segments(text, flagged),
                "auto_report_over_20",
            )
    except Exception:
        app.logger.exception("dataset_logger workspace detect hook failed")

    return jsonify(
        {
            "ai_percentage": round(float(ai_percentage), 1),
            "is_ai": float(ai_percentage) >= 50.0,
            "flagged_sentences": flagged,
            "flagged_parts": len(flagged),
            "text_words": int(text_words),
            "ai_words": int(ai_words),
            "credits_charged": cost,
            "coins_charged": cost,
            "balance": wallet.get_balance(user_id),
        }
    )


@app.post("/api/workspace/citations/search")
@economy_auth.email_verified_required
def api_workspace_citations_search():
    """Search scholarly works via the CitationService (Crossref provider).

    The frontend never talks to Crossref directly and never sees provider
    fields — only normalized works plus formatted in-text/reference strings.
    Search itself is free; inserting a citation charges via ``/citations/use``.
    """
    payload = request.get_json(silent=True) or {}
    query = str(payload.get("query") or "").strip()
    style = str(payload.get("style") or "APA 7").strip()
    try:
        limit = int(payload.get("limit") or 6)
    except (TypeError, ValueError):
        limit = 6
    if not query:
        return jsonify({"error": "query is required"}), 400

    try:
        result = citation_service.search(query, style=style, limit=limit)
    except Exception as exc:  # noqa: BLE001 — provider network/parse errors
        app.logger.warning("citation search failed: %s", exc)
        return jsonify({"error": f"Citation search failed: {exc}"}), 502

    return jsonify(result)


@app.post("/api/workspace/citations/use")
@economy_auth.email_verified_required
def api_workspace_citations_use():
    """Charge 2 credits when a citation is inserted / reference added."""
    payload = request.get_json(silent=True) or {}
    action = str(payload.get("action") or "insert").strip().lower()
    if action not in ("insert", "reference", "add_reference"):
        action = "insert"

    cost = feature_cost("cite")
    charged = _charge_current_user(
        "cite",
        cost,
        meta={"action": action, "doi": str(payload.get("doi") or "")[:120] or None},
    )
    if not isinstance(charged, tuple):
        return charged
    user_id, _tx = charged
    return jsonify(
        {
            "success": True,
            "credits_charged": cost,
            "coins_charged": cost,
            "balance": wallet.get_balance(user_id),
        }
    )


@app.post("/api/humanizer/session")
def api_humanizer_session_create():
    payload = request.get_json(silent=True) or {}
    for key in ("draft", "requirement_json", "blueprint"):
        if not isinstance(payload.get(key), dict):
            return jsonify({"error": f"{key} object is required"}), 400
    try:
        session = humanizer_engine.create_session(
            draft=payload["draft"],
            requirement_json=payload["requirement_json"],
            blueprint=payload["blueprint"],
            project_id=payload.get("project_id"),
        )
    except ValueError as exc:
        return jsonify({"error": str(exc)}), 400
    return jsonify(session.to_dict()), 201


@app.get("/api/humanizer/session/<session_id>")
def api_humanizer_session_get(session_id: str):
    try:
        session = humanizer_engine.get_session(session_id)
    except KeyError:
        return jsonify({"error": "Humanizer session not found"}), 404
    return jsonify(session.to_dict())


@app.post("/api/humanizer/session/<session_id>/advance")
def api_humanizer_session_advance(session_id: str):
    try:
        session = humanizer_engine.advance_paragraph(session_id)
    except KeyError:
        return jsonify({"error": "Humanizer session not found"}), 404
    except ValueError as exc:
        return jsonify({"error": str(exc)}), 502
    return jsonify(session.to_dict())


@app.post("/api/humanizer/session/<session_id>/merge")
def api_humanizer_session_merge(session_id: str):
    payload = request.get_json(silent=True) or {}
    try:
        draft = humanizer_engine.merge_humanized_draft(session_id, title=payload.get("title"))
    except KeyError:
        return jsonify({"error": "Humanizer session not found"}), 404
    except ValueError as exc:
        return jsonify({"error": str(exc)}), 400
    return jsonify(draft.to_dict()), 201


@app.post("/api/assignment/projects/<project_id>/humanizer/start")
def api_assignment_humanizer_start(project_id: str):
    try:
        session = project_service.start_humanizer(project_id)
    except KeyError:
        return jsonify({"error": "Project or draft not found"}), 404
    except ValueError as exc:
        return jsonify({"error": str(exc)}), 400
    return jsonify(session.to_dict()), 201


@app.post("/api/assignment/projects/<project_id>/humanizer/advance")
def api_assignment_humanizer_advance(project_id: str):
    payload = request.get_json(silent=True) or {}
    humanizer_session = payload.get("humanizer_session")
    if humanizer_session is not None and not isinstance(humanizer_session, dict):
        return jsonify({"error": "humanizer_session must be an object"}), 400
    try:
        session = project_service.advance_humanizer(project_id, humanizer_session=humanizer_session)
    except KeyError:
        return jsonify({"error": "Humanizer session not found"}), 404
    except ValueError as exc:
        return jsonify({"error": str(exc)}), 502
    try:
        from services.dataset_logger import log_humanization_event

        uid = economy_auth.current_user_id()
        completed = list(getattr(session, "completed_paragraph_ids", None) or [])
        if uid and completed:
            last_id = completed[-1]
            try:
                para = session.paragraph_by_id(last_id)
            except KeyError:
                para = None
            if para is not None:
                log_humanization_event(
                    int(uid),
                    "assignment",
                    getattr(para, "original_text", None),
                    getattr(para, "humanized_text", None),
                )
    except Exception:
        app.logger.exception("dataset_logger assignment hook failed")
    return jsonify(session.to_dict())


@app.post("/api/assignment/projects/<project_id>/humanizer/merge")
def api_assignment_humanizer_merge(project_id: str):
    payload = request.get_json(silent=True) or {}
    humanizer_session = payload.get("humanizer_session")
    if humanizer_session is not None and not isinstance(humanizer_session, dict):
        return jsonify({"error": "humanizer_session must be an object"}), 400
    try:
        draft = project_service.merge_humanized_draft(project_id, humanizer_session=humanizer_session)
    except KeyError:
        return jsonify({"error": "Humanizer session not found"}), 404
    except ValueError as exc:
        return jsonify({"error": str(exc)}), 400
    return jsonify(draft.to_dict()), 201


@app.get("/api/assignment/projects/<project_id>/humanizer")
def api_assignment_humanizer_get(project_id: str):
    try:
        session = project_service.get_humanizer_session(project_id)
    except KeyError:
        return jsonify({"error": "Humanizer session not found"}), 404
    except Exception as exc:  # noqa: BLE001
        trace(
            "api.humanizer.get.error",
            project_id=project_id,
            error=str(exc),
            error_type=type(exc).__name__,
        )
        return jsonify({"error": "Failed to load humanizer session. Please retry."}), 500
    return jsonify(session.to_dict())


@app.get("/api/assignment/projects/<project_id>/humanized-draft")
def api_assignment_humanized_draft_get(project_id: str):
    try:
        draft = project_service.get_humanized_draft(project_id)
    except KeyError:
        return jsonify({"error": "Humanized draft not found"}), 404
    return jsonify(draft.to_dict())


@app.post("/api/ai-detection/session")
def api_ai_detection_session_create():
    payload = request.get_json(silent=True) or {}
    for key in ("humanized_draft", "requirement_json"):
        if not isinstance(payload.get(key), dict):
            return jsonify({"error": f"{key} object is required"}), 400
    try:
        from services.ai_detection_engine.models import DetectionThresholds

        thresholds = DetectionThresholds.from_dict(payload.get("thresholds"))
        session = ai_detection_engine.create_session(
            humanized_draft=payload["humanized_draft"],
            requirement_json=payload["requirement_json"],
            project_id=payload.get("project_id"),
            thresholds=thresholds,
            humanizer_paragraph_ids=payload.get("humanizer_paragraph_ids"),
        )
    except ValueError as exc:
        return jsonify({"error": str(exc)}), 400
    return jsonify(session.to_dict()), 201


@app.get("/api/ai-detection/session/<session_id>")
def api_ai_detection_session_get(session_id: str):
    try:
        session = ai_detection_engine.get_session(session_id)
    except KeyError:
        return jsonify({"error": "Detection session not found"}), 404
    return jsonify(session.to_dict())


@app.post("/api/ai-detection/session/<session_id>/advance")
def api_ai_detection_session_advance(session_id: str):
    try:
        session = ai_detection_engine.advance_paragraph(session_id)
    except KeyError:
        return jsonify({"error": "Detection session not found"}), 404
    return jsonify(session.to_dict())


@app.post("/api/ai-detection/session/<session_id>/finalize")
def api_ai_detection_session_finalize(session_id: str):
    try:
        session = ai_detection_engine.finalize_session(session_id)
    except KeyError:
        return jsonify({"error": "Detection session not found"}), 404
    report = ai_detection_engine.get_report(session.report_id)
    return jsonify({"session": session.to_dict(), "detection_report": report.to_dict()})


@app.post("/api/assignment/projects/<project_id>/ai-detection/start")
def api_assignment_ai_detection_start(project_id: str):
    try:
        session = project_service.start_ai_detection(project_id)
    except KeyError:
        return jsonify({"error": "Project or humanized draft not found"}), 404
    except ValueError as exc:
        return jsonify({"error": str(exc)}), 400
    return jsonify(session.to_dict()), 201


@app.post("/api/assignment/projects/<project_id>/ai-detection/advance")
def api_assignment_ai_detection_advance(project_id: str):
    payload = request.get_json(silent=True) or {}
    detection_session = payload.get("detection_session")
    if detection_session is not None and not isinstance(detection_session, dict):
        return jsonify({"error": "detection_session must be an object"}), 400
    try:
        session = project_service.advance_ai_detection(project_id, detection_session=detection_session)
    except KeyError as exc:
        return jsonify({"error": str(exc) or "Detection session not found"}), 404
    except ValueError as exc:
        return jsonify({"error": str(exc)}), 502
    except Exception as exc:  # noqa: BLE001
        trace(
            "api.detection.advance.error",
            project_id=project_id,
            error=str(exc),
            error_type=type(exc).__name__,
        )
        return jsonify({"error": str(exc) or "Detection step failed. Please try again."}), 502
    return jsonify(session.to_dict())


@app.post("/api/assignment/projects/<project_id>/ai-detection/finalize")
def api_assignment_ai_detection_finalize(project_id: str):
    payload = request.get_json(silent=True) or {}
    detection_session = payload.get("detection_session")
    if detection_session is not None and not isinstance(detection_session, dict):
        return jsonify({"error": "detection_session must be an object"}), 400
    try:
        report = project_service.finalize_ai_detection(project_id, detection_session=detection_session)
    except KeyError as exc:
        return jsonify({"error": str(exc) or "Detection session not found"}), 404
    except ValueError as exc:
        return jsonify({"error": str(exc)}), 400
    except Exception as exc:  # noqa: BLE001
        trace(
            "api.detection.finalize.error",
            project_id=project_id,
            error=str(exc),
            error_type=type(exc).__name__,
        )
        return jsonify({"error": "Detection finalize failed. Please try again."}), 502
    return jsonify({"detection_report": report.to_dict()})


@app.post("/api/assignment/projects/<project_id>/ai-detection/prepare-retry")
def api_assignment_ai_detection_prepare_retry(project_id: str):
    try:
        project_service.prepare_detection_retry(project_id)
    except KeyError as exc:
        return jsonify({"error": str(exc)}), 404
    except ValueError as exc:
        return jsonify({"error": str(exc)}), 400
    return jsonify({"ok": True})


@app.get("/api/assignment/projects/<project_id>/ai-detection")
def api_assignment_ai_detection_get(project_id: str):
    try:
        session = project_service.get_detection_session(project_id)
    except KeyError:
        return jsonify({"error": "Detection session not found"}), 404
    return jsonify(session.to_dict())


@app.get("/api/assignment/projects/<project_id>/detection-report")
def api_assignment_detection_report_get(project_id: str):
    try:
        report = project_service.get_detection_report(project_id)
    except KeyError:
        return jsonify({"error": "Detection report not found"}), 404
    return jsonify(report.to_dict())


@app.post("/api/delivery/package")
def api_delivery_package_prepare():
    """Standalone Delivery Engine — packages prior pipeline outputs only."""
    payload = request.get_json(silent=True) or {}
    final_draft = payload.get("final_draft")
    requirement_json = payload.get("requirement_json")
    research_plan = payload.get("research_plan")
    blueprint = payload.get("blueprint")
    review_report = payload.get("review_report")
    detection_report = payload.get("detection_report")
    if not all(
        isinstance(item, dict)
        for item in (
            final_draft,
            requirement_json,
            research_plan,
            blueprint,
            review_report,
            detection_report,
        )
    ):
        return jsonify({"error": "All pipeline artifacts are required as objects"}), 400
    package = delivery_engine.prepare_package(
        final_draft=final_draft,
        requirement_json=requirement_json,
        research_plan=research_plan,
        blueprint=blueprint,
        review_report=review_report,
        detection_report=detection_report,
        project_id=payload.get("project_id"),
        revision_attempts=int(payload.get("revision_attempts", 0)),
        humanization_attempts=int(payload.get("humanization_attempts", 0)),
        completion_time=payload.get("completion_time"),
    )
    return jsonify(package.to_dict())


@app.post("/api/assignment/projects/<project_id>/citations/generate")
def api_assignment_citations_generate(project_id: str):
    try:
        pack = project_service.run_citation_generation(project_id)
    except KeyError as exc:
        return jsonify({"error": str(exc)}), 404
    except Exception as exc:  # noqa: BLE001
        return jsonify({"error": str(exc)}), 502
    return jsonify({"citation_pack": pack.to_dict()})


@app.post("/api/assignment/projects/<project_id>/format")
def api_assignment_format(project_id: str):
    try:
        formatted = project_service.run_formatting(project_id)
    except KeyError as exc:
        return jsonify({"error": str(exc)}), 404
    except Exception as exc:  # noqa: BLE001
        return jsonify({"error": str(exc)}), 500
    return jsonify({"formatted_document": formatted})


@app.post("/api/assignment/projects/<project_id>/validate-requirements")
def api_assignment_validate_requirements(project_id: str):
    try:
        report = project_service.run_requirement_validation(project_id)
    except KeyError as exc:
        return jsonify({"error": str(exc)}), 404
    except Exception as exc:  # noqa: BLE001
        return jsonify({"error": str(exc)}), 500
    # Always 200: passed/failed is business data in validation_report, not HTTP conflict.
    return jsonify({"validation_report": report})


@app.post("/api/assignment/projects/<project_id>/stages/<stage_name>/retry")
def api_assignment_stage_retry(project_id: str, stage_name: str):
    stage = _parse_pipeline_stage(stage_name)
    if stage is None:
        return jsonify({"error": f"Unknown stage: {stage_name}"}), 400
    try:
        result = project_service.retry_stage(project_id, stage)
    except KeyError as exc:
        return jsonify({"error": str(exc)}), 404
    except ValueError as exc:
        return jsonify({"error": str(exc)}), 400
    except Exception as exc:  # noqa: BLE001
        return jsonify({"error": str(exc)}), 500
    if hasattr(result, "to_dict"):
        payload = result.to_dict()
    elif isinstance(result, dict):
        payload = result
    else:
        payload = {"ok": True}
    return jsonify({"stage": stage.value, "result": payload})


@app.post("/api/assignment/projects/<project_id>/delivery")
def api_assignment_delivery_prepare(project_id: str):
    try:
        package = project_service.run_delivery(project_id)
    except KeyError as exc:
        return jsonify({"error": str(exc)}), 404
    except ValueError as exc:
        return jsonify({"error": str(exc)}), 400
    except Exception as exc:  # noqa: BLE001
        trace(
            "api.delivery.error",
            project_id=project_id,
            error=str(exc),
            error_type=type(exc).__name__,
        )
        return jsonify({"error": f"Delivery packaging failed: {exc}"}), 500
    return jsonify(package.to_dict())


@app.get("/api/assignment/projects/<project_id>/delivery-package")
def api_assignment_delivery_package_get(project_id: str):
    try:
        package = project_service.get_delivery_package(project_id)
    except KeyError:
        return jsonify({"error": "Delivery package not found"}), 404
    return jsonify(package.to_dict())


@app.get("/api/assignment/projects/<project_id>/revision-chat")
@economy_auth.login_required
def api_assignment_revision_chat_get(project_id: str):
    try:
        bundle = project_service.get_project(project_id)
    except KeyError as exc:
        return jsonify({"error": str(exc)}), 404
    user_id = economy_auth.current_user_id()
    if bundle.project.user_id and str(bundle.project.user_id) != str(user_id):
        return jsonify({"error": "Forbidden"}), 403
    return jsonify(project_service.get_revision_chat(project_id))


@app.put("/api/assignment/projects/<project_id>/chat-transcript")
@economy_auth.login_required
def api_assignment_chat_transcript_put(project_id: str):
    """Persist the Assignment chat thread so history reloads on open."""
    try:
        bundle = project_service.get_project(project_id)
    except KeyError as exc:
        return jsonify({"error": str(exc)}), 404
    user_id = economy_auth.current_user_id()
    if bundle.project.user_id and str(bundle.project.user_id) != str(user_id):
        return jsonify({"error": "Forbidden"}), 403
    payload = request.get_json(silent=True) or {}
    messages = payload.get("messages")
    if messages is None:
        messages = []
    try:
        saved = project_service.save_chat_transcript(project_id, messages)
    except ValueError as exc:
        return jsonify({"error": str(exc)}), 400
    return jsonify({"ok": True, "messages": saved, "count": len(saved)})


@app.post("/api/assignment/projects/<project_id>/revision-chat")
@economy_auth.login_required
def api_assignment_revision_chat_post(project_id: str):
    try:
        bundle = project_service.get_project(project_id)
    except KeyError as exc:
        return jsonify({"error": str(exc)}), 404
    user_id = economy_auth.current_user_id()
    if bundle.project.user_id and str(bundle.project.user_id) != str(user_id):
        return jsonify({"error": "Forbidden"}), 403
    payload = request.get_json(silent=True) or {}
    message = str(payload.get("message") or "").strip()
    try:
        result = project_service.apply_client_revision(project_id, message)
    except ValueError as exc:
        return jsonify({"error": str(exc)}), 400
    except Exception as exc:  # noqa: BLE001
        trace(
            "api.revision_chat.error",
            project_id=project_id,
            error=str(exc),
            error_type=type(exc).__name__,
        )
        return jsonify({"error": f"Revision failed: {exc}"}), 500
    return jsonify(result)


@app.get("/api/assignment/projects/<project_id>/download")
def api_assignment_project_download(project_id: str):
    """Download the client deliverable (docx or pdf only — no JSON package)."""
    try:
        package = project_service.get_delivery_package(project_id)
    except KeyError:
        try:
            package = project_service.run_delivery(project_id)
        except ValueError as exc:
            return jsonify({"error": str(exc)}), 400
        except KeyError as exc:
            return jsonify({"error": str(exc)}), 404
        except Exception as exc:  # noqa: BLE001
            trace(
                "api.assignment_download.error",
                project_id=project_id,
                error=str(exc),
                error_type=type(exc).__name__,
            )
            return jsonify({"error": f"Delivery packaging failed: {exc}"}), 500

    file_path, mime, download_name = _resolve_client_deliverable(package, project_id)
    if file_path is None or not file_path.exists():
        try:
            package = project_service.run_delivery(project_id)
        except Exception as exc:  # noqa: BLE001
            return jsonify({"error": f"Deliverable is not available: {exc}"}), 404
        file_path, mime, download_name = _resolve_client_deliverable(package, project_id)
    if file_path is None or not file_path.exists():
        return jsonify({"error": "Deliverable is not available"}), 404
    return send_file(
        file_path,
        mimetype=mime or "application/octet-stream",
        as_attachment=True,
        download_name=download_name or file_path.name,
    )


@app.get("/api/delivery/packages/<package_id>/download")
def api_delivery_package_download(package_id: str):
    try:
        package = project_service.find_delivery_package(package_id)
    except KeyError:
        return jsonify({"error": "Package not found"}), 404
    file_path, mime, download_name = _resolve_client_deliverable(
        package, str(package.project_id or "local")
    )
    if file_path is None or not file_path.exists():
        return jsonify({"error": "Deliverable is not available"}), 404
    return send_file(
        file_path,
        mimetype=mime or "application/octet-stream",
        as_attachment=True,
        download_name=download_name or file_path.name,
    )


@app.get("/api/delivery/files/<file_id>")
def api_delivery_file_download(file_id: str):
    try:
        file_record = delivery_engine.get_file(file_id)
    except KeyError:
        return jsonify({"error": "File not found"}), 404
    file_path = Path(file_record.storage_path)
    if not file_path.exists():
        return jsonify({"error": "Delivery file is not available"}), 404
    return send_file(
        file_path,
        mimetype=file_record.mime_type or "application/octet-stream",
        as_attachment=True,
        download_name=file_record.filename,
    )


def _resolve_client_deliverable(package, project_id: str):
    """Return (path, mime, download_name) for the single client-facing file."""
    project_dir = PROJECT_STORAGE_ROOT / str(getattr(package, "project_id", None) or project_id) / "delivery"
    client_name = getattr(package, "client_filename", None)
    if client_name:
        candidate = project_dir / str(client_name)
        if candidate.is_file():
            mime = "application/pdf" if candidate.suffix.lower() == ".pdf" else (
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            )
            return candidate, mime, candidate.name

    for entry in getattr(package, "files", None) or []:
        ftype = str(getattr(entry, "file_type", "") or "")
        if ftype not in {"final_assignment_docx", "final_assignment_pdf"}:
            continue
        path = Path(getattr(entry, "storage_path", "") or "")
        if path.is_file():
            return path, getattr(entry, "mime_type", None), getattr(entry, "filename", None) or path.name

    # Fallbacks for older packages (prefer docx over zip/json).
    for pattern in ("*.docx", "*.pdf"):
        matches = sorted(p for p in project_dir.glob(pattern) if p.is_file())
        if matches:
            path = matches[0]
            mime = "application/pdf" if path.suffix.lower() == ".pdf" else (
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            )
            return path, mime, path.name
    return None, None, None


def _safe_download_name(package) -> str:
    if package and getattr(package, "project_summary", None):
        raw = str(getattr(package.project_summary, "project_name", "") or "")
        cleaned = re.sub(r"[^A-Za-z0-9._-]+", "-", raw.strip()).strip("-")
        if cleaned:
            return cleaned
    return "Assignment"


@app.post("/api/research/plan")
def api_research_plan_build():
    """Standalone Research Engine entrypoint — no raw files, only parsed inputs."""
    payload = request.get_json(silent=True) or {}
    requirement_json = payload.get("requirement_json")
    if not isinstance(requirement_json, dict):
        return jsonify({"error": "requirement_json object is required"}), 400
    parsed_raw = payload.get("parsed_documents") or []
    if not isinstance(parsed_raw, list):
        return jsonify({"error": "parsed_documents must be an array"}), 400
    documents = [ParsedDocument.from_dict(item) for item in parsed_raw]
    plan = research_engine.build_plan(
        requirement_json=requirement_json,
        parsed_documents=documents,
        project_id=payload.get("project_id"),
    )
    return jsonify(plan.to_dict()), 201


@app.get("/api/assignment/projects/<project_id>/pipeline")
def api_assignment_project_pipeline(project_id: str):
    try:
        pipeline_project = assignment_pipeline.get_project(project_id)
    except KeyError:
        return jsonify({"error": "Project not found"}), 404
    return jsonify(pipeline_project.to_dict())


@app.get("/api/project/<project_id>/timeline")
def api_project_timeline(project_id: str):
    try:
        timeline = project_service.get_project_timeline(project_id)
    except KeyError:
        return jsonify({"error": "Project timeline not found"}), 404
    return jsonify({"project_id": project_id, "timeline": timeline})


@app.get("/api/project/<project_id>/status")
def api_project_status(project_id: str):
    try:
        status = project_service.get_project_lifecycle_status(project_id)
    except KeyError:
        return jsonify({"error": "Project status not found"}), 404
    return jsonify({"project_id": project_id, **status})


@app.post("/api/assignment/projects/<project_id>/stages/<stage_name>/start")
def api_assignment_stage_start(project_id: str, stage_name: str):
    stage = _parse_pipeline_stage(stage_name)
    if stage is None:
        return jsonify({"error": f"Unknown stage: {stage_name}"}), 400
    try:
        project = assignment_pipeline.start_stage(project_id, stage)
    except KeyError:
        return jsonify({"error": "Project not found"}), 404
    project_service.sync_pipeline_state(project_id)
    return jsonify(project.to_dict())


@app.post("/api/assignment/projects/<project_id>/stages/<stage_name>/complete")
def api_assignment_stage_complete(project_id: str, stage_name: str):
    stage = _parse_pipeline_stage(stage_name)
    if stage is None:
        return jsonify({"error": f"Unknown stage: {stage_name}"}), 400
    payload = request.get_json(silent=True) or {}
    result = StageResult(
        output=payload.get("output") if isinstance(payload.get("output"), dict) else {},
        artifacts=payload.get("artifacts") if isinstance(payload.get("artifacts"), dict) else {},
        requirement_json=payload.get("requirement_json")
        if isinstance(payload.get("requirement_json"), dict)
        else None,
        pricing=payload.get("pricing") if isinstance(payload.get("pricing"), dict) else None,
    )
    try:
        project = assignment_pipeline.complete_stage(project_id, stage, result)
    except KeyError:
        return jsonify({"error": "Project not found"}), 404
    project_service.sync_pipeline_state(project_id)
    return jsonify(project.to_dict())


@app.post("/api/assignment/projects/<project_id>/stages/<stage_name>/fail")
def api_assignment_stage_fail(project_id: str, stage_name: str):
    stage = _parse_pipeline_stage(stage_name)
    if stage is None:
        return jsonify({"error": f"Unknown stage: {stage_name}"}), 400
    payload = request.get_json(silent=True) or {}
    error = (payload.get("error") or "Stage failed").strip()
    try:
        project = assignment_pipeline.fail_stage(project_id, stage, error)
    except KeyError:
        return jsonify({"error": "Project not found"}), 404
    project_service.sync_pipeline_state(project_id)
    return jsonify(project.to_dict())


@app.post("/api/assignment/projects/<project_id>/run")
def api_assignment_stage_run(project_id: str):
    """Run a registered stage handler when wired; otherwise marks stage as running."""
    payload = request.get_json(silent=True) or {}
    stage = None
    if payload.get("stage"):
        stage = _parse_pipeline_stage(str(payload["stage"]))
        if stage is None:
            return jsonify({"error": f"Unknown stage: {payload['stage']}"}), 400
    try:
        project = assignment_pipeline.run_stage(project_id, stage)
    except KeyError:
        return jsonify({"error": "Project not found"}), 404
    project_service.sync_pipeline_state(project_id)
    return jsonify(project.to_dict())


@app.route("/turnitin")
@economy_auth.email_verified_required
def turnitin():
    return render_template(
        "turnitin.html",
        nav_active="turnitin",
        turnitin_cost=feature_cost("turnitin"),
    )


def _turnitin_row_api(row: dict[str, Any]) -> dict[str, Any]:
    return turnitin_service.to_api_row(row)


@app.get("/api/turnitin/reports")
@economy_auth.email_verified_required
def api_turnitin_reports():
    user_id = economy_auth.current_user_id()
    rows = turnitin_service.store.list_for_user(user_id)
    return jsonify({"success": True, "reports": [_turnitin_row_api(r) for r in rows]})


@app.get("/api/turnitin/submissions/<submission_id>")
@economy_auth.email_verified_required
def api_turnitin_submission(submission_id: str):
    user_id = economy_auth.current_user_id()
    row = turnitin_service.store.get_for_user(submission_id, user_id)
    if row is None:
        return jsonify({"success": False, "error": "Not found"}), 404
    return jsonify({"success": True, "report": _turnitin_row_api(row)})


@app.delete("/api/turnitin/submissions/<submission_id>")
@economy_auth.email_verified_required
def api_turnitin_delete(submission_id: str):
    user_id = economy_auth.current_user_id()
    if not turnitin_service.store.delete_for_user(submission_id, user_id):
        return jsonify({"success": False, "error": "Not found"}), 404
    return jsonify({"success": True})


@app.post("/api/turnitin/submissions/<submission_id>/highlights")
@economy_auth.email_verified_required
def api_turnitin_request_highlights(submission_id: str):
    """Queue a PlagDetect AI Highlights job (optional, user-initiated)."""
    user_id = economy_auth.current_user_id()
    row = turnitin_service.store.get_for_user(submission_id, user_id)
    if row is None:
        return jsonify({"success": False, "error": "Not found"}), 404
    if row.get("status") != "completed":
        return jsonify({"success": False, "error": "Wait for the check to finish first."}), 400
    if not row.get("external_id"):
        return jsonify({"success": False, "error": "PlagDetect submission id missing."}), 400

    hl_status = (row.get("highlights_status") or "").strip().lower()
    if hl_status in ("queued", "running"):
        return jsonify({"success": False, "error": "AI Highlights already in progress."}), 409
    if row.get("has_highlights_report"):
        return jsonify({"success": True, "report": _turnitin_row_api(row), "already": True})

    meta = row.get("meta") or {}
    # Allow re-fetch when score exists but PDF is missing.
    if row.get("ai_highlights") is not None and not row.get("has_highlights_report"):
        pass
    elif meta.get("ai_score_display") != "*%":
        return jsonify(
            {
                "success": False,
                "error": "AI Highlights are available when the AI score shows *%.",
            }
        ), 400

    try:
        from services.browser.providers import plagdetect as pd

        pd._validate_urls()
        report_dir = str(turnitin_service.report_dir(submission_id))
        ensure_engine_started()
        job = job_manager.create(
            "plagdetect",
            "highlights",
            {
                "external_id": row["external_id"],
                "report_dir": report_dir,
                "submission_id": submission_id,
            },
            max_retries=0,
        )
        turnitin_service.store.update(
            submission_id,
            highlights_status="queued",
            highlights_job_id=job.id,
        )
        turnitin_service.watch_highlights_job(
            submission_id=submission_id,
            job_id=job.id,
            job_manager=job_manager,
        )
    except Exception as exc:  # noqa: BLE001
        turnitin_service.store.update(submission_id, highlights_status="failed")
        app.logger.exception("turnitin/highlights failed")
        return jsonify({"success": False, "error": f"Could not queue highlights: {exc}"}), 500

    updated = turnitin_service.store.get_for_user(submission_id, user_id)
    return jsonify(
        {
            "success": True,
            "submission_id": submission_id,
            "highlights_status": "queued",
            "job_id": job.id,
            "report": _turnitin_row_api(updated or row),
        }
    )


@app.post("/api/turnitin/submissions/<submission_id>/fetch-reports")
@economy_auth.email_verified_required
def api_turnitin_fetch_reports(submission_id: str):
    """Re-download missing Similarity / AI / Highlights PDFs from PlagDetect."""
    user_id = economy_auth.current_user_id()
    row = turnitin_service.store.get_for_user(submission_id, user_id)
    if row is None:
        return jsonify({"success": False, "error": "Not found"}), 404
    if row.get("status") != "completed":
        return jsonify({"success": False, "error": "Wait for the check to finish first."}), 400
    if not row.get("external_id"):
        return jsonify({"success": False, "error": "PlagDetect submission id missing."}), 400

    body = request.get_json(silent=True) or {}
    want_sim = body.get("similarity", True)
    want_ai = body.get("ai", True)
    want_hl = body.get("highlights", False)

    fetch_similarity = bool(want_sim) and not row.get("has_similarity_report")
    fetch_ai = bool(want_ai) and not row.get("has_ai_report")
    fetch_highlights = bool(want_hl) and not row.get("has_highlights_report")

    if not fetch_similarity and not fetch_ai and not fetch_highlights:
        return jsonify({"success": True, "report": _turnitin_row_api(row), "already": True})

    try:
        from services.browser.providers import plagdetect as pd

        pd._validate_urls()
        report_dir = str(turnitin_service.report_dir(submission_id))
        ensure_engine_started()
        job = job_manager.create(
            "plagdetect",
            "fetch_reports",
            {
                "external_id": row["external_id"],
                "report_dir": report_dir,
                "submission_id": submission_id,
                "fetch_similarity": fetch_similarity,
                "fetch_ai": fetch_ai,
                "fetch_highlights": fetch_highlights,
            },
            max_retries=0,
        )
        turnitin_service.watch_fetch_reports_job(
            submission_id=submission_id,
            job_id=job.id,
            job_manager=job_manager,
        )
    except Exception as exc:  # noqa: BLE001
        app.logger.exception("turnitin/fetch-reports failed")
        return jsonify({"success": False, "error": f"Could not queue report download: {exc}"}), 500

    return jsonify(
        {
            "success": True,
            "submission_id": submission_id,
            "job_id": job.id,
            "fetching": {
                "similarity": fetch_similarity,
                "ai": fetch_ai,
                "highlights": fetch_highlights,
            },
            "report": _turnitin_row_api(row),
        }
    )


@app.get("/api/turnitin/submissions/<submission_id>/report/<kind>")
@economy_auth.email_verified_required
def api_turnitin_report_download(submission_id: str, kind: str):
    from services.turnitin_service.store import resolve_report_path

    user_id = economy_auth.current_user_id()
    row = turnitin_service.store.get_for_user(submission_id, user_id)
    if row is None:
        return jsonify({"success": False, "error": "Not found"}), 404
    if kind not in ("similarity", "ai", "highlights"):
        return jsonify({"success": False, "error": "Invalid report type."}), 400
    key_map = {
        "similarity": "similarity_report_path",
        "ai": "ai_report_path",
        "highlights": "ai_highlights_report_path",
    }
    key = key_map[kind]
    path = resolve_report_path(row.get(key))
    if not path:
        return jsonify({"success": False, "error": "Report not ready yet."}), 404
    download_name = f"{Path(row['filename']).stem}_{kind}_report.pdf"
    return send_file(path, mimetype="application/pdf", as_attachment=True, download_name=download_name)


@app.post("/api/turnitin/check")
@economy_auth.email_verified_required
def api_turnitin_check():
    """Upload a document, charge coins, and queue a PlagDetect browser job."""
    user_id = economy_auth.current_user_id()
    f = request.files.get("file")
    if f is None or not f.filename:
        return jsonify({"success": False, "error": "A document is required."}), 400

    filename = Path(f.filename).name
    exclude_bibliography = request.form.get("exclude_bibliography", "0") in ("1", "true", "on", "yes")
    exclude_quotes = request.form.get("exclude_quotes", "0") in ("1", "true", "on", "yes")

    cost = feature_cost("turnitin")
    submission_id = uuid.uuid4().hex[:12]
    charged = _charge_current_user("turnitin", cost, ref_id=submission_id, meta={"filename": filename})
    if not isinstance(charged, tuple):
        return charged
    user_id, _tx = charged

    try:
        from services.browser.providers import plagdetect as pd

        pd._validate_urls()
        upload_path = turnitin_service.save_upload(submission_id, filename, f.read())
        report_dir = str(turnitin_service.report_dir(submission_id))
        ensure_engine_started()
        job = job_manager.create(
            "plagdetect",
            "check",
            {
                "file_path": upload_path,
                "exclude_bibliography": exclude_bibliography,
                "exclude_quotes": exclude_quotes,
                "report_dir": report_dir,
                "submission_id": submission_id,
            },
            max_retries=0,
        )
        turnitin_service.store.create(
            submission_id=submission_id,
            user_id=user_id,
            filename=filename,
            upload_path=upload_path,
            exclude_bibliography=exclude_bibliography,
            exclude_quotes=exclude_quotes,
            job_id=job.id,
        )
        turnitin_service.watch_job(
            submission_id=submission_id,
            job_id=job.id,
            user_id=user_id,
            cost=cost,
            job_manager=job_manager,
            wallet=wallet,
            refund_fn=_refund_safe,
        )
    except Exception as exc:  # noqa: BLE001
        _refund_safe(user_id, cost, "turnitin", ref_id=submission_id)
        app.logger.exception("turnitin/check failed")
        return jsonify({"success": False, "error": f"Could not queue check: {exc}"}), 500

    _record_usage_safe(
        user_id,
        feature=FEATURE_TURNITIN,
        credits_used=cost,
        provider="PlagDetect",
        request_id=job.id,
    )

    return jsonify(
        {
            "success": True,
            "submission_id": submission_id,
            "filename": filename,
            "status": "queued",
            "job_id": job.id,
            "coins_charged": cost,
            "balance": wallet.get_balance(user_id),
        }
    )


@app.get("/api/browser/providers/plagdetect/config")
def api_browser_plagdetect_config():
    from services.browser.providers import plagdetect as pd

    cfg = pd.plagdetect_config()
    return jsonify({"success": True, **cfg})


@app.get("/api/browser/providers/plagdetect/status")
def api_browser_plagdetect_status():
    try:
        ensure_engine_started()
        from services.browser.providers import plagdetect as pd

        payload = _browser_submit(pd.get_session_status, timeout=60)
        return jsonify({"success": True, **payload})
    except Exception as exc:  # noqa: BLE001
        app.logger.exception("plagdetect status failed")
        return jsonify({"success": False, "error": str(exc)}), 500


@app.get("/api/browser/providers/plagdetect/login")
def api_browser_plagdetect_login():
    try:
        ensure_engine_started()
        from services.browser.providers import plagdetect as pd

        payload = _browser_submit(pd.start_interactive_login, timeout=60)
        return jsonify(payload)
    except Exception as exc:  # noqa: BLE001
        app.logger.exception("plagdetect login failed")
        return jsonify({"success": False, "error": str(exc)}), 500


@app.route("/pricing")
def pricing():
    """Coin packages — clear pricing without subscription pressure."""
    return render_template(
        "pricing.html",
        nav_active="pricing",
        welcome_bonus=WELCOME_BONUS,
        paddle_configured=paddle_configured(),
        paddle_client_token=paddle_client_token(),
        paddle_environment=paddle_environment(),
    )


def _git_revision() -> str:
    import subprocess

    try:
        return subprocess.check_output(
            ["git", "rev-parse", "--short", "HEAD"],
            cwd=Path(__file__).resolve().parent,
            text=True,
            stderr=subprocess.DEVNULL,
        ).strip()
    except (OSError, subprocess.CalledProcessError):
        return "unknown"


@app.get("/api/version")
def api_version():
    """Lets you verify which code revision is running (local vs production)."""
    return jsonify(
        {
            "revision": _git_revision(),
            "engine": "reconstruction+style_engine",
        }
    )


def _telegram_credentials() -> tuple[str, str] | tuple[None, None]:
    token = (os.environ.get("TELEGRAM_TOKEN") or os.environ.get("TELEGRAM_BOT_TOKEN") or "").strip()
    chat_id = (os.environ.get("CHAT_ID") or os.environ.get("TELEGRAM_CHAT_ID") or "").strip()
    if not token or not chat_id:
        return None, None
    return token, chat_id


def _send_telegram_text(text: str) -> tuple[dict[str, Any] | None, str | None, int]:
    """POST sendMessage. Returns (tg_json, error_message, http_status_on_error)."""
    token, chat_id = _telegram_credentials()
    if not token or not chat_id:
        return None, "Telegram is not configured. Set TELEGRAM_TOKEN and CHAT_ID environment variables.", 503

    body = text if len(text) <= TELEGRAM_TEXT_MAX_LEN else text[: TELEGRAM_TEXT_MAX_LEN - 1] + "…"
    url = f"https://api.telegram.org/bot{token}/sendMessage"
    try:
        tg_res = requests.post(
            url,
            json={"chat_id": chat_id, "text": body},
            timeout=TELEGRAM_SEND_MESSAGE_TIMEOUT_S,
        )
    except requests.Timeout:
        return None, "Telegram did not respond in time.", 502
    except requests.RequestException:
        app.logger.exception("feedback: Telegram request failed")
        return None, "Could not reach Telegram.", 502

    try:
        tg_data = tg_res.json()
    except ValueError:
        return None, "Unexpected response from Telegram.", 502

    if not tg_res.ok or not tg_data.get("ok"):
        desc = tg_data.get("description") if isinstance(tg_data, dict) else None
        app.logger.warning(
            "feedback: Telegram API error %s — %s",
            tg_res.status_code,
            desc or (tg_res.text[:200] if tg_res.text else ""),
        )
        return None, "Could not deliver message.", 502
    return tg_data if isinstance(tg_data, dict) else {}, None, 200


def _feedback_from_request():
    """
    Parse JSON {\"message\": \"...\"}, store SupportMessage, send to Telegram.
    Requires a signed-in user so replies can be routed back via User ID.
    Returns (flask Response, http_status).
    """
    try:
        from services.economy.support_chat import (
            format_telegram_outbound,
            save_support_message,
        )

        user = economy_auth.current_user()
        if not user:
            return (
                jsonify(
                    {
                        "success": False,
                        "error": "AUTH_REQUIRED",
                        "message": "Please sign in to use support chat.",
                    }
                ),
                401,
            )

        raw = request.get_data(cache=True)
        payload = request.get_json(silent=True)
        if payload is None:
            if raw and request.is_json:
                return jsonify({"error": "Invalid JSON body."}), 400
            payload = {}
        if not isinstance(payload, dict):
            return jsonify({"error": "Body must be a JSON object."}), 400
        message = (payload.get("message") or "").strip()
        if not message:
            return jsonify({"error": "Message is required."}), 400

        user_id = int(user["id"])
        saved = save_support_message(user_id=user_id, sender="user", message=message)
        outbound = format_telegram_outbound(
            message=message,
            user_id=user_id,
            email=str(user.get("email") or ""),
            name=str(user.get("name") or ""),
        )
        tg_data, err, err_status = _send_telegram_text(outbound)
        if err:
            # Message is already saved so the transcript is not lost if Telegram is down.
            app.logger.error("feedback: Telegram delivery failed after save: %s", err)
            if err_status == 503:
                return jsonify({"error": err, "saved": True, "message": saved.to_dict()}), 503
            return (
                jsonify(
                    {
                        "error": err,
                        "saved": True,
                        "message": saved.to_dict(),
                        "success": False,
                    }
                ),
                err_status,
            )

        try:
            from services.economy.support_chat import bind_telegram_message

            result = (tg_data or {}).get("result") if isinstance(tg_data, dict) else None
            tg_mid = (result or {}).get("message_id") if isinstance(result, dict) else None
            if tg_mid is not None:
                bind_telegram_message(
                    telegram_message_id=int(tg_mid),
                    user_id=user_id,
                    support_message_id=int(saved.id) if saved.id else None,
                )
        except Exception:  # noqa: BLE001
            app.logger.exception("feedback: failed to bind telegram message_id")

        app.logger.info("feedback: delivered via Telegram user_id=%s msg_id=%s", user_id, saved.id)
        return jsonify({"success": True, "message": saved.to_dict()}), 200
    except Exception:  # noqa: BLE001
        app.logger.exception("feedback: unexpected error")
        return jsonify({"error": "Something went wrong processing feedback."}), 500


@app.route("/api/feedback", methods=["POST"])
def api_feedback():
    """POST JSON {\"message\": \"...\"}; stores chat + forwards to Telegram when configured."""
    body, status = _feedback_from_request()
    return body, status


@app.post("/feedback")
def feedback():
    """Backward-compatible alias for POST /api/feedback."""
    return api_feedback()


@app.get("/api/chat/messages")
@economy_auth.login_required
def api_chat_messages():
    """Return support chat history for the signed-in user (polling endpoint)."""
    from services.economy.support_chat import list_support_messages

    user_id = economy_auth.current_user_id()
    if user_id is None:
        return jsonify({"success": False, "error": "AUTH_REQUIRED"}), 401
    after_raw = (request.args.get("after_id") or "").strip()
    after_id = int(after_raw) if after_raw.isdigit() else None
    messages = list_support_messages(int(user_id), after_id=after_id)
    return jsonify(
        {
            "success": True,
            "messages": [m.to_dict() for m in messages],
        }
    )


@app.post("/api/telegram-webhook")
def api_telegram_webhook():
    """Telegram Bot API webhook — admin Reply → SupportMessage(sender=admin).

    Set the bot webhook to this URL. Optional ``TELEGRAM_WEBHOOK_SECRET`` must
    match ``?secret=`` or ``X-Telegram-Bot-Api-Secret-Token`` when configured.
    """
    from services.economy.support_chat import (
        normalize_chat_id,
        parse_admin_reply_from_update,
        save_support_message,
    )

    expected_secret = (os.environ.get("TELEGRAM_WEBHOOK_SECRET") or "").strip()
    if expected_secret:
        got = (
            (request.args.get("secret") or "").strip()
            or (request.headers.get("X-Telegram-Bot-Api-Secret-Token") or "").strip()
        )
        if got != expected_secret:
            app.logger.warning("telegram webhook: bad secret")
            return jsonify({"ok": False, "error": "unauthorized"}), 401

    payload = request.get_json(silent=True)
    if not isinstance(payload, dict):
        app.logger.info("telegram webhook: non-json body ignored")
        return jsonify({"ok": True, "ignored": True, "reason": "non_json"}), 200

    message = payload.get("message") or payload.get("edited_message") or {}
    _, admin_chat_id = _telegram_credentials()
    if admin_chat_id and isinstance(message, dict):
        chat = message.get("chat") or {}
        incoming_chat = normalize_chat_id(chat.get("id"))
        expected_chat = normalize_chat_id(admin_chat_id)
        if incoming_chat and expected_chat and incoming_chat != expected_chat:
            # Do not hard-drop: still try to route if User ID / map resolve.
            # Mismatched CHAT_ID formatting used to silently kill all replies.
            app.logger.warning(
                "telegram webhook: chat_id mismatch incoming=%s expected=%s (continuing)",
                incoming_chat,
                expected_chat,
            )

    parsed = parse_admin_reply_from_update(payload)
    if not parsed:
        has_reply = isinstance(message, dict) and isinstance(message.get("reply_to_message"), dict)
        app.logger.info(
            "telegram webhook: ignored has_reply=%s keys=%s",
            has_reply,
            sorted(payload.keys()),
        )
        return jsonify({"ok": True, "ignored": True, "reason": "not_a_routable_reply"}), 200

    try:
        saved = save_support_message(
            user_id=int(parsed["user_id"]),
            sender="admin",
            message=str(parsed["message"]),
        )
    except ValueError as exc:
        app.logger.warning("telegram webhook: persist rejected: %s", exc)
        return jsonify({"ok": True, "ignored": True, "reason": str(exc)}), 200
    except Exception:  # noqa: BLE001
        app.logger.exception("telegram webhook: persist failed")
        return jsonify({"ok": False, "error": "persist_failed"}), 500

    app.logger.info(
        "telegram webhook: admin reply saved user_id=%s msg_id=%s via=%s",
        saved.user_id,
        saved.id,
        parsed.get("via"),
    )
    return jsonify({"ok": True, "message": saved.to_dict(), "via": parsed.get("via")}), 200


@app.post("/parse-requirements")
def parse_requirements_view():
    """
    Extract formatting hints from free-form text (brief, OCR, etc.).
    Uses Gemini when GOOGLE_API_KEY is set; otherwise returns heuristic mock data.
    """
    payload = request.get_json(silent=True) or {}
    text = (payload.get("text") or "").strip()
    if not text:
        return jsonify({"error": "Provide non-empty text in a JSON body: {\"text\": \"...\"}"}), 400

    using_local_parser = not (os.environ.get("GOOGLE_API_KEY") or "").strip()

    try:
        requirements = parse_requirements(text)
    except Exception as e:  # noqa: BLE001
        app.logger.exception("parse-requirements failed")
        return jsonify({"error": f"Could not parse requirements: {str(e)}"}), 502

    form = form_autofill_from_parsed(requirements)
    return jsonify({"requirements": requirements, "form": form, "mock": using_local_parser})


def _extract_requirement_text_from_upload(file_storage) -> str:
    """Extract text from supported assignment brief uploads."""
    filename = (file_storage.filename or "").lower()
    ext = os.path.splitext(filename)[1]
    if ext not in REQUIREMENTS_DOC_EXT:
        raise ValueError("Unsupported file type. Supported formats: PDF, DOCX, TXT, JPG, PNG.")

    raw = file_storage.read()
    if not raw:
        raise ValueError("The uploaded requirements file is empty.")

    if ext in REQUIREMENTS_TEXT_EXT:
        return raw.decode("utf-8", errors="replace")

    if ext == ".docx":
        return extract_text_from_document_bytes(raw, filename)

    if ext in REQUIREMENTS_IMAGE_EXT:
        return extract_text_from_image_stream(io.BytesIO(raw))

    if ext == ".pdf":
        return extract_text_from_document_bytes(raw, filename)

    raise ValueError("Unsupported requirements file.")


@app.post("/api/extract-brief-text")
def api_extract_brief_text():
    """
    Extract plain text from an uploaded assignment brief (.pdf, .docx, .txt, .jpg, .png).
    Multipart field name: file
    """
    f = request.files.get("file")
    if not f or not f.filename:
        return jsonify({"error": 'Upload a brief file as form field "file".'}), 400

    try:
        text = _extract_requirement_text_from_upload(f).strip()
    except ValueError as e:
        return jsonify({"error": str(e)}), 400
    except RuntimeError as e:
        app.logger.warning("extract-brief-text: %s", e)
        return jsonify({"error": str(e)}), 503
    except Exception:  # noqa: BLE001
        app.logger.exception("extract-brief-text failed")
        return jsonify({"error": "Could not read the brief file."}), 400

    if len(text) > MAX_TEXT_CHARS:
        return jsonify({"error": f"Brief text is too long (max {MAX_TEXT_CHARS:,} characters)."}), 400

    return jsonify({"text": text, "filename": f.filename})


@app.post("/api/extract-requirements")
def api_extract_requirements():
    """
    Extract precise formatting requirements from pasted text and/or uploaded brief.
    Multipart: optional requirements_text, optional file (.docx, .pdf, .txt, .md, .jpg, .png)
    """
    pasted = (request.form.get("requirements_text") or "").strip()
    f = request.files.get("file")
    uploaded_text = ""

    if f and f.filename:
        try:
            uploaded_text = _extract_requirement_text_from_upload(f).strip()
        except ValueError as e:
            return jsonify({"error": str(e)}), 400
        except RuntimeError as e:
            app.logger.warning("extract-requirements: %s", e)
            return jsonify({"error": str(e)}), 503
        except Exception:  # noqa: BLE001
            app.logger.exception("extract-requirements failed")
            return jsonify({"error": "Could not read the requirements file."}), 400

    text = "\n\n".join(x for x in (pasted, uploaded_text) if x).strip()
    if not text:
        return jsonify({"error": "Paste requirements or upload a supported brief."}), 400
    if len(text) > MAX_TEXT_CHARS:
        return jsonify({"error": f"Requirements text is too long (max {MAX_TEXT_CHARS:,} characters)."}), 400

    try:
        requirements = parse_requirements(text)
    except Exception as e:  # noqa: BLE001
        app.logger.exception("extract-requirements parse failed")
        return jsonify({"error": f"Could not parse requirements: {str(e)}"}), 502

    return jsonify(
        {
            "requirements": requirements,
            "form": form_autofill_from_parsed(requirements),
            "source_text": text,
            "source_text_chars": len(text),
            "mock": not (os.environ.get("GOOGLE_API_KEY") or "").strip(),
        }
    )


@app.post("/api/requirements-ocr")
def api_requirements_ocr():
    """
    OCR a JPEG/PNG and return extracted text for the requirements parser.
    Multipart field name: image
    """
    f = request.files.get("image")
    if not f or not f.filename:
        return jsonify({"error": 'Upload an image as form field "image" (JPEG or PNG).'}), 400

    name = (f.filename or "").lower()
    if not any(name.endswith(ext) for ext in REQUIREMENTS_IMAGE_EXT):
        return jsonify({"error": "Only JPEG and PNG images are supported."}), 400

    try:
        raw = f.read()
        if not raw:
            return jsonify({"error": "Empty file."}), 400
        text = extract_text_from_image_stream(io.BytesIO(raw))
    except ValueError as e:
        return jsonify({"error": str(e)}), 400
    except RuntimeError as e:
        app.logger.warning("requirements-ocr: %s", e)
        return jsonify({"error": str(e)}), 503
    except Exception as e:  # noqa: BLE001
        app.logger.exception("requirements-ocr failed")
        return jsonify({"error": f"OCR failed: {str(e)}"}), 500

    if not text.strip():
        return (
            jsonify(
                {
                    "error": "No text could be read from the image. "
                    "Try a clearer photo or paste text manually."
                }
            ),
            422,
        )

    return jsonify({"text": text})


@app.post("/api/reference")
def api_reference():
    """
    Generate a reference citation.
    JSON: mode (url|doi|isbn|title|manual|paste), style, and mode-specific fields.
    Backward compatible: url and/or text without mode → url mode.
    """
    payload = request.get_json(silent=True) or {}
    mode = (payload.get("mode") or "").strip().lower()
    style = (payload.get("style") or "APA").strip()

    if not mode:
        if payload.get("doi"):
            mode = "doi"
        elif payload.get("isbn"):
            mode = "isbn"
        elif payload.get("title"):
            mode = "title"
        elif payload.get("manual"):
            mode = "manual"
        elif payload.get("paste") or payload.get("text"):
            mode = "paste" if payload.get("paste") else ("url" if payload.get("url") else "paste")
        elif payload.get("url"):
            mode = "url"
        else:
            return jsonify({"error": "Provide input fields or a mode (url, doi, isbn, title, manual, paste)."}), 400

    if style.upper() not in ENGINE_CITATION_STYLES:
        style = "APA"

    try:
        result = generate_citation(
            mode=mode,
            style=style,
            url=(payload.get("url") or "").strip() or None,
            doi=(payload.get("doi") or "").strip() or None,
            isbn=(payload.get("isbn") or "").strip() or None,
            title=(payload.get("title") or "").strip() or None,
            author=(payload.get("author") or "").strip() or None,
            manual=payload.get("manual") if isinstance(payload.get("manual"), dict) else None,
            paste=(payload.get("paste") or payload.get("text") or "").strip() or None,
        )
    except ValueError as e:
        return jsonify({"error": str(e)}), 400
    except requests.RequestException as e:
        app.logger.exception("reference generation failed")
        return jsonify({"error": f"Could not retrieve metadata: {str(e)}"}), 502
    except Exception:  # noqa: BLE001
        app.logger.exception("reference generation failed")
        return jsonify({"error": "Could not generate citation."}), 500

    return jsonify(result)


@app.post("/api/intext-citation")
def api_intext_citation():
    """Generate in-text citations, footnotes, and endnotes."""
    payload = request.get_json(silent=True) or {}
    author = (payload.get("author") or "").strip()
    year = (payload.get("year") or "n.d.").strip()
    page = (payload.get("page") or "").strip() or None
    style = (payload.get("style") or "APA").strip()
    quote = bool(payload.get("direct_quote"))
    if not author:
        return jsonify({"error": "Author is required."}), 400
    return jsonify(generate_intext(author=author, year=year, page=page, style=style, quote=quote))


@app.post("/api/format-references")
def api_format_references():
    """
    Alphabetize a citation list and pick section title for APA / MLA / Harvard.
    Body: {\"citations\": [\"...\"], \"style\": \"APA\" | \"MLA\" | \"Harvard\"}
    """
    payload = request.get_json(silent=True) or {}
    cites = payload.get("citations")
    if not isinstance(cites, list):
        return jsonify({"error": "'citations' must be a JSON array of strings."}), 400
    style_raw = (payload.get("style") or "APA").strip()
    if style_raw.upper() not in CITATION_STYLES:
        style_raw = "APA"
    lines = [str(x) for x in cites if str(x).strip()]
    if not lines:
        return jsonify({"error": "Provide at least one non-empty citation string."}), 400
    heading, sorted_lines = prepare_reference_section(lines, style_raw)
    block = heading + "\n" + "\n".join(sorted_lines)
    return jsonify(
        {
            "section_title": heading,
            "citations": sorted_lines,
            "text": block,
        }
    )


@app.post("/api/extract-document")
def api_extract_document():
    """
    Extract plain text from an uploaded .docx or .pdf for the formatting workflow.
    Multipart field name: file
    """
    f = request.files.get("file")
    if not f or not f.filename:
        return jsonify({"error": 'Upload a .docx or .pdf file as form field "file".'}), 400
    if not is_supported_document_upload(f.filename, f.mimetype):
        return jsonify({"error": "Invalid file type. Upload a .docx or .pdf file."}), 400

    try:
        raw = f.read()
        text = extract_text_from_document_bytes(raw, f.filename, f.mimetype)
    except ValueError as e:
        return jsonify({"error": str(e)}), 400
    except RuntimeError as e:
        app.logger.warning("extract-document: %s", e)
        return jsonify({"error": str(e)}), 503
    except Exception:  # noqa: BLE001
        app.logger.exception("extract-document failed")
        return jsonify({"error": "Could not read the uploaded file."}), 400

    if len(text) > MAX_TEXT_CHARS:
        return jsonify({"error": f"Document text is too long (max {MAX_TEXT_CHARS:,} characters)."}), 400

    return jsonify({"text": text, "filename": f.filename})


@app.post("/api/check-document")
def api_check_document():
    """
    Smart document check: requirements + text/docx/pdf → score, categories, issue cards.
    Multipart: requirements, pasted_text, document_type, optional file (.docx or .pdf).
    """
    requirements = (request.form.get("requirements") or "").strip()
    pasted = (request.form.get("pasted_text") or "").strip()
    doc_type = (request.form.get("document_type") or "other").strip()
    parsed_requirements = None
    raw_parsed = request.form.get("parsed_requirements")
    if raw_parsed:
        try:
            parsed_requirements = json.loads(raw_parsed)
        except (json.JSONDecodeError, TypeError):
            parsed_requirements = None

    if len(requirements) > MAX_TEXT_CHARS:
        return jsonify({"error": f"Requirements text is too long (max {MAX_TEXT_CHARS:,} characters)."}), 400
    if len(pasted) > MAX_TEXT_CHARS:
        return jsonify({"error": f"Document text is too long (max {MAX_TEXT_CHARS:,} characters)."}), 400

    doc: Document | None = None
    f = request.files.get("file")

    if f and f.filename:
        if not is_supported_document_upload(f.filename, f.mimetype):
            return jsonify({"error": "Invalid file type. Upload a .docx or .pdf file."}), 400
        try:
            raw = f.read()
            if not raw:
                return jsonify({"error": "The uploaded file is empty."}), 400
            ext = upload_extension(f.filename, f.mimetype)
            if ext == ".docx":
                doc = Document(io.BytesIO(raw))
            doc_text = extract_text_from_document_bytes(raw, f.filename, f.mimetype)
        except ValueError as e:
            return jsonify({"error": str(e)}), 400
        except RuntimeError as e:
            app.logger.warning("check-document: %s", e)
            return jsonify({"error": str(e)}), 503
        except Exception:  # noqa: BLE001
            app.logger.exception("check-document: invalid upload")
            return jsonify({"error": "Could not read the uploaded file."}), 400
    else:
        doc_text = ""

    text = pasted
    if doc_text:
        if not text:
            text = doc_text
        elif text != doc_text:
            # Prefer uploaded file content when both are present.
            text = doc_text

    if not text:
        return jsonify(
            {"error": "Provide non-empty text or upload a .docx or .pdf with readable content."}
        ), 400

    # Academic Check is free for everyone (no login, no coins).
    try:
        result = check_document(
            text=text,
            requirements=requirements,
            doc=doc,
            document_type=doc_type,
            parsed_requirements=parsed_requirements,
        )
    except Exception:  # noqa: BLE001
        app.logger.exception("check-document failed")
        return jsonify({"error": "Document check failed. Please try again."}), 500

    if result.get("error"):
        return jsonify({"error": result["error"]}), 400

    return jsonify(result)


@app.post("/api/structure-recovery")
def api_structure_recovery():
    """
    Reconstruct academic document structure from pasted text or uploaded .docx / .pdf.
    Multipart: pasted_text, document_type (optional), optional file (.docx or .pdf).
    """
    pasted = (request.form.get("pasted_text") or "").strip()
    doc_type = (request.form.get("document_type") or "other").strip()

    if len(pasted) > MAX_TEXT_CHARS:
        return jsonify({"error": f"Document text is too long (max {MAX_TEXT_CHARS:,} characters)."}), 400

    doc: Document | None = None
    f = request.files.get("file")

    if f and f.filename:
        if not is_supported_document_upload(f.filename, f.mimetype):
            return jsonify({"error": "Invalid file type. Upload a .docx or .pdf file."}), 400
        try:
            raw = f.read()
            if not raw:
                return jsonify({"error": "The uploaded file is empty."}), 400
            ext = upload_extension(f.filename, f.mimetype)
            if ext == ".docx":
                doc = Document(io.BytesIO(raw))
            doc_text = extract_text_from_document_bytes(raw, f.filename, f.mimetype)
        except ValueError as e:
            return jsonify({"error": str(e)}), 400
        except RuntimeError as e:
            app.logger.warning("structure-recovery: %s", e)
            return jsonify({"error": str(e)}), 503
        except Exception:  # noqa: BLE001
            app.logger.exception("structure-recovery: invalid upload")
            return jsonify({"error": "Could not read the uploaded file."}), 400
    else:
        doc_text = ""

    text = pasted
    if doc_text:
        if not text:
            text = doc_text
        elif text != doc_text:
            text = doc_text

    if not text and doc is None:
        return jsonify(
            {"error": "Provide non-empty text or upload a .docx or .pdf with readable content."}
        ), 400

    try:
        result = recover_structure(text=text or None, doc=doc, document_type=doc_type)
    except Exception:  # noqa: BLE001
        app.logger.exception("structure-recovery failed")
        return jsonify({"error": "Structure recovery failed. Please try again."}), 500

    if result.get("error"):
        if result.get("ai_failure"):
            return jsonify(result), 503
        return jsonify({"error": result["error"]}), 400

    return jsonify(result)


@app.post("/api/preview-formatted")
def preview_formatted():
    """Server-side After preview — same formatting pipeline as /api/format."""
    payload = request.get_json(silent=True) or {}
    text = (payload.get("text") or "").strip()
    if not text:
        return jsonify({"error": "No text to preview."}), 400
    settings = payload.get("settings") or {}
    job = parse_job(settings)
    document_type = (payload.get("document_type") or settings.get("document_type") or "").strip() or None
    required_sections: list[str] = []
    brief = (payload.get("requirements_text") or settings.get("requirements_text") or "").strip()
    if brief:
        from formatter.requirement_headings import extract_format_section_labels

        required_sections = extract_format_section_labels(brief)
    try:
        clean_spaces = (
            _truthy(settings, "clean_extra_spaces")
            if "clean_extra_spaces" in settings
            else True
        )
        clean_breaks = (
            _truthy(settings, "clean_extra_linebreaks")
            if "clean_extra_linebreaks" in settings
            else False
        )
        html = build_formatted_preview_html(
            text,
            job,
            document_type=document_type,
            required_sections=required_sections if job.requirement_headings else None,
            cleaning_spaces=clean_spaces,
            cleaning_breaks=clean_breaks,
        )
    except Exception as exc:
        app.logger.exception("preview-formatted failed")
        return jsonify({"error": str(exc)}), 500
    return jsonify({"html": html})


@app.post("/api/format")
def format_document():
    try:
        job = parse_job(request.form)

        file_storage = request.files.get("file")
        pasted_raw = request.form.get("pasted_text") or ""
        clean_spaces = _truthy(request.form, "clean_extra_spaces")
        clean_breaks = _truthy(request.form, "clean_extra_linebreaks")

        if file_storage and file_storage.filename:
            if not is_supported_document_upload(file_storage.filename, file_storage.mimetype):
                return (
                    jsonify(
                        {
                            "error": "Invalid file type. Upload a .docx or .pdf file.",
                        }
                    ),
                    400,
                )
            raw = file_storage.read()
            if not raw:
                return jsonify({"error": "The uploaded file is empty."}), 400
            try:
                doc = build_document_from_upload(
                    raw,
                    file_storage.filename,
                    mimetype=file_storage.mimetype,
                    cleaning_spaces=clean_spaces,
                    cleaning_breaks=clean_breaks,
                )
            except ValueError as e:
                return jsonify({"error": str(e)}), 400
            except RuntimeError as e:
                app.logger.warning("format: %s", e)
                return jsonify({"error": str(e)}), 503
        elif pasted_raw.strip():
            doc = build_document_from_inputs(
                pasted_raw=pasted_raw,
                file_bytes=None,
                cleaning_spaces=clean_spaces,
                cleaning_breaks=clean_breaks,
            )
        else:
            return (
                jsonify(
                    {
                        "error": "Please upload a .docx or .pdf file, or paste some non-empty text.",
                    }
                ),
                400,
            )

        if not document_has_visible_text(doc):
            return (
                jsonify(
                    {
                        "error": "The document looks empty after loading. Add text and try again.",
                    }
                ),
                400,
            )

        fallback_paragraphs = [p.text.strip() for p in doc.paragraphs if (p.text or "").strip()]
        if not fallback_paragraphs and pasted_raw.strip():
            fallback_paragraphs = paragraphs_from_text(pasted_raw)

        paragraph_assignments: list[ParagraphHeadingAssignment] | None = None
        recovery_mode = ""
        ai_powered = False
        doc_type = (request.form.get("document_type") or "other").strip() or None
        brief = (request.form.get("requirements_text") or "").strip()
        required_sections: list[str] = []
        if brief:
            from formatter.requirement_headings import extract_format_section_labels

            required_sections = extract_format_section_labels(brief)

        if job.auto_headings or job.requirement_headings:
            recon = reconstruct_document_before_format(
                doc,
                document_type=doc_type,
                required_sections=required_sections if job.requirement_headings else None,
                prefer_ai=job.auto_headings,
            )
            paragraph_assignments = recon.assignments
            recovery_mode = recon.recovery_mode
            ai_powered = recon.ai_powered

        before_cover_paragraph_count = len(doc.paragraphs)
        cover_storage = request.files.get("cover_file")
        cover_doc = None
        if cover_storage and cover_storage.filename:
            if not is_supported_document_upload(cover_storage.filename, cover_storage.mimetype):
                return (
                    jsonify(
                        {
                            "error": "Cover page must be a .docx or .pdf file.",
                        }
                    ),
                    400,
                )
            cover_raw = cover_storage.read()
            if not cover_raw:
                return jsonify({"error": "The uploaded cover page file is empty."}), 400
            try:
                cover_doc = build_document_from_upload(
                    cover_raw,
                    cover_storage.filename,
                    mimetype=cover_storage.mimetype,
                    cleaning_spaces=False,
                    cleaning_breaks=False,
                )
            except ValueError as e:
                return jsonify({"error": f"Cover page: {e}"}), 400
            except RuntimeError as e:
                app.logger.warning("format cover: %s", e)
                return jsonify({"error": str(e)}), 503

        cover = None if cover_doc is not None else parse_cover_page(
            request.form, fallback_paragraphs=fallback_paragraphs
        )
        if cover_doc is not None:
            prepend_cover_document(doc, cover_doc)
            if paragraph_assignments:
                inserted = len(doc.paragraphs) - before_cover_paragraph_count
                if inserted > 0:
                    paragraph_assignments = (
                        [ParagraphHeadingAssignment()] * inserted + paragraph_assignments
                    )
        elif cover:
            prepend_cover_page(doc, cover, font_family=job.font_family)
            if paragraph_assignments:
                inserted = len(doc.paragraphs) - before_cover_paragraph_count
                if inserted > 0:
                    paragraph_assignments = (
                        [ParagraphHeadingAssignment()] * inserted + paragraph_assignments
                    )

        structure_debug = (
            os.environ.get("STRUCTURE_RECOVERY_DEBUG", "").strip().lower() in {"1", "true", "yes"}
            or _truthy(request.form, "structure_recovery_debug")
        )

        debug_report = format_document_full(
            doc,
            job,
            paragraph_assignments,
            structure_debug=structure_debug,
            recovery_mode=recovery_mode,
            ai_powered=ai_powered,
        )

        ref_lines = [r.strip() for r in request.form.getlist("references") if r.strip()]
        if ref_lines:
            style = (request.form.get("citation_style") or "APA").strip()
            if style.upper() not in CITATION_STYLES:
                style = "APA"
            section_heading, sorted_refs = prepare_reference_section(ref_lines, style)
            append_references_section(doc, job, sorted_refs, section_title=section_heading)

        out = io.BytesIO()
        doc.save(out)
        out.seek(0)

        response = send_file(
            out,
            mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            as_attachment=True,
            download_name="formatted_document.docx",
        )
        if debug_report:
            response.headers["X-Structure-Recovery-Debug"] = json.dumps(debug_report.to_dict())
        return response
    except Exception as e:  # noqa: BLE001
        app.logger.exception("Format failed")
        return jsonify({"error": f"Could not format document: {str(e)}"}), 500


def formatter_v2_enabled() -> bool:
    """Feature flag for the parallel Formatter V2 pipeline (default off)."""
    return os.environ.get("FORMATTER_V2_ENABLED", "0").strip().lower() in {
        "1",
        "true",
        "yes",
        "on",
    }


@app.get("/format-v2")
def format_v2_page():
    """Formatter V2 UI. Gated by ``FORMATTER_V2_ENABLED`` (404 when off)."""
    if not formatter_v2_enabled():
        return jsonify({"error": "Not found"}), 404
    return render_template("format_v2.html", nav_active="home")


@app.get("/api/format-v2/profile/<style>")
def api_format_v2_profile(style: str):
    """Return StyleProfile (+ flattened form defaults) as JSON for the V2 UI."""
    if not formatter_v2_enabled():
        return jsonify({"error": "Not found"}), 404
    try:
        from formatter_v2.web_api import profile_payload

        return jsonify(profile_payload(style))
    except ValueError as exc:
        return jsonify({"error": str(exc)}), 404
    except Exception as exc:  # noqa: BLE001
        app.logger.exception("format-v2 profile failed")
        return jsonify({"error": f"Could not load profile: {exc}"}), 500


@app.post("/api/extract-requirements-v2")
def api_extract_requirements_v2():
    """Parse an assignment brief via Formatter V2 smartform.

    Multipart: ``requirements_text`` (required). Returns
    ``{overrides, evidence_by_field, warnings, prompt_version}``.
    """
    if not formatter_v2_enabled():
        return jsonify({"error": "Not found"}), 404

    brief = (request.form.get("requirements_text") or "").strip()
    if not brief:
        return jsonify({"error": "Paste requirements text first."}), 400
    if len(brief) > MAX_TEXT_CHARS:
        return (
            jsonify(
                {
                    "error": (
                        f"Requirements text is too long "
                        f"(max {MAX_TEXT_CHARS:,} characters)."
                    )
                }
            ),
            400,
        )

    style_raw = (
        request.form.get("format_style")
        or request.form.get("style")
        or "harvard"
    )

    try:
        from formatter_v2.pipeline import resolve_style_name
        from formatter_v2.profiles import load_profile
        from formatter_v2.smartform import extract_requirements, to_user_overrides
        from formatter_v2.smartform.extract import PROMPT_VERSION
        from formatter_v2.smartform.gemini_client import GeminiSmartformClient

        style = resolve_style_name(style_raw)
        profile = load_profile(style)
        extracted = extract_requirements(brief, GeminiSmartformClient())
        if extracted.style is not None:
            profile = load_profile(extracted.style)
        prefill = to_user_overrides(extracted, profile)
        return jsonify(
            {
                "overrides": prefill.overrides.model_dump(mode="json", exclude_none=True),
                "evidence_by_field": prefill.evidence_by_field,
                "warnings": list(extracted.warnings),
                "unsupported": list(extracted.unsupported),
                "prompt_version": PROMPT_VERSION,
                "style": (extracted.style.value if extracted.style else style.value),
            }
        )
    except Exception as exc:  # noqa: BLE001
        app.logger.exception("extract-requirements-v2 failed")
        return jsonify({"error": f"Could not analyse brief: {exc}"}), 502


@app.post("/api/format-v2")
def format_document_v2_route():
    """V2 formatter behind FORMATTER_V2_ENABLED.

    Accepts the same document inputs as ``/api/format`` (file or pasted_text)
    plus optional ``overrides`` JSON and ``format_style``.

    Returns JSON ``{document_id, summary, rejected, notices, overrides}``.
    Download the DOCX via ``GET /api/format-v2/download/<document_id>``.
    """
    if not formatter_v2_enabled():
        return jsonify({"error": "Not found"}), 404

    try:
        from formatter_v2.pipeline import format_document_v2
        from formatter_v2.web_api import parse_user_overrides_from_form

        file_storage = request.files.get("file")
        pasted_raw = request.form.get("pasted_text") or ""
        clean_spaces = _truthy(request.form, "clean_extra_spaces")
        clean_breaks = _truthy(request.form, "clean_extra_linebreaks")

        source: object
        if file_storage and file_storage.filename:
            if not is_supported_document_upload(file_storage.filename, file_storage.mimetype):
                return (
                    jsonify(
                        {
                            "error": "Invalid file type. Upload a .docx or .pdf file.",
                        }
                    ),
                    400,
                )
            raw = file_storage.read()
            if not raw:
                return jsonify({"error": "The uploaded file is empty."}), 400
            try:
                # Reuse V1 loaders for PDF→DOCX / cleanup only; formatting is V2.
                doc = build_document_from_upload(
                    raw,
                    file_storage.filename,
                    mimetype=file_storage.mimetype,
                    cleaning_spaces=clean_spaces,
                    cleaning_breaks=clean_breaks,
                )
            except ValueError as e:
                return jsonify({"error": str(e)}), 400
            except RuntimeError as e:
                app.logger.warning("format-v2: %s", e)
                return jsonify({"error": str(e)}), 503
            source = doc
        elif pasted_raw.strip():
            source = [line for line in pasted_raw.replace("\r\n", "\n").split("\n")]
        else:
            return (
                jsonify(
                    {
                        "error": "Please upload a .docx or .pdf file, or paste some non-empty text.",
                    }
                ),
                400,
            )

        try:
            overrides = parse_user_overrides_from_form(request.form)
        except ValueError as e:
            return jsonify({"error": str(e)}), 400

        style = (
            request.form.get("format_style")
            or request.form.get("style_preset")
            or request.form.get("citation_style")
            or (overrides.style.value if overrides.style else None)
            or "harvard"
        )
        result = format_document_v2(source, overrides, style)
        from formatter_v2.document_store import get_document_store
        from formatter_v2.web_api import format_v2_response_payload

        document_id = get_document_store().save(result.docx_bytes)
        return jsonify(
            format_v2_response_payload(
                document_id=document_id,
                overrides=overrides,
                notices=result.notices,
            )
        )
    except Exception as e:  # noqa: BLE001
        app.logger.exception("Format V2 failed")
        return jsonify({"error": f"Could not format document: {str(e)}"}), 500


@app.get("/api/format-v2/download/<document_id>")
def format_v2_download_route(document_id: str):
    """Download a formatted DOCX stored by ``POST /api/format-v2`` or chat."""
    if not formatter_v2_enabled():
        return jsonify({"error": "Not found"}), 404

    from formatter_v2.document_store import get_document_store

    path = get_document_store().resolve(document_id)
    if path is None:
        return jsonify({"error": "Document not found or expired."}), 404

    return send_file(
        path,
        mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        as_attachment=True,
        download_name="formatted_document.docx",
    )


@app.post("/api/format-v2/chat")
def format_v2_chat_route():
    """Apply a post-format chat edit.

    Returns JSON ``{document_id, summary, rejected, notices, overrides}``.
    Download the DOCX via ``GET /api/format-v2/download/<document_id>``.
    """
    if not formatter_v2_enabled():
        return jsonify({"error": "Not found"}), 404

    message = (request.form.get("message") or "").strip()
    if not message:
        return jsonify({"error": "message is required"}), 400

    try:
        from formatter_v2.chat.apply import apply_chat_edit
        from formatter_v2.chat.gemini_client import GeminiChatClient
        from formatter_v2.document_store import get_document_store
        from formatter_v2.pipeline import format_document_v2
        from formatter_v2.web_api import (
            format_v2_response_payload,
            load_format_v2_source_from_form,
            parse_user_overrides_from_form,
        )

        try:
            current_overrides = parse_user_overrides_from_form(request.form)
        except ValueError as e:
            return jsonify({"error": str(e)}), 400

        try:
            source = load_format_v2_source_from_form(
                request.form,
                request.files,
                build_document_from_upload=build_document_from_upload,
                is_supported_document_upload=is_supported_document_upload,
            )
        except ValueError as e:
            return jsonify({"error": str(e)}), 400
        except RuntimeError as e:
            app.logger.warning("format-v2 chat: %s", e)
            return jsonify({"error": str(e)}), 503

        style = (
            request.form.get("format_style")
            or request.form.get("style_preset")
            or (current_overrides.style.value if current_overrides.style else None)
            or "harvard"
        )

        new_overrides, summary, rejected = apply_chat_edit(
            message,
            current_overrides,
            style,
            GeminiChatClient(),
        )

        result = format_document_v2(source, new_overrides, style)
        document_id = get_document_store().save(result.docx_bytes)
        return jsonify(
            format_v2_response_payload(
                document_id=document_id,
                overrides=new_overrides,
                notices=result.notices,
                summary=summary,
                rejected=[
                    {"request": r.request, "reason": r.reason} for r in rejected
                ],
            )
        )
    except Exception as e:  # noqa: BLE001
        app.logger.exception("Format V2 chat failed")
        return jsonify({"error": f"Could not apply edit: {str(e)}"}), 500


@app.get("/api/test/zerogpt")
def api_test_zerogpt():
    """Connectivity test endpoint for raw ZeroGPT Detection API response."""
    sample_text = "This essay discusses artificial intelligence in higher education."
    try:
        zerogpt_client.login()
        raw = zerogpt_client.detect(sample_text)
        return jsonify(raw), 200
    except (ZeroGPTError, ZeroGPTProviderError) as exc:
        return jsonify({"error": str(exc)}), 502


@app.post("/api/ai/orchestrator/review")
def api_ai_orchestrator_review():
    """Run AIOrchestrator-compatible review flow for Review Engine."""
    payload = request.get_json(silent=True) or {}
    text = str(payload.get("text") or "").strip()
    if not text:
        return jsonify({"error": "text is required"}), 400
    return jsonify(orchestrator_review(text, client=zerogpt_client)), 200


@app.get("/api/test/zerogpt-humanizer")
def api_test_zerogpt_humanizer():
    """Connectivity test endpoint for raw ZeroGPT Humanizer API response."""
    sample_text = "Artificial intelligence improves productivity in academic writing."
    try:
        zerogpt_client.login()
        raw = zerogpt_client.humanize(sample_text)
        return jsonify(raw), 200
    except (ZeroGPTError, ZeroGPTProviderError) as exc:
        return jsonify({"error": str(exc)}), 502


@app.get("/api/test/env")
def api_test_env():
    return jsonify(
        {
            "email_exists": os.getenv("ZEROGPT_EMAIL") is not None,
            "password_exists": os.getenv("ZEROGPT_PASSWORD") is not None,
            "api_key_exists": os.getenv("ZEROGPT_API_KEY") is not None,
            "cwd": os.getcwd(),
            "dotenv_found": os.path.exists(".env"),
        }
    ), 200


@app.get("/api/test/gemini")
def api_test_gemini():
    """Temporary debug endpoint: raw Gemini API connectivity test."""
    api_key = (os.environ.get("GOOGLE_API_KEY") or os.environ.get("GEMINI_API_KEY") or "").strip()
    from services.gemini_client import _gemini_auth, gemini_model

    model = gemini_model()
    url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent"
    body = {
        "contents": [
            {
                "parts": [
                    {
                        "text": "Reply with exactly OK",
                    }
                ]
            }
        ]
    }

    auth_params, auth_headers = _gemini_auth(api_key)
    res = requests.post(
        url,
        params=auth_params,
        headers=auth_headers,
        json=body,
        timeout=30,
    )
    content_type = (res.headers.get("Content-Type") or "").lower()
    if "application/json" in content_type:
        response_body = res.json()
    else:
        response_body = res.text

    return jsonify(
        {
            "http_status": res.status_code,
            "response_body": response_body,
            "response_headers": dict(res.headers),
        }
    ), 200


# ---------------------------------------------------------------------------
# Browser production execution engine (JobManager + BrowserWorker + Monitor)
# ---------------------------------------------------------------------------
import threading as _threading  # noqa: E402

browser_metrics = None
job_manager = None
browser_worker = None
health_monitor = None
_engine_started = False
_engine_lock = _threading.Lock()


def _register_browser_providers() -> None:
    """Register providers with BrowserService (cheap; no Chrome launch)."""
    try:
        from services.browser.browser_service import BrowserService
        from services.browser.providers.plagdetect import PlagDetectProvider
        from services.browser.providers.stealthwriter import StealthWriterProvider

        service = BrowserService.instance()
        service.register_provider(StealthWriterProvider())
        service.register_provider(PlagDetectProvider())
    except Exception as exc:  # noqa: BLE001
        print(f"[browser] provider registration skipped: {exc}", flush=True)


def ensure_engine_started() -> None:
    """Idempotently start the browser worker + health monitor.

    The worker thread owns Playwright and auto-starts Chrome inside its own
    thread on first run, satisfying the "zero manual actions" requirement.
    """
    global browser_metrics, job_manager, browser_worker, health_monitor, _engine_started
    if _engine_started:
        return
    with _engine_lock:
        if _engine_started:
            return
        from services.browser.browser_service import BrowserService
        from services.browser.health_monitor import HealthMonitor
        from services.browser.jobs.job_manager import JobManager
        from services.browser.jobs.metrics import Metrics
        from services.browser.jobs.worker import BrowserWorker

        service = BrowserService.instance()
        browser_metrics = Metrics()
        browser_worker = BrowserWorker(service, None, browser_metrics)
        job_manager = JobManager(enqueue=browser_worker.enqueue_job, metrics=browser_metrics)
        browser_worker.attach_job_manager(job_manager)
        browser_worker.start()

        health_monitor = HealthMonitor(service, job_manager, browser_metrics, browser_worker)
        health_monitor.start()

        _engine_started = True
        print("[browser] production engine started (worker + health monitor)", flush=True)


def _browser_submit(fn, timeout=None):
    """Run a browser-touching callable on the worker thread (Playwright owner)."""
    ensure_engine_started()
    return browser_worker.submit(fn, timeout)


def _install_browser_submitter() -> None:
    from services.browser.thread_affinity import set_browser_submitter

    set_browser_submitter(lambda fn, timeout=None: _browser_submit(fn, timeout=timeout))


# Providers are known as soon as the app module loads, even before Chrome starts.
_register_browser_providers()
_install_browser_submitter()

try:
    from formatter_v2.document_store import ensure_store_started

    ensure_store_started(testing=app.config.get("TESTING", False))
except Exception as _format_v2_store_exc:  # noqa: BLE001
    print(f"[format-v2] document store init skipped: {_format_v2_store_exc}", flush=True)


if __name__ == "__main__":
    ensure_engine_started()
    app.run(debug=True, port=5001, use_reloader=False)
